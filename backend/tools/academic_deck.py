"""
tools/academic_deck.py

Assembles the 10-slide final presentation deck (.pptx) from AI-generated
slide JSON and light-mode chart images.

The deck rebuild (May 28 2026) is JSON-driven: a single Academic Writer
call (deck_generation_prompt() through harness_narrative()) returns
content for all ten slides — {slide_number, title, bullets, table_data,
speaker_notes} — and build_presentation_deck() lays them out. The ten
canonical SLIDE_TITLES and per-slide chart roles (SLIDE_CHARTS) are
fixed; the AI fills the content from the live data context.

  render_deck_charts()  — renders the legacy five canvas/deck charts as
    light-mode PNGs with matplotlib (white background, navy ink — NOT the
    platform's dark theme). Kept because chart_render.py's canvas-render
    path depends on it; the rebuilt deck renders its per-slide charts
    through the chart_render deck renderers instead. matplotlib is lazy
    and guarded — unavailable matplotlib / missing data returns None and
    the slide degrades to a [DATA PENDING] note.

  build_presentation_deck(slides, charts)  — lays out the ten slides in a
    professional navy/white theme and returns the .pptx bytes. Pure
    assembly: no LLM calls, no database reads. Always emits exactly ten
    slides with the canonical titles; a missing slide / table / chart
    degrades to a [DATA PENDING] note rather than failing.

Every slide carries the AI DRAFT footer and a speaker-notes verification
caveat — the deck is a first draft for the team to refine before the
July 1 presentation.
"""
from __future__ import annotations

import io
import re
from datetime import datetime, timezone
from typing import Any

import structlog
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from tools.academic_export import (
    table_drawdown, table_factor_loadings, table_regime_conditional,
    table_summary_statistics,
)

log = structlog.get_logger(__name__)

# ── Validated production constants (June 22 2026 -- Path A revised) ──
# These are the locked submission figures threaded through every document
# generator (brief, appendix, deck, script). They backstop the AI prompts
# so the writer never invents a number, and they flow through the
# substitution layer (tools/numeric_substitution.build_substitution_table)
# so every {{TOKEN}} in the generated prose resolves deterministically.
#
# OOS_SHARPE figures: LOCKED TO THE DECEMBER 2025 ACADEMIC SUBMISSION
# ---------------------------------------------------------------------
# Walk-forward OOS Sharpe values frozen at the December 2025 data lock.
# These are the figures the cohort peer review (June 3) and the panel
# (July 1 / July 3) defend on the record; the Council Performance Record
# page on the platform shows both side by side:
#   Academic Submission (locked Dec 2025):  blend 0.86, benchmark 0.43
#   Live Figure (through May 2026):         blend 1.24, benchmark 0.73
# The platform footnote is authoritative: "The submitted figures are
# the December 2025 data lock... not used in the executive brief or
# final presentation -- the academic submission stands as the record."
# So the brief / deck / appendix OOS HEADLINE cites 0.86 vs 0.43.
#
# The live full-period Sharpes (REGIME_SWITCHING.sharpe_ratio = 0.6291,
# BENCHMARK.sharpe_ratio = 0.5370) appear separately via the per-
# strategy tokens ({{REGIME_SWITCHING_SHARPE}} etc.) which read from
# the strategy cache directly -- they are NOT derived from these
# constants and appear in the full-period performance tables only.
#
# History: PR #370 initially proposed Path A as "update OOS_SHARPE to
# live full-period" (0.6291 / 0.5370). User reverted that part of the
# proposal after surfacing the Performance Record's two-state design.
# Keeping the OOS_SHARPE pair at the December lock preserves the
# panel-defense alignment; the rest of Path A (max drawdown live,
# equal-weight Sharpe live, new OOS window constants) is unchanged.
#
# DO NOT manually edit OOS_SHARPE_REGIME_CONDITIONAL or
# OOS_SHARPE_BENCHMARK without an explicit decision to break the
# December 2025 academic submission record. Updating these in
# response to a fresh data ingestion is a SUBMISSION INTEGRITY
# VIOLATION; the academic record defends these specific values.
OOS_SHARPE_REGIME_CONDITIONAL = 0.8576  # Dec 2025 academic submission
OOS_SHARPE_BENCHMARK = 0.4341          # Dec 2025 academic submission
OOS_SHARPE_EQUAL_WEIGHT = 0.5728       # live EQUAL_WEIGHT.sharpe_ratio
CORRELATION_PRE_2022 = -0.05           # avg of 12m rolling corr pre-2022
CORRELATION_POST_2022 = 0.57           # avg of 12m rolling corr post-2022
PLAY_BY_PLAY_EVENTS = 9                # rebalance_events.csv row count

# Max drawdown over the full study period (Jul 2002 - May 2026).
# BENCHMARK matches the cache exactly (rounds from -0.5256 to -0.526);
# REGIME_SWITCHING bumped from -0.253 to -0.2974 to match the live
# cache REGIME_SWITCHING.max_drawdown under Path A.
MAX_DRAWDOWN_BENCHMARK = -0.526
MAX_DRAWDOWN_REGIME_CONDITIONAL = -0.2974
PLAY_BY_PLAY_ADD_VALUE = 2

# June 22 2026 (Path A scope) -- locked submission figures threaded
# through validated_constants so the brief story plan resolver, the
# deck/appendix substitution tables, and the per-section LLM prompts
# all see them. Prompts that previously hardcoded the literal "40
# months" / "14%" / etc. have been replaced with {{OOS_WINDOW_MONTHS}}
# and {{OOS_WINDOW_PCT_OF_STUDY}} placeholders that resolve here.
OOS_WINDOW_MONTHS = 53                  # Feb 2022 - May 2026 inclusive
OOS_WINDOW_PCT_OF_STUDY = 18.5          # 53 / 287 ≈ 18.5% of full window
CURRENT_REGIME = "BULL"                 # live HMM read at submission lock
CURRENT_EQUITY_WEIGHT = 0.80            # BULL regime equity from
                                        # REGIME_WEIGHTS in
                                        # tools/backtester.py:888
                                        # (BULL = 80% eq / 20% IG)

# Regime state is LIVE-SOURCED, never a constant. The current regime, its HMM
# posterior confidence, and the live blend weights drift as new monthly data
# lands, so they would go stale before the presentation. They are wired into
# the deck context block at generation time from detect_current_regime() (the
# same source the Forward Projection tile and CIO card use). Slide 6's spec
# below instructs the model to read current_regime / regime_confidence /
# blend_weights from the context rather than any baked-in value.

# The real platform universe — confirmed May 28 2026. The earlier-draft
# TLT / DJP / Trend Following / Dynamic Risk Parity names are NOT used.
_STATIC_STRATEGIES = ["CLASSIC_60_40", "EQUAL_WEIGHT", "RISK_PARITY",
                      "BLACK_LITTERMAN"]
_DYNAMIC_STRATEGIES = ["MOMENTUM_ROTATION", "REGIME_SWITCHING", "VOL_TARGETING",
                       "MIN_VARIANCE", "MAX_SHARPE_ROLLING"]

# June 6 2026 — Six-slide rewrite. The previous 10-slide structure split
# the answer across an introduction (slides 1-2), evidence (3-4), method
# (5-6), validation (7-8) and limitations + conclusions (9-10) — a
# structure that delayed the verdict and read as a paper, not a pitch.
# The new structure leads with the answer (slide 1), establishes the
# practical problem (slide 2), shows the evidence (slide 3), discloses
# the honest validation result (slide 4), names the five team-only
# decisions that the platform could not make (slide 5, addresses the
# 3/5 division-of-labor score from the midpoint), and closes with the
# live recommendation (slide 6, the only slide that pulls live data
# every regeneration — regime + blend updated on presentation morning).
#
# The builder always emits exactly these six slides; when the AI JSON
# is missing or unparseable (the test environment, or an LLM outage)
# the slide still appears with its canonical title and a [DATA PENDING]
# body. SLIDE_CHARTS maps a slide number to the deck-chart role the
# builder embeds on it.
# LOCKED TITLES (June 22 2026 -- SO WHAT framing pass).
# Every title answers either "what does this prove?" (answer
# slide) or "what is the question?" (setup slide). The Sonnet
# per-slide writer must use these verbatim or as the
# token-resolved equivalent -- no alternative titles. SLIDE_TITLES
# is a plain Python list (not an f-string) so {{TOKEN}} markers
# stay literal and pass through to the substitution layer.
SLIDE_TITLES = [
    # 1 -- answer slide: states the finding up front
    "Yes -- Regime-Conditional Beats 100% Equity Out-of-Sample",
    # 2 -- structural agenda (no so-what needed)
    "Agenda",
    # 3 -- combined Investment Case slide. June 27 2026: merged
    # what were previously slide 3 ("Three Strategies, One
    # Question" setup) + slide 4 ("The Numbers: 0.86 vs 0.43 OOS
    # verdict") into a single split-panel slide -- IS setup on
    # the left, OOS verdict table on the right. Mirrors Molly's
    # reference deck and shrinks the platform deck from 12 to 11
    # slides to match her structure. The +1 chart-slot offset
    # caused by the agenda insert (PR #287) is now removed.
    "The Investment Case: Setup and OOS Verdict",
    # 4 -- answer slide: the cause of the 2022 correlation break
    "Why Static Allocation Failed in 2022",
    # 5 -- answer slide: the drawdown finding
    "Capital Preservation: Half the Drawdown, Half the Recovery Time",
    # 6 -- answer slide (rhetorical question with the answer)
    "Does It Hold Up Out-of-Sample? Yes.",
    # 7 -- live regime watchpoints; two tokens resolve at gen time
    "Live Regime Signal: {{CURRENT_REGIME}} at {{REGIME_CONFIDENCE}} Confidence",
    # 8 -- intellectual honesty: 2 of 9 failure scenarios
    "What the Model Gets Wrong: 2 of 9",
    # 9 -- AI methodology before the live demo (PR #375 flip)
    "How We Used AI: What Worked and What We Learned",
    # 10 -- live demo (PR #375 flip)
    "Live Demo -- analyticsdesk.app",
    # 11 -- recommendation
    "The Answer: Yes, With Conditions",
]
DECK_SLIDE_COUNT = len(SLIDE_TITLES)

# Slides that MUST emit at least one interpretive bullet (PR 3,
# June 27 2026). Slides 4, 7, 9, 12 are intentionally absent --
# those are table-heavy proof slides where the table fully
# carries the evidence, so empty bullets is acceptable. The
# per-slide LLM call retries ONCE when bullets come back empty
# for any slide in this set; if still empty after retry the
# bullet block is silently skipped (no [DATA PENDING] text).
SLIDES_REQUIRING_BULLETS: frozenset[int] = frozenset(
    {1, 3, 5, 6, 8, 10, 11})

# slide_number → chart role. The role string is resolved to a chart_render
# renderer + arguments in main._render_deck_slide_charts. Slides not listed
# carry no chart -- the slide body renders bullets + optional table only.
#
# June 22 2026 -- 12-slide structure (agenda insert + AI/demo flip).
#
# June 27 2026 -- reconciled with editor_content.DECK_SLIDE_CHART_KEYS
# so the generation-time PPTX export and the editor canvas show the
# SAME chart on each slide. Before reconciliation the two maps
# disagreed on four slides (4, 6, 7, 8): main only had 5 / 6 / 12,
# and slide 6 was 'strategy_comparison_oos_sharpe' here vs
# 'cumulative_returns' in the canvas. The editor canvas is the
# source of truth for what Molly sees + presents, so generation
# follows it. The assertion at the bottom of editor_content.py
# (where both maps are visible together) fires at module load to
# keep the invariant honest.
SLIDE_CHARTS: dict[int, str] = {
    # June 27 2026 -- collapsed to 11 slides + remapped chart
    # assignments to match Molly's reference deck. Slides 3 + 4
    # were merged into one Investment Case slide; chart slots
    # downstream shifted up by -1. After the collapse only slide
    # 4 (Why Static Allocation Failed in 2022) carries a chart;
    # slides 5-11 are tables / cards / verdict panels (rendered
    # in PR B). All slide numbers below match Molly's deck 1:1.
    4: "rolling_correlation",
}


# ── Generation prompt (the per-slide framing) ─────────────────────────────────
# DECK_GENERATION_PROMPT is the preamble carrying the project framing + the
# audience + the story-arc seed. slide_generation_prompt() slices ONE slide's
# spec out of SLIDE_SPECIFICATIONS, prepends the preamble (with the story-arc
# seed parameterised for that slide), and asks for a single-object JSON
# response. Bridge #98 / #100 (June 7 2026) rebuilt the deck from 6 slides
# to 11 and switched generation from the academic-writer harness to a single
# direct Sonnet call per slide (no peer-discussant evaluator, no Gemini, no
# Opus arbiter -- those were leaking peer-review text into the slide output).
DECK_GENERATION_PROMPT = """\
You are generating slide content for an 18-20 minute final investment \
presentation to an academic panel for FNA 670 (Financial Strategies and \
Analytics) at the McColl School of Business, Queens University of Charlotte. \
The audience is Dr. Katerina Panttser and a graduate-level industry panel.

The central question of the project is: \
"Does diversification outperform 100% equity?" Every slide must contribute \
to answering that question with the investable conclusion front and center.

Midpoint feedback was explicit: (a) lead with the answer, not the methodology; \
(b) simplify the strategy set to three (Static / Dynamic / Benchmark) for the \
panel; (c) emphasise the economic story around regime-switching, not the HMM \
math; (d) keep the framing at the executive level; (e) discuss AI use \
explicitly. Forest Capital purchases individual stocks and bonds; ETF \
proxies in this analysis represent asset-class signals only.

The team is:
Bob Thao -- quantitative analysis, regime hypothesis, economic significance \
threshold.
Michael Ruurds -- platform engineering, architecture, OOS window design, \
asset scope decisions.
Molly Murdock -- presentation, validation framework, 9-event play-by-play \
scorecard, peer review.

Using ONLY the data provided in the context block, generate the slide \
content. Do not invent numbers. If a value is missing from context, write \
[DATA PENDING].

CRITICAL: All numeric values in your response must come exactly from the \
data context provided. Do not estimate, interpolate, round beyond two \
decimal places, or substitute figures from prior knowledge. Sharpe ratios, \
max drawdowns, CAGR figures, and correlation coefficients MUST be quoted \
from the data block verbatim or else replaced with [DATA PENDING]. Every \
numeric attribution is verified against the source cache by a post-\
generation audit; a Sharpe attributed to the wrong strategy or a drawdown \
that is not the cache's figure will fail the numeric cross-reference \
check and flag in the audit panel. Numeric accuracy is non-negotiable."""


# Story-arc seed: every per-slide prompt prepends this with the slide-N
# value substituted in. The seed tells the model where it is in the overall
# narrative -- critical for keeping slide-by-slide tone consistent without
# re-reading every other slide's content.
STORY_ARC_SEED = """\
This is slide {slide_number} of {total_slides} in the final FNA 670 \
investment presentation. The story arc is:
  1. The investment question and direct answer.
  2. Agenda -- structural roadmap for the panel.
  3. The Investment Case -- three strategies + OOS verdict in a single \
split-panel slide (IS setup on the left, OOS verdict table on the right).
  4. Why regime-switching works (the 2022 correlation break).
  5. Capital preservation in bear regimes.
  6. Out-of-sample validation (resolves the overfitting concern).
  7. Macro context -- live regime signal + watchpoints.
  8. Limitations and honest assessment (the 2/9 play-by-play).
  9. AI methodology -- what worked, what didn't.
  10. Live platform demo (analyticsdesk.app).
  11. Final recommendation.

You are generating slide {slide_number}: "{slide_title}". Stay in arc \
position -- do not anticipate slides ahead or re-state slides already \
covered. The panel reads the deck in order."""

# The slide specifications. Constants are interpolated; performance numbers for
# the strategy tables are deliberately left for gather_document_data() to
# supply (the prompt instructs the model to pull them from context, never
# hardcode them). Strategy names are the real platform names.
SLIDE_SPECIFICATIONS = f"""\
SLIDE SPECIFICATIONS

SLIDE FORMAT CONSTRAINTS (non-negotiable, apply to every slide):
- BULLET DISCIPLINE: Target 2-3 bullets per slide. Hard ceiling 3
  (2 for slides 4, 6, 7, 8, 9, 12 where a table carries the
  evidence). One strong bullet beats two weak ones. Silence beats
  padding. The slide_plan's max_bullets field is a CEILING, not a
  target -- max_bullets=2 means "no more than 2", never
  "write exactly 2".
- NON-EMPTY BULLETS REQUIREMENT (slides 1, 3, 5, 6, 8, 10, 11 --
  June 27 2026 PR 3): These slides MUST emit at least one
  interpretive bullet, even when the slide content is primarily
  structured (split-panel cards, card grid, feature rows). The
  panel reader's eye flow depends on the bullet narrative even
  when the card or table carries the headline number. Floor for
  these specific slides is 1; ceiling still 2-3. Slides 4, 7, 9,
  12 (table-heavy proof slides where the table fully carries the
  evidence) may emit an empty bullets array -- floor stays zero
  for those four slides only.
- Each bullet must be a "because" or "which means", never a
  "what". The title and the table state the WHAT. Bullets
  interpret. Wrong: "OOS Sharpe of the blend: 0.86" (the title
  already said this). Right: "Nearly double the benchmark's
  risk-adjusted return on data the model never trained on".
- Maximum 12 words per bullet. Fragments only -- no full
  sentences.
- One headline number or stat per slide, referenced explicitly
  in the title or first bullet.
- No sub-bullets under any circumstance.
- Speaker notes carry the full explanation. The slide carries
  the signal.
- BULLET STRING DISCIPLINE: Bullets are PLAIN STRINGS in the
  bullets array. Do NOT prefix any bullet with "- ", "* ",
  "•", or any other marker character. The PPTX renderer adds
  the bullet glyph automatically; a leading "- " renders as
  "• - Foo" on the slide and reads as a defect. Wrong: "- The
  blend outperforms...". Right: "The blend outperforms...".
- TABLE OUTPUT DISCIPLINE (applies to slides with a Required
  table: spec -- slides 4, 7, 9, 12): emit table_data as
  STRUCTURED JSON {{"headers": [...], "rows": [[...]]}}. The
  schema described in each slide's "Required table:" block
  uses pipe-delimited prose to convey the column NAMES;
  translate that schema into the structured object on output.
  Do NOT emit pipe characters in any string. Do NOT place
  the table in the bullets array. Do NOT wrap header or cell
  values in surrounding pipes. Wrong:
    "table_data": "| Strategy | Sharpe |\\n| Blend | 0.86 |"
    "bullets": ["| Strategy | Sharpe |", ...]
    {{"headers": ["| Strategy", "Sharpe |"], ...}}
  Right:
    "table_data": {{"headers": ["Strategy", "Sharpe"],
                    "rows": [["Blend", "0.86"], ...]}}

TOKEN FORMAT NOTE (non-negotiable, apply to every slide):
- {{{{OOS_SHARPE_IMPROVEMENT_PCT}}}} already resolves to a complete
  formatted string including the + prefix AND the % suffix (e.g.
  "+98%"). Do NOT add any surrounding + or % characters around
  this token; they duplicate the formatting already in the token
  value. Wrong: "+{{{{OOS_SHARPE_IMPROVEMENT_PCT}}}}%" (produces
  "++98%%"). Wrong: "{{{{OOS_SHARPE_IMPROVEMENT_PCT}}}}%" (produces
  "+98%%"). Right: write the token alone -- "...an improvement of
  {{{{OOS_SHARPE_IMPROVEMENT_PCT}}}}..." renders as
  "...an improvement of +98%...".
- The same rule applies to any *_PCT token whose value already
  carries a built-in % suffix. If a token's name suggests it is
  already a percentage, treat the resolved value as complete --
  do not append additional formatting characters around it.

June 27 2026 -- collapsed to 11-slide structure (Investment Case slide 3
combines the old "Three Strategies" setup + "The Numbers" OOS verdict
into one split-panel slide to match Molly's reference deck). The deck
opens with the OOS proof point (slide 1), walks the agenda (slide 2),
lays the combined investment case (slide 3), lays remaining evidence
in slides 4-8, explains AI methodology BEFORE the live demo (slide 9),
runs the demo (slide 10), and closes with the investable recommendation
(slide 11). The AI-before-demo order is deliberate: the panel needs
context on how the council works before they watch it operate live.

Slide 1 -- Yes -- Regime-Conditional Beats 100% Equity Out-of-Sample
max_bullets: 3
Message: The regime-conditional blend outperforms the 100% equity \
benchmark on a risk-adjusted basis in the out-of-sample period.
Required bullets (no more than 3, answer-first):
- Yes. The dynamic regime-conditional blend beats the 100% equity \
benchmark on out-of-sample Sharpe.
- OOS Sharpe of the blend: {{strategy_performance.regime_conditional.oos_sharpe}}. \
OOS Sharpe of the benchmark: {{strategy_performance.benchmark.oos_sharpe}}. \
Pull both numbers from context -- do NOT invent.
- Current regime context: state the live regime label from \
current_regime as a single word (BEAR / BULL / TRANSITION).
Required table: none. Single large stat card -- the OOS Sharpe of the \
blend, the benchmark Sharpe, and the percentage advantage.
Chart: none -- the answer is the data point.
Speaker notes: "Before we get to methodology, here is the answer. \
The regime-conditional blend outperforms 100% equity on risk-adjusted \
returns in the out-of-sample period. The rest of the deck is the evidence \
behind that answer." Time: ~1 minute.
Guardrails: No more than 3 data points. Do not mention HMM, strategy \
codes, or factor loadings on this slide. Visible answer in 5 seconds.

Slide 2 -- Agenda
max_bullets: 3
Message: Walk the panel through the deck structure before diving into \
evidence. This slide is purely structural -- no data, no tokens, no \
chart -- so the audience can anchor each upcoming slide to a known \
agenda item.
Required bullets (six items, section labels only, no detail under each):
- The Investment Case -- three strategies, one question
- The Evidence -- OOS performance and drawdown
- Why Static Failed -- the 2022 correlation break
- Out-of-Sample Validation -- what the model never saw
- Honest Limitations -- 2 of 9, what the model gets wrong
- AI Methodology, Live Demo, and Recommendation
Required table: none.
Chart: none.
Speaker notes: "Walk through the agenda briefly. Total presentation is \
18-20 minutes. Panel will have questions after slide 11." \
Time: ~30 seconds.
Guardrails: No data, no substitution tokens, no chart on this slide. \
Pure structure. Do not add detail bullets under each agenda item; the \
panel sees the detail when each section lands. This slide is excluded \
from brief grounding (SLIDES_EXCLUDED_FROM_BRIEF_GROUNDING) because it \
is structural, not analytical.

Slide 3 -- The Investment Case: Setup and OOS Verdict
max_bullets: 3
Layout: Split panel. LEFT (50% width) -- IS setup: three strategy \
stat callout cards stacked vertically, each card carrying the \
strategy name, its IS Sharpe figure, and a "full period" label. The \
Regime-Conditional card is the highlighted card. RIGHT (50% width) -- \
OOS verdict table: Strategy | OOS Sharpe | Max Drawdown | OOS CAGR. The \
Regime-Conditional row is the highlighted row. A "+98% vs benchmark" \
callout sits below the table.
Message: Frame the three strategies the panel will hear about, then \
answer the headline question with the OOS verdict in the same slide. \
Static = Classic 60/40. Dynamic = Regime-conditional blend. Benchmark = \
100% S&P 500. Everything else is supporting evidence. The dynamic \
strategy beats the benchmark on every risk-adjusted metric in the \
out-of-sample period; static diversification also outperforms but by \
less.
Required bullets (left-panel stat cards):
- Static (Classic 60/40): a fixed 60% equity / 40% bond mix, rebalanced \
monthly. IS Sharpe from context.summary_statistics. The traditional \
diversification answer.
- Dynamic (Regime-Conditional Blend): allocation shifts with the live \
regime read -- the platform's recommendation. IS Sharpe from \
context.summary_statistics. Highlighted card.
- Benchmark (100% S&P 500): the question's baseline. IS Sharpe from \
context.summary_statistics. The question is whether diversification \
beats this.
Required table (right panel):
Headers: Strategy | OOS Sharpe | Max Drawdown | OOS CAGR
Rows: Three rows in this exact order: Dynamic Blend (highlighted), \
Classic 60/40, 100% Equity Benchmark. Pull every figure from \
context.strategy_performance. Use plain English strategy names on \
this slide.
Footnote: "Out-of-sample window: {{OOS_WINDOW_MONTHS}}+ months post-2022 \
correlation break. Figures based on December 2025 data lock. Academic \
submission figures."
Chart: none -- the split-panel layout (cards + table) is the visual.
Speaker notes: "We focus on three strategies for this panel: the \
benchmark (100% equity), the static diversifier (Classic 60/40), and \
the dynamic blend (regime-conditional). These three are the submission \
record; the analytical appendix carries the same three strategies at \
higher detail. The cards on the left show the in-sample Sharpe across \
the full study period. The table on the right is the out-of-sample \
verdict -- the dynamic blend leads on Sharpe and drawdown, Classic \
60/40 helps but the edge is smaller, the benchmark is third on every \
risk-adjusted metric." Time: ~3 minutes (combines what were two \
separate slides in the prior deck structure).
Guardrails: Use ONLY the three submission strategies (BENCHMARK, \
CLASSIC_60_40, REGIME_SWITCHING); other strategy codes are out of \
scope. Use plain English on this slide. Do NOT mix IS and OOS in the \
right-panel table -- it is OOS only; the left-panel cards are IS only.

Slide 4 -- Why Static Allocation Failed in 2022
max_bullets: 3
Message: Pre-2022 equities and bonds were negatively correlated -- \
diversification was a free hedge. Post-2022 the correlation flipped to \
+0.68. Regime-conditional allocation adapts; static does not.
Required bullets:
- Equity-bond correlation pre-2022: {CORRELATION_PRE_2022:+.2f} \
(negative -- bonds hedged equity falls).
- Post-2022: {CORRELATION_POST_2022:+.2f} (structural inversion). Bonds \
and equities now fall together.
- Static 60/40 was designed for the pre-2022 regime. It cannot adapt to \
the new correlation environment.
- Dynamic regime-conditional allocation shifts when the regime shifts.
Required table: none -- the chart carries the slide.
Chart: rolling_correlation -- twelve-month equity-bond correlation \
2002-2026, vertical line at Jan 2022, annotation showing the +0.68 \
post-2022 figure. Right-panel bar chart NOT required in the JSON; the \
builder renders the canonical rolling-correlation matplotlib output.
Speaker notes: "Pre-2022 the equity-bond correlation was {CORRELATION_PRE_2022:+.2f}. \
That negative correlation was the diversification answer. In 2022 it \
flipped to {CORRELATION_POST_2022:+.2f}. The static 60/40 answer was \
designed for the old regime. The dynamic blend adapts." Time: ~2 minutes.
Guardrails: Causal story, not just numbers. No HMM technical detail. \
Just the economic intuition: correlation broke, static suffered, regime \
detection adapted.

Slide 5 -- Capital Preservation: Half the Drawdown, Half the Recovery Time
max_bullets: 2  (table-heavy slide; let the table do the work)
Message: The platform's edge is capital preservation when it matters most, \
not bull-market outperformance.
Required bullets:
- The blend's drawdown profile is meaningfully better in bear regimes; \
the benchmark and 60/40 are both punished in 2008-2009 and 2022.
- Max drawdown: blend vs benchmark vs 60/40 -- pull from context.drawdown_comparison.
- Honest caveat: the council MISSED the April 2025 Liberation Day \
V-shaped recovery. Capital preservation is the edge, not crisis \
prediction.
Required table:
Headers: Strategy | Max Drawdown | 2008-09 DD | 2022 DD
Rows: Dynamic Blend, Classic 60/40, 100% Equity Benchmark. Pull figures \
from context.drawdown_comparison. If a value is missing, write \
[DATA PENDING] in that cell.
Chart: strategy_comparison_oos_sharpe -- cumulative-returns proxy, \
three lines, 2008-09 and 2022 bear regions shaded.
Speaker notes: "Where does the dynamic blend earn its Sharpe? In bear \
regimes. The blend cuts the worst drawdown roughly in half versus the \
benchmark. Liberation Day in April 2025 was a miss for us -- we did not \
call the V-shaped recovery. Our edge is sustained directional regimes, \
not sharp reversals." Time: ~2 minutes.
Guardrails: Be honest about the April 2025 miss. Cumulative returns \
chart starts at $1.00 in 2002.

Slide 6 -- Does It Hold Up Out-of-Sample? Yes.
max_bullets: 2  (table-heavy slide; let the table do the work)
Message: Strategy was designed on pre-2022 data and tested on \
{{OOS_WINDOW_MONTHS}}+ months \
of post-2022 data it never saw. It beats the benchmark OOS -- this \
addresses the overfitting concern directly.
Required bullets:
- Design window (in-sample): 2002 through 2021-end. The strategy was \
calibrated on this period.
- Test window (out-of-sample): January 2022 through the May 2026 \
data lock -- ~{{OOS_WINDOW_MONTHS}} months the strategy never saw.
- Dynamic Blend OOS Sharpe: {{strategy_performance.regime_conditional.oos_sharpe}}.
- Benchmark OOS Sharpe: {{strategy_performance.benchmark.oos_sharpe}}.
- OOS performance is genuine -- the strategy parameters were NOT tuned \
on post-2022 data.
Required table:
Headers: Strategy | IS Sharpe (pre-2022) | OOS Sharpe (post-2022)
Rows: Dynamic Blend, Classic 60/40, 100% Equity Benchmark. Pull both \
columns from context.strategy_performance.
Chart: none -- the comparison table IS the slide.
Speaker notes: "The overfitting concern: did we tune on the same data we \
report on? No. Pre-2022 is the design window. Post-2022 is \
{{OOS_WINDOW_MONTHS}}+ months of \
genuine out-of-sample. The dynamic blend beats the benchmark on both. \
The static 60/40 wins IS, ties OOS." Time: ~2 minutes.
Guardrails: Address overfitting head-on. Panel WILL ask.

Slide 7 -- Live Regime Signal: {{{{CURRENT_REGIME}}}} at {{{{REGIME_CONFIDENCE}}}} Confidence
max_bullets: 2  (table-heavy slide; the watchpoint grid is the evidence)
Message: Connect the live platform regime signal to real macroeconomic \
conditions. The current regime classification reflects specific \
observable factors -- not a black box.

Title contract: reproduce the title VERBATIM in your JSON, including \
the {{{{CURRENT_REGIME}}}} placeholder. Do NOT substitute the regime \
label yourself -- the platform's substitution layer resolves the token \
to the live HMM read at generation time so a future regeneration with \
a new regime classification automatically updates the slide title.

Required bullets:
- Live regime: {{current_regime}} ({{regime_confidence}}). The signal \
updates as macro conditions shift.
- VIX, 10Y-2Y spread, HY credit spread, equity trend -- five watch-tiles \
showing what the live signal is reading right now.
- This is the only slide with live data; everything else is the locked \
academic figures from December 2025.
Required table:
Headers: Watchpoint | Current | Direction | Threshold
Rows: five watch-tiles -- pull from context.live_regime_signals. If the \
context lacks a watch-tile, write [DATA PENDING] in that row.
Chart: none -- the watchpoint tile grid carries the slide.
Speaker notes: "The live regime read updates as macro conditions move. \
Here is what the platform is reading right now. {{current_regime}} -- \
{{regime_confidence}}. This is the only slide in the deck with live \
data; the rest are the December 2025 academic-submission figures. \
[PRESENTER REVIEW -- Molly: before the July 1 panel, review the \
contextual event references on this slide (Fed meeting dates, \
geopolitical events, specific macro data points cited in the \
watchpoint descriptions) and update any that are stale. The \
watchpoint VALUES are live and self-updating via the substitution \
layer, but the NARRATIVE CONTEXT around them is from generation \
time -- a 2026-06-22 generation that references the \"recent FOMC \
decision\" will still say that on the July 1 panel even if a newer \
meeting has happened by then.]" Time: ~2 minutes.
Guardrails: This slide is FOR DISCUSSION ONLY. Add the label: \
"Live signal as of [generation date] -- for discussion, not academic \
submission figures."

Slide 8 -- What the Model Gets Wrong: 2 of 9
max_bullets: 2  (table-heavy slide; the failure scenarios table is the evidence)
Message: Intellectual honesty is a strength. Edge is capital preservation \
in sustained bear regimes. NOT designed to call sharp V-shaped reversals. \
Acknowledge explicitly.
Required bullets:
- We tested the council against {PLAY_BY_PLAY_EVENTS} named market events \
committed to the database BEFORE the council was asked about them.
- Honest result: value added in {PLAY_BY_PLAY_ADD_VALUE} of \
{PLAY_BY_PLAY_EVENTS} events.
- Misses include April 2025 Liberation Day (V-shaped recovery -- \
council was risk-off, missed the bounce).
- The regime filter is optimised for SUSTAINED directional regimes, not \
short-duration reversals.
Required table:
Headers: Event | Date | Council Signal | Outcome
Rows: every event in context.play_by_play_events. One row per event. \
Event names only, no quantitative outcome scores -- the slide is a \
scorecard not a P&L attribution.
Chart: none -- the table IS the scorecard.
Speaker notes: "{PLAY_BY_PLAY_ADD_VALUE} of {PLAY_BY_PLAY_EVENTS}. We \
are not selling crisis prediction. We are selling allocation discipline \
that compounds. The model's edge is sustained bear regimes -- not sharp \
reversals like April 2025 Liberation Day." Time: ~1.5 minutes.
Guardrails: Do NOT spin the misses. {PLAY_BY_PLAY_ADD_VALUE}/{PLAY_BY_PLAY_EVENTS} \
honest framing is academically stronger than cherry-picking wins.

Slide 9 -- How We Used AI: What Worked and What We Learned
max_bullets: 3
Message: Rubric explicitly requires discussion of AI use. The panel \
needs context on how the council works BEFORE they watch the live \
demo on slide 10. The generator-evaluator council with dissenting \
agents is the platform differentiator. Honest about what failed.
Required bullets (two columns -- what worked + what we learned):
- What worked: multi-model validation (Sonnet + Opus + Gemini + Grok), \
regime-keyed caching, harness evaluation with the dissenting Risk Manager.
- What worked: deterministic Python recomputation of every reported \
number -- the LLMs draft prose, Python computes the numbers.
- What worked: documentation generation -- this deck, the brief, the \
appendix all from the same data layer.
- What we learned: LLM arithmetic is unreliable -- replaced every \
numerical computation with deterministic Python.
- What we learned: early prompts produced sycophantic outputs -- the \
dissenting Risk Manager agent specifically counters this.
- What we learned: the council is the analytical engine, not just a \
writing tool.
Required table: none -- two-column bullets carry the slide.
Chart: none.
Speaker notes: "The rubric asks about AI use. We used AI critically, \
not blindly. What worked: multi-model validation, deterministic \
recomputation, dissenting agents that argue back. What we learned: LLM \
arithmetic is unreliable, early prompts were sycophantic, and the \
council is the analytical engine -- not just a writing helper. The \
next slide will show this in action live." Time: ~2 minutes.
Guardrails: Honest and reflective, not promotional. Do NOT list every \
AI model used. The callout: "Every number in this presentation was \
verified by deterministic Python recomputation, not LLM arithmetic."

Slide 10 -- Live Demo -- analyticsdesk.app
max_bullets: 3
Message: Live demo. The slide-10 AI methodology context primes the \
panel for what they are about to watch. Show three things on the \
live platform: live regime + CIO recommendation, council output with \
dissenting view, document generation.
Required bullets:
- Live regime detection + CIO recommendation (Investment Outlook page).
- Council output -- five agents, generator-evaluator harness, dissenting \
view from the Risk Manager.
- Document generation -- this deck, the executive brief, the analytical \
appendix all from the same data layer.
- URL for the panel: analyticsdesk.app.
Required table: none.
Chart: none -- this is the live demo slide. The panel watches a live \
browser walkthrough of the platform.
Speaker notes: "The platform is the analytical engine behind everything \
you have seen so far. You just saw how the council works in concept; \
now you see it live. Three things to show: regime detection, the \
council with its dissenting view, and document generation. URL is \
analyticsdesk.app." Time: 3.5 minutes (live demo, not slide content).
Guardrails: Live-demo slide. NO data tables. The demo is the live \
browser pivot.

Slide 11 -- The Answer: Yes, With Conditions
max_bullets: 2  (table-heavy slide; the conditions table is the evidence)
Message: Return to the central question with a direct investable answer. \
Yes diversification outperforms 100% equity -- but only when allocation \
is regime-conditional. Static helps but is insufficient post-2022. \
Recommendation for Forest Capital is the current BEAR regime blend.
Required bullets (no more than 4 -- the conclusion is the slide):
- Question restated: "Does diversification outperform 100% equity?"
- Answer: Yes -- regime-conditional diversification outperforms; static \
helps but is insufficient post-2022.
- Current recommended blend (live, refreshes on regeneration): pull \
{{blend_weights}} from context. Express as implied asset allocation \
(equities X%, bonds Y%) using compute_implied_asset_allocation on the \
weights.
- These are asset-class signals -- Forest Capital purchases individual \
stocks and bonds; ETF proxies represent asset-class direction only.
Required table:
Headers: Asset Class | Allocation
Rows: Equities (compute from blend_weights), Bonds (compute), Cash \
residual if any. Two-decimal percentage values.
Chart: efficient_frontier with the live blend point marked.
Speaker notes: "The answer to the panel's question: yes, diversification \
outperforms 100% equity -- when the allocation is regime-conditional. \
Static 60/40 helps but is not sufficient post-2022. The current BEAR \
regime recommendation is {{blend_weights}}, which works out to roughly \
[implied asset allocation]. Forest Capital purchases individual stocks \
and bonds -- these are asset-class signals." Time: 1 minute (close).
Guardrails: End on the investable conclusion, NOT on limitations or \
caveats. Final image the panel sees should be the answer to the \
question. Academic disclaimer present but visually subordinate.

"""


def deck_generation_prompt() -> str:
    """The full generation task text — preamble plus slide specifications —
    handed to harness_narrative() as the deck-generation instruction."""
    return f"{DECK_GENERATION_PROMPT}\n\n{SLIDE_SPECIFICATIONS}"


def _slice_slide_spec(slide_number: int) -> str:
    """Returns the slide-N section out of SLIDE_SPECIFICATIONS.

    SLIDE_SPECIFICATIONS is a single block of text with "Slide N --"
    headers. We split on those headers and return the requested slide's
    body. Bridge #95 -- per-slide generation needs ONE slide's spec at
    a time so the LLM doesn't have to hold all six in working memory
    AND so each call's output fits comfortably under max_tokens.
    """
    if not (1 <= slide_number <= DECK_SLIDE_COUNT):
        raise ValueError(
            f"slide_number must be 1..{DECK_SLIDE_COUNT}; got {slide_number}")
    text = SLIDE_SPECIFICATIONS
    marker = f"Slide {slide_number} --"
    start = text.find(marker)
    if start == -1:
        raise ValueError(
            f"Slide {slide_number} spec missing from SLIDE_SPECIFICATIONS")
    # The next slide's marker terminates this slide's body. The last
    # slide runs to end-of-string.
    if slide_number < DECK_SLIDE_COUNT:
        end = text.find(f"Slide {slide_number + 1} --", start)
        if end == -1:
            end = len(text)
    else:
        end = len(text)
    return text[start:end].rstrip()


def slide_generation_prompt(slide_number: int) -> str:
    """Bridge #98 / #100 -- per-slide generation prompt. Returns a prompt
    that asks the model to emit a SINGLE slide's JSON object for slide
    {slide_number}. The prompt is sent directly to Sonnet (no harness,
    no evaluator, no Gemini, no Opus arbiter -- those were leaking
    peer-review text into the slide output via the academic-review
    evaluator's feedback retry loop).

    Structure:
      1. DECK_GENERATION_PROMPT preamble (project framing, audience,
         midpoint feedback, team roles).
      2. STORY_ARC_SEED with slide_number / total_slides / slide_title
         substituted. Tells the model where it sits in the 11-slide arc.
      3. The slide's spec block from SLIDE_SPECIFICATIONS.
      4. Explicit single-object JSON contract -- no wrapped list, no
         markdown fences, no preamble.

    Token budget: with one slide's content the response stays at
    ~500-1200 tokens, comfortably under a 2000 max_tokens cap.
    """
    spec = _slice_slide_spec(slide_number)
    title = SLIDE_TITLES[slide_number - 1]
    arc_seed = STORY_ARC_SEED.format(
        slide_number=slide_number,
        total_slides=DECK_SLIDE_COUNT,
        slide_title=title)
    return (
        f"{DECK_GENERATION_PROMPT}\n\n"
        f"{arc_seed}\n\n"
        f"SLIDE SPEC:\n{spec}\n\n"
        f"Output ONLY a JSON object with these keys (no preamble, no "
        f"markdown fences, no wrapping list, no commentary):\n"
        f'{{\n'
        f'  "slide_number": {slide_number},\n'
        f'  "title": "{title}",\n'
        f'  "bullets": ["...", "..."],\n'
        f'  "table_data": null or {{"headers": [...], "rows": [[...], ...]}},\n'
        f'  "speaker_notes": "..."\n'
        f'}}\n\n'
        # June 25 2026 -- table-data discipline.
        "TABLE DISCIPLINE: table_data must always use the structured "
        "{headers: [...], rows: [[...]]} object. NEVER emit pipe-"
        "delimited markdown tables ('| col1 | col2 |\\n| --- | --- |"
        "\\n| a | b |') as bullets or as a string in table_data. The "
        "PPTX renderer expects the structured object; a markdown "
        "string lands in a text frame as raw text and reads as "
        "garbage on the slide.\n\n"
        + _chart_bullet_discipline_for_slide(slide_number)
    )


def _chart_bullet_discipline_for_slide(slide_number: int) -> str:
    """June 25 2026 -- per-slide chart-discipline instruction. Slides
    in SLIDE_CHARTS have an embedded chart image AT GENERATION TIME;
    the LLM doesn't know that and tends to emit bullets that
    describe the chart ('the chart shows...', 'as illustrated in
    Figure 1...', 'see chart at right'). Both the bullet AND the
    image render -- visual duplication. This injects an explicit
    instruction for chart-bearing slides so the bullets stick to
    interpretive context (what the chart MEANS) rather than
    describing it (what it SHOWS). The slide spec's own bullets
    list already covers the desired content; this is reinforcement.
    Non-chart slides get an empty string -- no instruction needed."""
    if slide_number not in SLIDE_CHARTS:
        return ""
    return (
        "CHART DISCIPLINE: This slide has an embedded chart image. "
        "Do not include bullets that describe, reference, or "
        "summarize the chart visually. Bullets on this slide must "
        "provide INTERPRETIVE context only -- what the chart means "
        "for the investment thesis, not what the chart shows. "
        "Forbidden phrases: 'see chart', 'shown above', 'the chart "
        "shows', 'as illustrated', 'depicted in', 'visualized in', "
        "'Figure N shows'. The chart speaks for itself; the bullets "
        "anchor the so-what.\n\n")


# LLM preamble patterns that should never appear in a slide bullet.
# When the model emits a short conversational opener before the JSON
# (which the JSON-extraction step strips) the same conversational tone
# sometimes bleeds into a bullet -- "I cannot generate the requested
# slide", "Note: ...", "Sorry, ...". A short prefix scan after parsing
# replaces those with a clearly-flagged regen marker so the deck never
# carries raw apology text to the audience.
_BULLET_PREAMBLE_PATTERNS = (
    "i ", "i'm ", "i am ", "note:", "sorry", "as an ai",
    "i apologize", "i cannot", "i can't", "unfortunately",
)
_BULLET_PREAMBLE_REPLACEMENT = (
    "[Content generation error -- regenerate deck]")


def _bullet_looks_like_preamble(bullet: str) -> bool:
    """A bullet that opens with a known LLM preamble pattern is the
    model talking ABOUT the slide rather than producing slide content.
    Case-insensitive prefix check; the leading bullet glyph (if any)
    is stripped before the comparison."""
    text = (bullet or "").lstrip("-•▪ \t").lower()
    return any(text.startswith(p) for p in _BULLET_PREAMBLE_PATTERNS)


def _scrub_bullet_preambles(obj: dict) -> dict:
    """Replace any bullet that opens with a known LLM preamble pattern
    with the regen marker. Logs each substitution so a deck with a
    polluted prose surface is visible in Render logs without having to
    open the .pptx.
    """
    bullets = obj.get("bullets")
    if not isinstance(bullets, list):
        return obj
    cleaned: list[str] = []
    for b in bullets:
        if isinstance(b, str) and _bullet_looks_like_preamble(b):
            log.warning(
                "deck_bullet_preamble_replaced",
                bullet_preview=str(b)[:120])
            cleaned.append(_BULLET_PREAMBLE_REPLACEMENT)
        else:
            cleaned.append(b if isinstance(b, str) else str(b))
    obj["bullets"] = cleaned
    return obj


def parse_single_slide_json(raw: str) -> dict | None:
    """Bridge #95 -- parse a per-slide LLM response into a single slide
    dict, or None on any failure (the caller writes a [DATA PENDING]
    placeholder for that slide). Mirrors _parse_deck_slides's tolerance
    of markdown fences and leading/trailing prose, but pulls the
    SINGLE object instead of the wrapped list.

    Hardened June 18 2026:
      * any preamble text before the first '{' is discarded entirely
        before parsing (the JSON-object extraction was lenient enough
        that a stray "Here is the slide:" prefix could survive when
        paired with an unusually-shaped body),
      * on parse failure we log the raw response truncated to 500
        chars at WARNING level so a regression is diagnosable without
        having to reproduce locally,
      * post-parse we scrub any bullet that opens with an LLM
        apology / preamble pattern -- those are the model talking
        ABOUT the slide rather than producing slide content and they
        otherwise leak verbatim into the final .pptx.
    """
    import json
    text = (raw or "").strip()
    if "{" not in text:
        return None
    if text.startswith("```"):
        text = text.strip("`")
        if text[:4].lower() == "json":
            text = text[4:]
    # Discard any preamble BEFORE the first opening brace. This is
    # belt-and-braces -- the slice-by-find below already trims it
    # implicitly -- but doing it explicitly here makes the intent
    # visible at the entry point.
    first_brace = text.find("{")
    if first_brace > 0:
        text = text[first_brace:]
    try:
        start, end = text.find("{"), text.rfind("}")
        if start == -1 or end == -1 or end < start:
            log.warning(
                "deck_slide_parse_no_object_braces",
                raw_preview=(raw or "")[:500])
            return None
        obj = json.loads(text[start:end + 1])
        if not isinstance(obj, dict):
            log.warning(
                "deck_slide_parse_non_object",
                raw_preview=(raw or "")[:500])
            return None
        return _scrub_bullet_preambles(obj)
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "deck_slide_parse_failed",
            error=str(exc),
            raw_preview=(raw or "")[:500])
        return None


# ── Professional navy/white theme — deliberately NOT the platform dark UI ─────
_NAVY = RGBColor(0x1A, 0x2A, 0x4A)
_NAVY_SOFT = RGBColor(0x2D, 0x4A, 0x6B)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_INK = RGBColor(0x1A, 0x1A, 0x2E)
_GREY = RGBColor(0x4A, 0x4A, 0x6A)
_ACCENT = RGBColor(0x1D, 0x4E, 0xD8)
_AMBER = RGBColor(0xB4, 0x53, 0x09)

# June 27 2026 -- Molly reference deck palette (PR B specialized
# renderers). Added alongside the existing _NAVY / _ACCENT block so
# the title slide + split panel + card grid + scorecard + feature
# rows all draw from a single palette source. Kept separate from
# the generic _NAVY/_ACCENT used by the legacy uniform renderer so
# a future palette refresh on Molly's side touches only this block.
_MOLLY_NAVY  = RGBColor(0x1E, 0x27, 0x61)  # title chrome + header bands
_MOLLY_TEAL  = RGBColor(0x02, 0x80, 0x90)  # accent (highlighted row / card)
_TEAL_LIGHT  = RGBColor(0xE8, 0xF4, 0xF8)  # alt card / row background
_NAVY_LIGHT  = RGBColor(0xEE, 0xF0, 0xF8)  # alt card / row background
_GOLD_LIGHT  = RGBColor(0xFF, 0xF8, 0xE1)  # 3rd feature-row background
_INK_DARK    = RGBColor(0x1A, 0x1A, 0x2E)  # body text on light backgrounds
# Council-signal palette (slide 8 scorecard cell_style_hook).
_SIGNAL_RED   = RGBColor(0xFF, 0x6B, 0x6B)  # BEAR
_SIGNAL_AMBER = RGBColor(0xF9, 0xA8, 0x25)  # TRANSITION
_SIGNAL_GREEN = RGBColor(0x4C, 0xAF, 0x50)  # BULL

_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)

_AI_DRAFT_FOOTER = ("AI DRAFT — REQUIRES HUMAN REVIEW · "
                    "first-draft deck; verify every figure before presenting")

# CAVEAT 4 — the verification note added to every slide's speaker notes.
_MOLLY_VERIFY_NOTE = (
    "[MOLLY — VERIFY BEFORE PRESENTING: Confirm the data on this slide "
    "matches the current platform output. Rewrite the talking points in "
    "your own voice. Do not read AI-generated text verbatim during the "
    "presentation.]"
)

# CAVEAT 1 + CAVEAT 5 — the review warning and submission checklist,
# carried in the title slide's speaker notes (the deck's equivalent of
# the .docx header warning box and end-of-document checklist).
_DECK_TITLE_NOTE = (
    "AI DRAFT — REVIEW REQUIRED. This deck was generated by AI from "
    "platform data. Before presenting:\n"
    "CITATIONS: verify every cited source exists and says what is "
    "claimed.\n"
    "STATISTICS: confirm every number against the platform Analytics "
    "page — if a value differs from the screen, the screen wins.\n"
    "YOUR VOICE: rewrite every talking point in your own words; do not "
    "read AI prose verbatim.\n"
    "HALLUCINATIONS: AI can produce plausible but incorrect claims — "
    "read every slide critically.\n\n"
    "SUBMISSION CHECKLIST — before the final deck:\n"
    "- All figures confirmed against the Analytics page\n"
    "- All [[VERIFY]] markers resolved and removed\n"
    "- All speaker notes rewritten in the presenter's own voice\n"
    "- AI DRAFT footer removed from the final version\n"
    "- Academic Review run against the final deck\n\n"
    + _MOLLY_VERIFY_NOTE
)

# matplotlib hex equivalents of the theme — light mode for print/projection.
_MPL_INK = "#1A1A2E"
_MPL_GREY = "#4A4A6A"
_MPL_GRID = "#E2E8F0"
_MPL_ACCENT = "#1D4ED8"
_MPL_SERIES = ["#1D4ED8", "#059669", "#B45309", "#7C3AED", "#DB2777",
               "#0891B2", "#CA8A04", "#15803D", "#9333EA", "#374151"]


# ── Chart rendering (matplotlib, light mode) ──────────────────────────────────

def render_deck_charts(
    data: dict[str, Any], sensitivity: dict[str, Any] | None = None,
    target_key: str | None = None,
    chart_config: dict | None = None,
) -> dict[str, bytes | None]:
    """
    Renders every deck chart to a light-mode PNG. Returns a dict keyed by
    slide role — rolling_correlation, cumulative_returns, risk_return,
    sensitivity, team_activity — with PNG bytes or None when the chart
    cannot be drawn (matplotlib missing, or no source data).

    June 26 2026 -- two new optional kwargs (legacy callers
    unaffected when both default to None):
      target_key:   the chart_key the caller actually wants. ONLY this
                    chart's axes have chart_config applied; the other
                    charts in the returned dict render at their
                    hardcoded defaults. Keeps the per-(target_key,
                    chart_config) cache entry distinct without
                    affecting the other slots.
      chart_config: ChartConfig dict for target_key. Applied via
                    _apply_chart_config (title / axis label / axis
                    bounds). None preserves byte-identical legacy
                    output.
    """
    from tools.chart_render import _apply_chart_config

    charts: dict[str, bytes | None] = {
        "rolling_correlation": None, "cumulative_returns": None,
        "risk_return": None, "sensitivity": None, "team_activity": None,
    }

    def _cfg_for(k: str) -> dict | None:
        """Returns chart_config only when this block matches target_key."""
        return chart_config if target_key == k else None
    try:
        import matplotlib
        matplotlib.use("Agg")  # headless — no display on the server
        import matplotlib.pyplot as plt
    except Exception as exc:  # noqa: BLE001
        log.warning("deck_charts_matplotlib_unavailable", error=str(exc))
        return charts

    def _finish(fig) -> bytes:
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor="white")
        plt.close(fig)
        return buf.getvalue()

    def _style(ax) -> None:
        ax.set_facecolor("white")
        ax.grid(True, color=_MPL_GRID, linewidth=0.7)
        ax.tick_params(colors=_MPL_GREY, labelsize=9)
        for spine in ax.spines.values():
            spine.set_color(_MPL_GRID)
        ax.title.set_color(_MPL_INK)
        ax.xaxis.label.set_color(_MPL_GREY)
        ax.yaxis.label.set_color(_MPL_GREY)

    # ── Rolling correlation — the 2022 regime break ───────────────────────
    try:
        rc = data.get("rolling_correlation") or {}
        pts = rc.get("points") or []
        if pts:
            import pandas as pd
            dates = [pd.to_datetime(p["date"]) for p in pts]
            fig, ax = plt.subplots(figsize=(8, 4.2))
            ax.plot(dates, [p.get("equity_ig") for p in pts],
                    color=_MPL_ACCENT, linewidth=1.6, label="Equity vs IG bonds")
            ax.plot(dates, [p.get("equity_hy") for p in pts],
                    color="#059669", linewidth=1.6, label="Equity vs HY bonds")
            ax.axhline(0, color=_MPL_GREY, linewidth=0.8)
            if rc.get("regime_break"):
                ax.axvline(pd.to_datetime(rc["regime_break"]), color=_AMBER_HEX,
                           linestyle="--", linewidth=1.4, label="2022 regime break")
            ax.set_title("Rolling 12-Month Equity–Bond Correlation", fontsize=11)
            ax.set_ylabel("Correlation")
            ax.legend(fontsize=8, frameon=False)
            _style(ax)
            _apply_chart_config(ax, _cfg_for("rolling_correlation"))
            charts["rolling_correlation"] = _finish(fig)
    except Exception as exc:  # noqa: BLE001
        log.warning("deck_chart_rolling_correlation_failed", error=str(exc))

    # ── Cumulative returns — growth of $1 ─────────────────────────────────
    try:
        cr = data.get("cumulative_returns") or {}
        pts = cr.get("points") or []
        strategies = cr.get("strategies") or []
        if pts and strategies:
            import pandas as pd
            dates = [pd.to_datetime(p["date"]) for p in pts]
            fig, ax = plt.subplots(figsize=(8, 4.2))
            for i, name in enumerate(strategies):
                ax.plot(dates, [p.get(name) for p in pts],
                        color=_MPL_SERIES[i % len(_MPL_SERIES)],
                        linewidth=1.3, label=name)
            ax.set_title("Cumulative Total Return — Growth of $1", fontsize=11)
            ax.set_ylabel("Growth of $1")
            ax.legend(fontsize=6.5, frameon=False, ncol=2)
            _style(ax)
            _apply_chart_config(ax, _cfg_for("cumulative_returns"))
            charts["cumulative_returns"] = _finish(fig)
    except Exception as exc:  # noqa: BLE001
        log.warning("deck_chart_cumulative_returns_failed", error=str(exc))

    # ── Risk-return scatter — strategy volatility vs CAGR ──────────────────
    try:
        results = data.get("strategy_results") or {}
        xs, ys, labels, best = [], [], [], None
        best_sharpe = float("-inf")
        for name, r in results.items():
            vol, cagr = r.get("volatility"), r.get("cagr")
            if isinstance(vol, (int, float)) and isinstance(cagr, (int, float)):
                xs.append(vol * 100)
                ys.append(cagr * 100)
                labels.append(r.get("strategy_name") or name)
                sh = r.get("sharpe_ratio")
                if isinstance(sh, (int, float)) and sh > best_sharpe:
                    best_sharpe, best = sh, len(xs) - 1
        if xs:
            fig, ax = plt.subplots(figsize=(8, 4.4))
            ax.scatter(xs, ys, color=_MPL_ACCENT, s=70, zorder=3)
            for i, lab in enumerate(labels):
                ax.annotate(lab, (xs[i], ys[i]), fontsize=6.5,
                            color=_MPL_GREY, xytext=(4, 4),
                            textcoords="offset points")
            if best is not None:
                ax.scatter([xs[best]], [ys[best]], color=_AMBER_HEX, s=130,
                           zorder=4, label="Highest Sharpe")
                ax.legend(fontsize=8, frameon=False)
            ax.set_title("Risk–Return Profile by Strategy", fontsize=11)
            ax.set_xlabel("Annualised volatility (%)")
            ax.set_ylabel("CAGR (%)")
            _style(ax)
            _apply_chart_config(ax, _cfg_for("risk_return"))
            charts["risk_return"] = _finish(fig)
    except Exception as exc:  # noqa: BLE001
        log.warning("deck_chart_risk_return_failed", error=str(exc))

    # ── Sensitivity — parameter robustness ────────────────────────────────
    try:
        strat_sens = (sensitivity or {}).get("strategies") or []
        if strat_sens:
            fig, ax = plt.subplots(figsize=(8, 4.2))
            for i, s in enumerate(strat_sens[:4]):
                xs = [p.get("value") for p in s.get("points", [])]
                ys = [p.get("sharpe") for p in s.get("points", [])]
                if xs and ys:
                    ax.plot(xs, ys, marker="o", markersize=3,
                            color=_MPL_SERIES[i % len(_MPL_SERIES)],
                            linewidth=1.4, label=s.get("strategy", f"Strategy {i+1}"))
            ax.set_title("Parameter Sensitivity — Sharpe vs Key Parameter",
                         fontsize=11)
            ax.set_ylabel("Sharpe ratio")
            ax.legend(fontsize=7, frameon=False)
            _style(ax)
            charts["sensitivity"] = _finish(fig)
    except Exception as exc:  # noqa: BLE001
        log.warning("deck_chart_sensitivity_failed", error=str(exc))

    # ── Team activity — per-member contribution ───────────────────────────
    try:
        members = (data.get("team_summary") or {}).get("per_member") or []
        if members:
            names = [m.get("user_name") or m.get("user", "—") for m in members]
            councils = [m.get("council_interactions", 0) for m in members]
            reviews = [m.get("academic_review_sessions", 0) for m in members]
            views = [m.get("page_views", 0) for m in members]
            fig, ax = plt.subplots(figsize=(7.5, 4.0))
            import numpy as np
            x = np.arange(len(names))
            ax.bar(x - 0.25, councils, 0.25, color=_MPL_ACCENT, label="Council")
            ax.bar(x, reviews, 0.25, color="#059669", label="Academic reviews")
            ax.bar(x + 0.25, views, 0.25, color="#B45309", label="Page views")
            ax.set_xticks(x)
            ax.set_xticklabels(names, fontsize=8)
            ax.set_title("Team Platform Engagement", fontsize=11)
            ax.legend(fontsize=8, frameon=False)
            _style(ax)
            charts["team_activity"] = _finish(fig)
    except Exception as exc:  # noqa: BLE001
        log.warning("deck_chart_team_activity_failed", error=str(exc))

    return charts


# matplotlib needs a hex string for the amber accent.
_AMBER_HEX = "#B45309"


# ── Slide primitives ──────────────────────────────────────────────────────────

def _blank(prs: Presentation):
    """A fully blank slide — every element is drawn manually for theme control."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color: RGBColor) -> None:
    """Fills the whole slide background with a solid colour."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, _SLIDE_W, _SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def _textbox(slide, left, top, width, height, text, *, size=18, bold=False,
             color=_INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Adds a single-paragraph textbox."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _title_bar(slide, text: str) -> None:
    """A navy band across the top of a content slide with the slide title."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _SLIDE_W, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = _WHITE


def _bullets(slide, items: list[str], *, left=Inches(0.7), top=Inches(1.4),
             width=Inches(12.0), height=Inches(5.4), size=18) -> None:
    """A bulleted list. Items already prefixed with '- ' are de-prefixed."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text = item.lstrip("-• ").strip()
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        bullet = p.add_run()
        bullet.text = "▪  "
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = _ACCENT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = _INK


def _image(slide, png: bytes | None, *, left, top, width, fallback="chart"):
    """Places a PNG, or a clearly-marked unavailable note when the
    chart renderer returned None.

    Hardened June 18 2026: every silent None now emits a WARNING log
    naming the chart slot, so a deck that came back with a missing
    chart is diagnosable from Render logs alone (previously the
    [DATA PENDING] text was the ONLY signal, requiring the operator
    to open the .pptx to know a chart had been skipped). The
    placeholder text was also reworded to name the chart slot
    explicitly -- "[Chart unavailable: rolling_correlation -- ..."
    is more actionable than the old "[DATA PENDING] -- chart
    unavailable" line which gave no slot context.
    """
    if png:
        slide.shapes.add_picture(io.BytesIO(png), left, top, width=width)
        return
    # June 27 2026 -- spec: NO placeholder text in exported PPTX.
    # Log + return without rendering anything in the chart slot.
    # The slide's other elements still render around the empty area.
    log.warning(
        "deck_chart_slot_unavailable",
        chart=str(fallback),
        note=("chart PNG missing; skipping render per spec (no "
              "[Chart unavailable] placeholder). Operator: warm "
              "analytics caches + regenerate."))


def _callout(slide, left, top, width, height, heading, value, *,
             color=_ACCENT) -> None:
    """A small coloured callout box — a label over a large value."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = heading
    r1.font.size = Pt(12)
    r1.font.color.rgb = _WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = _WHITE


def _table(slide, headers: list[str], rows: list[list[str]], *,
           left=Inches(0.7), top=Inches(1.4), width=Inches(12.0),
           max_rows=11,
           cell_style_hook=None) -> None:
    """A light-styled native PowerPoint table. [DATA PENDING] when empty.

    June 27 2026 -- cell_style_hook (PR B). When supplied, the hook
    is called for EVERY body cell with kwargs:

      hook(row_idx, col_idx, header, value, row_values) -> dict | None

    Returns a style-overrides dict OR None (use defaults). Supported
    override keys:

      'fill':      RGBColor for cell fill
      'text':      RGBColor for the run's font color
      'bold':      bool for the run's font weight
      'row_fill':  RGBColor that overrides the alternating-row fill
                   for the ENTIRE row (applied to every cell on the
                   row when ANY cell in the row returns it; the first
                   non-None row_fill wins per row).

    The hook receives the full row_values tuple so it can look up
    sibling-column values (e.g. inspect 'Council Signal' to pick the
    badge color for the signal column, OR shade the entire row based
    on an 'Outcome' column). Returning None for a cell uses the
    default header / alternating-row styling.

    The hook signature is keyword-friendly: the caller's hook can
    accept **kwargs to ignore fields it doesn't need."""
    if not rows:
        # June 27 2026 -- spec: [DATA PENDING] must never appear
        # in an exported PPTX. Log + return without rendering;
        # the slide's other elements (bullets / chart / title)
        # still render around the missing table area.
        log.warning(
            "deck_table_rows_empty_skipping",
            note=("table_data missing or empty; skipping the table "
                  "render per spec (no [DATA PENDING] placeholder). "
                  "Operator: warm analytics caches + regenerate."))
        return
    rows = rows[:max_rows]
    n_rows, n_cols = len(rows) + 1, len(headers)
    height = Inches(0.4 * n_rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = tbl_shape.table
    for c, label in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _NAVY
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = label
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = _WHITE
    for r, row in enumerate(rows, start=1):
        # First pass: ask the hook for every cell, collect overrides
        # + look for a row_fill that wins for the whole row.
        cell_overrides: list[dict | None] = []
        row_fill: "RGBColor | None" = None
        for c in range(n_cols):
            if c >= len(row):
                cell_overrides.append(None)
                continue
            ov = None
            if cell_style_hook is not None:
                try:
                    ov = cell_style_hook(
                        row_idx=r, col_idx=c,
                        header=headers[c] if c < len(headers) else "",
                        value=row[c],
                        row_values=tuple(row))
                except Exception as exc:  # noqa: BLE001
                    log.warning(
                        "deck_table_cell_style_hook_failed",
                        error=str(exc), row=r, col=c)
                    ov = None
            if (isinstance(ov, dict) and row_fill is None
                    and ov.get("row_fill") is not None):
                row_fill = ov["row_fill"]
            cell_overrides.append(ov if isinstance(ov, dict) else None)
        default_row_fill = (
            _WHITE if r % 2 else RGBColor(0xF1, 0xF5, 0xF9))
        effective_row_fill = (
            row_fill if row_fill is not None else default_row_fill)
        # Second pass: write the cells with overrides applied.
        for c, value in enumerate(row):
            if c >= n_cols:
                break
            cell = table.cell(r, c)
            cell.fill.solid()
            ov = cell_overrides[c] if c < len(cell_overrides) else None
            cell.fill.fore_color.rgb = (
                ov.get("fill") if (ov and ov.get("fill") is not None)
                else effective_row_fill)
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = str(value)
            run.font.size = Pt(10)
            run.font.color.rgb = (
                ov.get("text") if (ov and ov.get("text") is not None)
                else _INK)
            if ov and ov.get("bold"):
                run.font.bold = True


def _footer(slide, idx: int, total: int) -> None:
    """The AI DRAFT footer plus a slide counter, on every slide."""
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05),
                                   Inches(12.5), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{_AI_DRAFT_FOOTER}   ·   Slide {idx}/{total}"
    run.font.size = Pt(8)
    run.font.italic = True
    run.font.color.rgb = _AMBER


# ── Deck builder ──────────────────────────────────────────────────────────────

_DATA_PENDING_BULLET = (
    "[DATA PENDING] — generated from live analytics; warm the caches "
    "and regenerate the deck.")


def _normalize_slides(raw_slides: Any) -> list[dict[str, Any]]:
    """Map the AI's slide list onto exactly DECK_SLIDE_COUNT slides indexed
    1..10, each with a title. A missing or invalid slide keeps its canonical
    title and a [DATA PENDING] body, so the deck always has all ten slides
    even when the generation JSON is absent (test env / LLM outage)."""
    by_num: dict[int, dict[str, Any]] = {}
    for s in (raw_slides or []):
        if not isinstance(s, dict):
            continue
        n = s.get("slide_number")
        if isinstance(n, int) and 1 <= n <= DECK_SLIDE_COUNT and n not in by_num:
            by_num[n] = s
    out: list[dict[str, Any]] = []
    for n in range(1, DECK_SLIDE_COUNT + 1):
        s = dict(by_num.get(n) or {})
        s["slide_number"] = n
        if not str(s.get("title") or "").strip():
            s["title"] = SLIDE_TITLES[n - 1]
        bullets = [str(b).strip() for b in (s.get("bullets") or [])
                   if str(b).strip()]
        # June 27 2026 -- the legacy [_DATA_PENDING_BULLET] fallback
        # is gone (per spec: [DATA PENDING] must NEVER appear in an
        # exported PPTX). Leave bullets empty when the LLM didn't
        # emit any; the renderer logs + skips the bullet block.
        s["bullets"] = bullets
        out.append(s)
    return out


def _parse_markdown_table(s: str) -> tuple[list[str], list[list[str]]]:
    """June 25 2026 -- defence-in-depth pre-parser for the case
    where the LLM ignores the structured-table_data prompt and
    emits a pipe-delimited markdown table as a string. Returns
    (headers, rows) or ([], []) if the input can't be parsed.

    Recognised shapes:
      '| Col1 | Col2 |\\n| --- | --- |\\n| a | b |\\n| c | d |'
      '| Col1 | Col2 |\\n| a | b |'   (no separator row)

    Surrounding whitespace + leading/trailing pipes are stripped
    per cell. The separator row ('|---|---|' style) is detected by
    cell content being only dashes/colons/spaces + dropped."""
    if not s or "|" not in s:
        return [], []
    raw_lines = [ln.strip() for ln in s.splitlines() if ln.strip()]
    rows: list[list[str]] = []
    for ln in raw_lines:
        if not ln.startswith("|"):
            continue
        cells = [c.strip() for c in ln.strip("|").split("|")]
        # Separator row -- every cell is dashes/colons/spaces.
        if cells and all(
                c == "" or set(c) <= set("-: ") for c in cells):
            continue
        rows.append(cells)
    if len(rows) < 2:
        return [], []
    return rows[0], rows[1:]


_BULLET_MARKER_RE = re.compile(r"^\s*[-*•‣◦·]\s+")


def _strip_bullet_marker(bullet: str) -> str:
    """June 25 2026 -- defensive cleanup for bullet strings the
    LLM prefixed with a marker character. The PPTX renderer adds
    the bullet glyph automatically; a leading '- Foo' would
    render as '• - Foo' on the slide. Strips leading '-', '*',
    '•' (U+2022), '‣' (U+2023), '◦' (U+25E6), '·' (U+00B7) plus
    surrounding whitespace. Idempotent."""
    if not isinstance(bullet, str):
        return str(bullet)
    return _BULLET_MARKER_RE.sub("", bullet, count=1).strip()


def _strip_cell_pipes(cell: str) -> str:
    """June 25 2026 -- defensive cleanup for table cells the LLM
    wrapped in pipes (e.g. emitting {'headers': ['| Strategy',
    'Sharpe |']} instead of {'headers': ['Strategy', 'Sharpe']}).
    Strips leading/trailing '|' + surrounding whitespace."""
    if not isinstance(cell, str):
        return str(cell)
    return cell.strip().strip("|").strip()


# June 26 2026 -- explicit pipe-row regex used by the
# bullets-as-table fallback. Match after stripping any leading
# bullet marker so "- | A | B |" qualifies the same as "| A | B
# |". The trailing pipe is required so a single 'foo | bar'
# sentence with two stray pipes doesn't false-positive.
_PIPE_ROW_RE = re.compile(r"^\s*\|.*\|\s*$")


def _lift_pipe_table_from_bullets(
    bullets: list[str],
    slide_idx: int | None = None,
) -> tuple[list[str], list[list[str]], list[str]]:
    """June 25 2026 (hardened June 26 2026) -- bullets-as-table
    fallback. The prompt forbids placing a table in the bullets
    array, but a model that ignores the instruction emits each
    table row as a SEPARATE bullet item in the bullets list. This
    function scans ACROSS the bullets list (not within any single
    bullet) for 3+ consecutive items that look like pipe-table
    rows, lifts the run via _parse_markdown_table, and returns
    (headers, rows, remaining_bullets) where remaining_bullets
    has the table lines removed so they don't ALSO render as
    garbage prose in the text frame.

    Per-item match: each bullet is FIRST run through
    _strip_bullet_marker (so '- | A | B |' becomes '| A | B |')
    and THEN matched against _PIPE_ROW_RE which requires the line
    to both start AND end with '|'. A row joined with markers
    stripped goes into _parse_markdown_table.

    Logs at WARNING level when the lift fires so a future run can
    confirm whether the prompt-side discipline is being honoured
    or the renderer fallback is doing the work. Pass slide_idx
    to scope the warning to a specific slide number; None when
    the caller has no slide context.

    Returns ([], [], bullets) unchanged when no pipe-table run is
    detected (the common case once the prompt is being followed)."""
    if not bullets:
        return [], [], bullets

    # Per-bullet normalisation: strip any leading bullet marker
    # ("- ", "* ", "•", etc.) so a row formatted as "- | A | B |"
    # matches the same pipe-row pattern as a bare "| A | B |".
    # Idempotent -- already-clean rows pass through unchanged.
    normalised = [_strip_bullet_marker(str(b)) for b in bullets]

    def _is_pipe_row(s: str) -> bool:
        if not s.strip():
            return False
        return bool(_PIPE_ROW_RE.match(s))

    n = len(normalised)
    # Scan for the longest consecutive run of pipe-row bullets.
    best_start, best_len = -1, 0
    i = 0
    while i < n:
        if _is_pipe_row(normalised[i]):
            j = i
            while j < n and _is_pipe_row(normalised[j]):
                j += 1
            run_len = j - i
            if run_len > best_len:
                best_start, best_len = i, run_len
            i = j
        else:
            i += 1
    if best_len < 3:
        return [], [], bullets

    # Reconstruct a markdown table block from the run of
    # normalised pipe rows and parse it. Joining with newlines
    # gives _parse_markdown_table the same shape as a single
    # multi-line table_data string would have.
    block = "\n".join(normalised[best_start:best_start + best_len])
    headers, rows = _parse_markdown_table(block)
    if not headers or not rows:
        return [], [], bullets

    log.warning(
        "deck_bullets_as_table_fallback_fired",
        slide_idx=slide_idx,
        bullets_lifted=best_len,
        header_count=len(headers),
        row_count=len(rows),
        hint=(
            "The LLM emitted table rows as separate bullets "
            "instead of structured table_data. Prompt-side "
            "TABLE OUTPUT DISCIPLINE (academic_deck.py) is not "
            "being honoured for this slide. Renderer fallback "
            "lifted the run into a real PPTX table; consider "
            "investigating whether the prompt needs further "
            "tightening or the model is regressing."))

    # Remove the table bullets from the surviving list. Operate
    # on the ORIGINAL bullets (not normalised) so any non-table
    # bullets that survive aren't stripped of their markers.
    remaining = (
        bullets[:best_start] + bullets[best_start + best_len:])
    return headers, rows, remaining


def _slide_table(
    sl: dict[str, Any],
    slide_idx: int | None = None,
) -> tuple[list[str], list[list[str]], list[str] | None]:
    """Extract (headers, rows-of-strings, bullets_to_use) from a
    slide's table_data + bullets. Returns ([], [], None) when the
    slide carries no usable table -- the third element is None to
    signal 'no bullet edit needed' (the caller keeps the slide's
    original bullets list); when non-None it's the bullets list
    minus any rows that were lifted from a pipe-table-in-bullets
    fallback below.

    June 25 2026 -- three sources tried in order:
      1. table_data dict with structured {headers, rows} (canonical
         prompt-compliant shape)
      2. table_data string containing pipes (the LLM-emitted
         markdown-table string; parsed via _parse_markdown_table)
      3. bullets list with a 3+ consecutive pipe-row run (the LLM
         dumped the table into the bullets array; lifted via
         _lift_pipe_table_from_bullets and the lifted bullets are
         removed from the surviving bullets list so they don't
         ALSO render as text-frame garbage)

    Every cell value is run through _strip_cell_pipes for the
    case where the LLM wrapped headers / cells in surrounding
    pipes (e.g. {'headers': ['| Strategy', 'Sharpe |']})."""
    td = sl.get("table_data")
    # Source 1: structured dict shape.
    if isinstance(td, dict):
        headers = [_strip_cell_pipes(h)
                   for h in (td.get("headers") or [])]
        rows = [[_strip_cell_pipes(c) for c in r]
                for r in (td.get("rows") or [])
                if isinstance(r, (list, tuple))]
        if headers and rows:
            return headers, rows, None
    # Source 2: pipe-delimited string in table_data.
    if isinstance(td, str) and "|" in td:
        headers, rows = _parse_markdown_table(td)
        if headers and rows:
            return (
                [_strip_cell_pipes(h) for h in headers],
                [[_strip_cell_pipes(c) for c in r] for r in rows],
                None,
            )
    # Source 3: 3+ consecutive pipe-row bullets.
    bullets = sl.get("bullets") or []
    if isinstance(bullets, list):
        headers, rows, remaining = _lift_pipe_table_from_bullets(
            [str(b) for b in bullets], slide_idx=slide_idx)
        if headers and rows:
            return (
                [_strip_cell_pipes(h) for h in headers],
                [[_strip_cell_pipes(c) for c in r] for r in rows],
                remaining,
            )
    return [], [], None


# Chart-description phrases that should be stripped from bullets on
# slides that already carry an embedded chart image. The prompt
# forbids them (see _chart_bullet_discipline_for_slide); this is
# defence-in-depth for when the model ignores the instruction.
_CHART_REFERENCE_PHRASES: tuple[str, ...] = (
    "see chart", "shown above", "shown below", "the chart shows",
    "as illustrated", "depicted in", "visualized in",
    "figure 1 shows", "figure 2 shows", "figure 3 shows",
    "figure 4 shows", "the figure shows", "(see chart)",
)


def _strip_chart_bullets(
    bullets: list[str], slide_number: int,
) -> list[str]:
    """Remove bullets that describe a chart on slides whose
    slide_number is in SLIDE_CHARTS (the slide already has the
    chart image embedded; the bullet would be a duplicate).
    Lower-cases the bullet for substring detection so case
    variations don't slip through. Non-chart slides pass through
    unchanged."""
    if slide_number not in SLIDE_CHARTS:
        return bullets
    keep: list[str] = []
    for b in bullets:
        lowered = (b or "").lower()
        if any(p in lowered for p in _CHART_REFERENCE_PHRASES):
            continue
        keep.append(b)
    return keep


# ── PR B (June 27 2026) specialized per-slide renderers ─────────────────
#
# Five renderers add Molly-reference layouts to the deck builder. The
# dispatch lives at the top of _render_content_slide; the generic
# bullets / chart / table renderer continues to serve slides 2 and 5
# (Agenda + Capital Preservation -- intentionally uniform), so the
# legacy code path is preserved verbatim for those.
#
# All renderers take the same signature as _render_content_slide
# (prs, sl, chart_png, idx, total) so the dispatch can call any of
# them uniformly. Each is fully guarded -- one bad slide degrades
# to a [DATA PENDING] placeholder rather than crashing the deck.


def _card(slide, *, left, top, width, height,
          bg_color, header_text, body_text,
          highlight: bool = False) -> None:
    """Draw a single rounded card. Used by _render_card_grid_slide
    and _render_split_panel_slide's left-panel stat cards.

    highlight=True swaps the bg to _MOLLY_TEAL with white text and
    a thicker title -- used for the Regime-Conditional card so the
    panel reads the recommended choice at a glance."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = (
        _MOLLY_TEAL if highlight else bg_color)
    rect.line.fill.background()
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.14)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_color = _WHITE if highlight else _MOLLY_NAVY
    body_color = _WHITE if highlight else _INK_DARK
    p1 = tf.paragraphs[0]
    p1.space_after = Pt(4)
    r1 = p1.add_run()
    r1.text = header_text or ""
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = title_color
    if body_text:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body_text
        r2.font.size = Pt(13)
        r2.font.color.rgb = body_color


def _render_title_slide(prs, sl, chart_png, idx, total) -> None:
    """Slide 1 -- title chrome. Navy header bar full width carrying
    the deck title. Centered title large font. Subtitle line for the
    course / program. Presenter line for the team. No bullet content
    area. Per PR B spec."""
    s = _blank(prs)
    _bg(s, _WHITE)
    try:
        # Navy header bar full width (taller than the generic
        # _title_bar's 1.0" so the title sits in a true title-slide
        # band, not the content-slide chrome).
        bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, _SLIDE_W, Inches(3.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _MOLLY_NAVY
        bar.line.fill.background()
        bar.shadow.inherit = False
        tf = bar.text_frame
        tf.margin_left = Inches(0.8)
        tf.margin_right = Inches(0.8)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = sl.get("title") or SLIDE_TITLES[idx - 1]
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = _WHITE

        # Teal accent rule under the title bar.
        rule = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0),
            _SLIDE_W, Inches(0.08))
        rule.fill.solid()
        rule.fill.fore_color.rgb = _MOLLY_TEAL
        rule.line.fill.background()
        rule.shadow.inherit = False

        # Subtitle: course / program. Centered.
        _textbox(s, Inches(0.8), Inches(3.7),
                 _SLIDE_W - Inches(1.6), Inches(0.6),
                 "Forest Capital  /  McColl School of Business",
                 size=20, color=_MOLLY_NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)

        # Presenter line.
        _textbox(s, Inches(0.8), Inches(4.6),
                 _SLIDE_W - Inches(1.6), Inches(0.6),
                 "Group 1: Bob Thao, Michael Ruurds, Molly Murdock",
                 size=16, color=_INK_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
    except Exception as exc:  # noqa: BLE001
        log.warning("deck_title_slide_failed", error=str(exc))
    _footer(s, idx, total)


def _render_split_panel_slide(
        prs, sl, chart_png, idx, total) -> None:
    """Slide 3 -- The Investment Case split panel.

    LEFT 50%: three strategy stat callout cards stacked vertically.
    Each card: strategy name bold, IS Sharpe large, "full period"
    small label. The Regime-Conditional card is highlighted in
    _MOLLY_TEAL with white text.

    RIGHT 50%: OOS results table (Strategy / OOS Sharpe / Max
    Drawdown / OOS CAGR) with the navy header band + alternating
    rows + the Regime-Conditional row highlighted teal.

    "+98% vs benchmark" callout below the table.

    Data source: the slide's bullets array carries the three IS
    setup lines (one per strategy); the slide's table_data carries
    the OOS verdict table. Bullets shape:
      ["Strategy Name | IS Sharpe", ...] -- the '|' separates
      the card header from the body. When the LLM emits plain
      prose bullets the renderer still draws them; the highlight
      pass keys off the Regime-Conditional string."""
    s = _blank(prs)
    _bg(s, _WHITE)
    title = sl.get("title") or SLIDE_TITLES[idx - 1]
    try:
        _title_bar(s, title)

        # ── LEFT PANEL (50%) -- stat cards ──────────────────────
        left_x = Inches(0.5)
        card_w = Inches(6.0)
        cards_top = Inches(1.4)
        card_h = Inches(1.7)
        card_gap = Inches(0.15)
        # June 27 2026 -- spec: NEVER emit [DATA PENDING].
        # When the LLM returned no bullets, log + skip the
        # stat-card panel entirely; the right-panel OOS table
        # still renders. Bullets that exist get rendered; we no
        # longer pad to 3 with placeholders.
        raw_bullets = sl.get("bullets") or []
        bullets = [_strip_bullet_marker(b) for b in raw_bullets]
        bullets = [b for b in bullets if b.strip()][:3]
        if not bullets:
            log.warning(
                "deck_split_panel_bullets_empty_skipping",
                slide=idx,
                note=("LLM emitted no bullets for the slide 3 "
                      "stat-card panel; skipping the panel + "
                      "rendering only the OOS table per spec "
                      "(no [DATA PENDING] placeholder)"))
        for i, raw in enumerate(bullets):
            # Card header / body split on '|'; fall back to a
            # heuristic when the bullet is a single sentence
            # (split on the first colon, or use the bullet as
            # body with a numbered fallback header).
            parts = raw.split("|", 1)
            if len(parts) == 2:
                header_text = parts[0].strip()
                body_text = parts[1].strip()
            elif ":" in raw:
                header_text, body_text = raw.split(":", 1)
                header_text = header_text.strip()
                body_text = body_text.strip()
            else:
                header_text = f"Strategy {i + 1}"
                body_text = raw.strip()
            highlight = (
                "regime" in header_text.lower()
                or "regime" in body_text.lower())
            bg = _TEAL_LIGHT if i % 2 == 0 else _NAVY_LIGHT
            top = Inches(1.4 + i * (1.7 + 0.15))
            _card(s, left=left_x, top=top,
                  width=card_w, height=card_h,
                  bg_color=bg,
                  header_text=header_text,
                  body_text=body_text,
                  highlight=highlight)

        # ── RIGHT PANEL (50%) -- OOS results table ─────────────
        headers, rows, _surviving = _slide_table(sl, slide_idx=idx)
        right_x = Inches(6.85)
        right_w = Inches(6.0)
        table_top = Inches(1.4)
        if headers and rows:
            # Highlight the Regime-Conditional row teal.
            def _row_hook(*, row_idx, col_idx, header,
                          value, row_values, **_):
                row_text = " ".join(
                    str(v).lower() for v in row_values)
                if ("regime" in row_text
                        or "dynamic" in row_text):
                    return {
                        "row_fill": _MOLLY_TEAL,
                        "text": _WHITE,
                        "bold": True,
                    }
                return None
            _table(s, headers, rows,
                   left=right_x, top=table_top,
                   width=right_w, max_rows=8,
                   cell_style_hook=_row_hook)
        else:
            # June 27 2026 -- spec: [DATA PENDING] must never
            # appear. Log + render only the left panel + callout.
            log.warning(
                "deck_split_panel_table_empty_skipping",
                slide=idx,
                note=("OOS verdict table_data missing; rendering "
                      "left panel + callout only per spec (no "
                      "[DATA PENDING] placeholder)"))

        # "+98% vs benchmark" callout below the table.
        callout_top = Inches(5.5)
        callout = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            right_x, callout_top, right_w, Inches(0.95))
        callout.fill.solid()
        callout.fill.fore_color.rgb = _MOLLY_TEAL
        callout.line.fill.background()
        callout.shadow.inherit = False
        ctf = callout.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = Inches(0.18)
        ctf.margin_right = Inches(0.18)
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = (
            "+98% Sharpe improvement vs benchmark "
            "across the OOS window")
        cr.font.size = Pt(16)
        cr.font.bold = True
        cr.font.color.rgb = _WHITE
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "deck_split_panel_slide_failed", error=str(exc))
    _footer(s, idx, total)


def _render_card_grid_slide(
        prs, sl, chart_png, idx, total) -> None:
    """Slides 4, 6, 7, 9, 11 -- two or three cards side by side.

    Each card: alternating _TEAL_LIGHT / _NAVY_LIGHT background,
    bold header, body text. For slide 4 (Why Static Failed) the
    layout splits 50/50: cards on the left, rolling_correlation
    chart on the right. The other card-grid slides use the full
    width.

    Card content comes from the slide's bullets array; same '|' /
    ':' split convention as the split-panel left panel."""
    s = _blank(prs)
    _bg(s, _WHITE)
    title = sl.get("title") or SLIDE_TITLES[idx - 1]
    try:
        _title_bar(s, title)
        # June 27 2026 -- spec: NEVER emit [DATA PENDING].
        # When the LLM returned no bullets, log + skip the card
        # area; the title bar + chart placeholder (if any) still
        # render. No placeholder text under any circumstances.
        raw_bullets = sl.get("bullets") or []
        bullets = [_strip_bullet_marker(b) for b in raw_bullets]
        bullets = [b for b in bullets if b.strip()][:3]
        if not bullets:
            log.warning(
                "deck_card_grid_bullets_empty_skipping",
                slide=idx,
                note=("LLM emitted no bullets for card-grid slide; "
                      "skipping the card area + rendering the chart "
                      "placeholder only per spec (no [DATA PENDING] "
                      "placeholder)"))
            # Render the chart slot then return -- no cards to draw.
            if idx in SLIDE_CHARTS:
                _image(s, chart_png,
                       left=Inches(0.6), top=Inches(1.5),
                       width=Inches(12.1),
                       fallback=SLIDE_CHARTS.get(idx, "chart"))
            _footer(s, idx, total)
            return

        # A slide carries a chart slot when SLIDE_CHARTS pins one,
        # whether or not the PNG actually arrived. _image still
        # drops a "[Chart unavailable: ...]" placeholder when the
        # PNG is None so the slot stays visible. Cards fit in the
        # left 50% when a slot is reserved; otherwise they spread
        # horizontally.
        has_chart_slot = idx in SLIDE_CHARTS
        cards_x = Inches(0.5)
        cards_top = Inches(1.4)
        cards_h = Inches(5.4)
        if has_chart_slot:
            cards_w = Inches(6.0)
        else:
            cards_w = Inches(12.33)
        card_count = len(bullets)
        # Cards stack vertically inside the cards_w column when a
        # chart slot is reserved (3 cards * 1.7" each + gaps fits
        # within 5.4"); otherwise lay out horizontally to fill the
        # width.
        if has_chart_slot:
            card_h = Inches(1.7)
            card_gap = 0.15
            for i, raw in enumerate(bullets):
                hdr, body = _split_card_line(raw, i + 1)
                bg = (
                    _TEAL_LIGHT if i % 2 == 0 else _NAVY_LIGHT)
                top = Inches(
                    1.4 + i * (1.7 + card_gap))
                _card(s, left=cards_x, top=top,
                      width=cards_w, height=card_h,
                      bg_color=bg,
                      header_text=hdr, body_text=body)
            _image(s, chart_png,
                   left=Inches(6.85), top=Inches(1.5),
                   width=Inches(6.0),
                   fallback=SLIDE_CHARTS.get(idx, "chart"))
        else:
            # Horizontal layout: divide cards_w equally.
            gap = Inches(0.25)
            total_gap = gap * (card_count - 1)
            card_w = (cards_w - total_gap) // max(1, card_count)
            for i, raw in enumerate(bullets):
                hdr, body = _split_card_line(raw, i + 1)
                bg = (
                    _TEAL_LIGHT if i % 2 == 0 else _NAVY_LIGHT)
                left = cards_x + (card_w + gap) * i
                _card(s, left=left, top=cards_top,
                      width=card_w, height=cards_h,
                      bg_color=bg,
                      header_text=hdr, body_text=body)
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "deck_card_grid_slide_failed", error=str(exc))
    _footer(s, idx, total)


def _split_card_line(raw: str, fallback_n: int) -> tuple[str, str]:
    """Split a single bullet string into (header, body). Tries '|'
    first, then ':', then falls back to a numbered header. Used by
    both _render_card_grid_slide and _render_split_panel_slide."""
    parts = raw.split("|", 1)
    if len(parts) == 2:
        return parts[0].strip(), parts[1].strip()
    if ":" in raw:
        h, b = raw.split(":", 1)
        return h.strip(), b.strip()
    return f"Point {fallback_n}", raw.strip()


def _render_scorecard_slide(
        prs, sl, chart_png, idx, total) -> None:
    """Slide 8 -- What the Model Gets Wrong: 2 of 9 scorecard.

    Standard _table renderer plus a cell_style_hook that:
      * Inspects a 'Council Signal' / 'Signal' column for BEAR /
        TRANSITION / BULL and paints the cell with the corresponding
        signal-palette fill + appropriate text color.
      * Shades 'Added value' outcome rows in _TEAL_LIGHT and leaves
        'Did not add value' rows on the default alternating fill.

    Column detection is case- + substring-tolerant so an LLM that
    emits 'Council signal' or 'Signal' both hit the badge path; same
    for 'Outcome' vs 'Result' on the row-shade column."""
    s = _blank(prs)
    _bg(s, _WHITE)
    title = sl.get("title") or SLIDE_TITLES[idx - 1]
    try:
        _title_bar(s, title)
        bullets = sl.get("bullets") or []
        bullets = [_strip_bullet_marker(b) for b in bullets]
        bullets = [b for b in bullets if b.strip()][:2]
        if bullets:
            _bullets(s, bullets,
                     left=Inches(0.5), top=Inches(1.2),
                     width=Inches(12.33), height=Inches(1.0),
                     size=14)
            table_top = Inches(2.3)
        else:
            table_top = Inches(1.4)
        headers, rows, _surviving = _slide_table(sl, slide_idx=idx)

        # Column-index lookups -- substring-tolerant.
        def _col_idx(*candidates: str) -> int | None:
            for cand in candidates:
                cand_l = cand.lower()
                for i, h in enumerate(headers):
                    if cand_l in str(h).lower():
                        return i
            return None
        signal_col = _col_idx("council signal", "signal")
        outcome_col = _col_idx("outcome", "result")

        def _scorecard_hook(*, row_idx, col_idx, header,
                            value, row_values, **_):
            v_lower = str(value).lower().strip()
            # Signal badge column.
            if signal_col is not None and col_idx == signal_col:
                if "bear" in v_lower:
                    return {"fill": _SIGNAL_RED, "text": _WHITE,
                            "bold": True}
                if "transition" in v_lower:
                    return {"fill": _SIGNAL_AMBER,
                            "text": _INK_DARK, "bold": True}
                if "bull" in v_lower:
                    return {"fill": _SIGNAL_GREEN,
                            "text": _WHITE, "bold": True}
            # Outcome-driven row shading.
            if outcome_col is not None:
                outcome_v = (
                    str(row_values[outcome_col])
                    if outcome_col < len(row_values) else "")
                if "added value" in outcome_v.lower():
                    return {"row_fill": _TEAL_LIGHT}
            return None

        if headers and rows:
            _table(s, headers, rows,
                   left=Inches(0.5), top=table_top,
                   width=Inches(12.33), max_rows=9,
                   cell_style_hook=_scorecard_hook)
        else:
            # June 27 2026 -- spec: [DATA PENDING] must never
            # appear. Log + render bullets-above-table area only.
            log.warning(
                "deck_scorecard_table_empty_skipping",
                slide=idx,
                note=("scorecard table_data missing; rendering "
                      "orientation bullets only per spec (no "
                      "[DATA PENDING] placeholder)"))
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "deck_scorecard_slide_failed", error=str(exc))
    _footer(s, idx, total)


def _render_feature_rows_slide(
        prs, sl, chart_png, idx, total) -> None:
    """Slide 10 -- Live Demo three-row feature block.

    Three full-width rows. Each row: icon circle on the left
    (_MOLLY_NAVY fill, numeric badge as the icon), bold feature
    title, description text. Row backgrounds alternate _TEAL_LIGHT
    / _NAVY_LIGHT / _GOLD_LIGHT to match Molly's three-row color
    scheme.

    Row content comes from the slide's bullets array. Each bullet
    is split on '|' (or ':') into (title, body) using the same
    convention as the card grid + split-panel."""
    s = _blank(prs)
    _bg(s, _WHITE)
    title = sl.get("title") or SLIDE_TITLES[idx - 1]
    try:
        _title_bar(s, title)
        # June 27 2026 -- spec: NEVER emit [DATA PENDING].
        # When the LLM returned no bullets, log + skip the feature
        # rows entirely; the title bar still renders. No padding
        # to 3 with placeholders.
        raw_bullets = sl.get("bullets") or []
        bullets = [_strip_bullet_marker(b) for b in raw_bullets]
        bullets = [b for b in bullets if b.strip()][:3]
        if not bullets:
            log.warning(
                "deck_feature_rows_bullets_empty_skipping",
                slide=idx,
                note=("LLM emitted no bullets for slide 10 feature "
                      "rows; skipping rows + rendering title only "
                      "per spec (no [DATA PENDING] placeholder)"))
            _footer(s, idx, total)
            return

        row_colors = (_TEAL_LIGHT, _NAVY_LIGHT, _GOLD_LIGHT)
        rows_top = Inches(1.3)
        row_h = Inches(1.75)
        row_gap = Inches(0.15)
        full_w = Inches(12.33)
        icon_w = Inches(1.2)

        for i, raw in enumerate(bullets):
            hdr, body = _split_card_line(raw, i + 1)
            top = Inches(
                1.3 + i * (1.75 + 0.15))
            bg = row_colors[i % len(row_colors)]
            # Row background rectangle.
            row_rect = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), top, full_w, row_h)
            row_rect.fill.solid()
            row_rect.fill.fore_color.rgb = bg
            row_rect.line.fill.background()
            row_rect.shadow.inherit = False
            # Icon circle (navy fill, white numeric badge).
            icon = s.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.75), top + Inches(0.275),
                Inches(1.2), Inches(1.2))
            icon.fill.solid()
            icon.fill.fore_color.rgb = _MOLLY_NAVY
            icon.line.fill.background()
            icon.shadow.inherit = False
            itf = icon.text_frame
            itf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ip = itf.paragraphs[0]
            ip.alignment = PP_ALIGN.CENTER
            ir = ip.add_run()
            ir.text = str(i + 1)
            ir.font.size = Pt(28)
            ir.font.bold = True
            ir.font.color.rgb = _WHITE
            # Header + body to the right of the icon.
            text_left = Inches(2.15)
            text_w = full_w - Inches(1.85)
            tbox = s.shapes.add_textbox(
                text_left, top + Inches(0.18),
                text_w, row_h - Inches(0.3))
            ttf = tbox.text_frame
            ttf.word_wrap = True
            tp = ttf.paragraphs[0]
            tp.space_after = Pt(4)
            tr = tp.add_run()
            tr.text = hdr
            tr.font.size = Pt(18)
            tr.font.bold = True
            tr.font.color.rgb = _MOLLY_NAVY
            if body:
                bp = ttf.add_paragraph()
                br = bp.add_run()
                br.text = body
                br.font.size = Pt(14)
                br.font.color.rgb = _INK_DARK
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "deck_feature_rows_slide_failed", error=str(exc))
    _footer(s, idx, total)


def _render_content_slide(prs, sl, chart_png, idx, total) -> None:
    """Lay out one content slide: title bar, bullets, optional table, optional
    chart. Adds exactly one slide. Slides in SLIDE_CHARTS always reserve a
    chart slot — a None PNG degrades to a [DATA PENDING] note. Fully guarded so
    a malformed slide never raises out of the builder.

    June 27 2026 (PR B) -- per-slide dispatch at the top routes the
    Molly-reference slides to their specialized renderers. Slides 2
    (Agenda) and 5 (Capital Preservation) fall through to the
    generic bullets / chart / table renderer below; everything else
    has a dedicated layout."""
    if idx == 1:
        return _render_title_slide(prs, sl, chart_png, idx, total)
    if idx == 3:
        return _render_split_panel_slide(
            prs, sl, chart_png, idx, total)
    if idx in (4, 6, 7, 9, 11):
        return _render_card_grid_slide(
            prs, sl, chart_png, idx, total)
    if idx == 8:
        return _render_scorecard_slide(
            prs, sl, chart_png, idx, total)
    if idx == 10:
        return _render_feature_rows_slide(
            prs, sl, chart_png, idx, total)
    s = _blank(prs)
    _bg(s, _WHITE)
    title = sl.get("title") or SLIDE_TITLES[idx - 1]
    try:
        _title_bar(s, title)
        # June 27 2026 -- spec: NEVER emit [DATA PENDING] in the
        # exported PPTX. Empty bullets now log + skip the bullet
        # block rather than emit a placeholder string. Same
        # treatment for post-strip empty + post-table-lift empty.
        bullets = list(sl.get("bullets") or [])
        # Strip chart-describing bullets on slides that already
        # have an embedded chart image (SLIDE_CHARTS).
        bullets = _strip_chart_bullets(bullets, idx)
        # _slide_table can lift a pipe-table out of the bullets
        # array (the 'LLM dumped the table into bullets' fallback).
        # When it does, the third return value is the bullets list
        # MINUS the lifted rows; None means 'no edit needed'.
        headers, rows, surviving = _slide_table(sl, slide_idx=idx)
        if surviving is not None:
            bullets = list(surviving)
        # Strip leading bullet-marker characters then drop empty.
        bullets = [_strip_bullet_marker(b) for b in bullets]
        bullets = [b for b in bullets if b.strip()]
        has_table = bool(headers and rows)
        has_chart = idx in SLIDE_CHARTS
        role = SLIDE_CHARTS.get(idx, "chart")
        if not bullets:
            log.warning(
                "deck_content_slide_bullets_empty_skipping",
                slide=idx, has_table=has_table, has_chart=has_chart,
                note=("LLM emitted no bullets for legacy content "
                      "slide; skipping bullet block + rendering "
                      "remaining elements (chart/table/title) per "
                      "spec (no [DATA PENDING] placeholder)"))

        # Layout dispatch -- when bullets is empty we DON'T call
        # _bullets at all (no placeholder text). Chart + table
        # still render at their sizes; the bullets-area space
        # becomes whitespace on slides 2 and 5.
        if has_chart and has_table:
            if bullets:
                _bullets(s, bullets, left=Inches(0.6), top=Inches(1.15),
                         width=Inches(12.1), height=Inches(1.55), size=13)
            _table(s, headers, rows, left=Inches(0.6), top=Inches(2.85),
                   width=Inches(6.1), max_rows=11)
            _image(s, chart_png, left=Inches(7.0), top=Inches(2.85),
                   width=Inches(5.8), fallback=role)
        elif has_chart:
            if bullets:
                _bullets(s, bullets, left=Inches(0.6), top=Inches(1.4),
                         width=Inches(4.9), height=Inches(5.2), size=15)
            _image(s, chart_png, left=Inches(5.7), top=Inches(1.5),
                   width=Inches(7.1), fallback=role)
        elif has_table:
            if bullets:
                _bullets(s, bullets, left=Inches(0.6), top=Inches(1.25),
                         width=Inches(12.1), height=Inches(1.95), size=14)
            _table(s, headers, rows, left=Inches(0.6), top=Inches(3.35),
                   width=Inches(12.1), max_rows=12)
        else:
            if bullets:
                _bullets(s, bullets, left=Inches(0.7), top=Inches(1.7),
                         width=Inches(11.9), height=Inches(5.0), size=18)
    except Exception as exc:  # noqa: BLE001 — one bad slide never fails the deck
        log.warning("deck_slide_body_failed", slide=idx, error=str(exc))

    _footer(s, idx, total)


def _dedupe_chart_assignments(
    charts: dict[int, bytes | None] | None,
) -> dict[int, bytes | None]:
    """June 26 2026 -- chart assignment deduplication.

    Two enforcement passes on the incoming {slide_number: png_bytes}
    dict before it reaches _render_content_slide:

    1. SLIDE_CHARTS slot enforcement -- a slide number that is NOT
       in SLIDE_CHARTS gets cleared to None even if the caller
       passed bytes. The renderer's downstream gate
       (`has_chart = idx in SLIDE_CHARTS`) already drops these,
       but clearing here makes the contract explicit and lets the
       dedup pass below operate only on intended slots.

    2. Bytes-identity dedup -- if the SAME PNG bytes object
       appears against multiple slide numbers (e.g. the caller's
       chart renderer dispatched the same role twice by accident),
       keep ONLY the slide whose SLIDE_CHARTS role canonically
       owns that chart. The 'canonical owner' is the lowest slide
       number that maps to the same role -- a stable tiebreaker
       that produces the same dedup regardless of dict-iteration
       order. Other slides with the duplicate bytes get cleared
       to None.

    Logs a WARNING when either pass clears a slot so the operator
    can see in Render logs that defence-in-depth fired. Returns a
    fresh dict; the input is not mutated.

    Why this exists: the symptom we're guarding against is the
    same chart PNG appearing on multiple slides (e.g. the
    'Risk-Return Profile by Strategy' chart showing up on slides
    4, 5, and 7). The PPTX path's `has_chart = idx in
    SLIDE_CHARTS` gate prevents stray chart rendering on its own,
    so the symptom requires a caller-side bug putting bytes on
    multiple keys. This dedup catches that defensively without
    requiring the caller to be perfect."""
    if not charts:
        return {}

    cleaned: dict[int, bytes | None] = {}
    for slide_num, png in charts.items():
        if slide_num not in SLIDE_CHARTS:
            if png is not None:
                log.warning(
                    "deck_chart_slot_outside_slide_charts",
                    slide_number=slide_num,
                    hint=("Chart bytes were passed for a slide "
                          "not in SLIDE_CHARTS. Cleared to None; "
                          "investigate the caller's chart "
                          "assignment logic."))
            cleaned[slide_num] = None
            continue
        cleaned[slide_num] = png

    # Bytes-identity dedup. For each duplicate bytes blob across
    # cleaned slots, keep only the canonical-owner slot.
    seen: dict[int, int] = {}  # id(png) -> slide_number (winner)
    for slide_num in sorted(cleaned.keys()):
        png = cleaned[slide_num]
        if png is None:
            continue
        key = id(png)
        if key in seen:
            winner = seen[key]
            log.warning(
                "deck_chart_duplicate_assignment",
                slide_number=slide_num,
                duplicate_of=winner,
                hint=("The same PNG bytes were assigned to two "
                      "slides. Cleared the later slot to None to "
                      "avoid duplicate chart rendering; the "
                      "canonical owner (lowest slide_number for "
                      "the role) keeps the chart."))
            cleaned[slide_num] = None
        else:
            seen[key] = slide_num
    return cleaned


def _substitute_pptx_text(
    prs: Any, substitution_table: dict[str, str] | None,
) -> int:
    """June 27 2026 -- post-build belt-and-suspenders pass.

    Walks every text frame on every slide (plus speaker notes) and
    replaces any remaining {{TOKEN}} placeholder with its value from
    substitution_table. Returns the number of token replacements
    performed.

    Why this exists:
      * The substitution table is normally applied at prompt-
        assembly time so Sonnet's slide JSON already carries
        substituted values. But when a renderer falls back to
        SLIDE_TITLES[idx-1] for a missing slide title (e.g. on a
        per-slide LLM failure), the un-substituted literal
        `{{REGIME_CONFIDENCE}}` placeholder leaks through.
      * This pass catches every remaining `{{...}}` placeholder
        and substitutes from the table.
      * Unknown tokens (not in the table) are LEFT INTACT --
        downstream operator-facing audit (a future addition)
        would flag those; for now they remain visible so the
        operator notices a typo'd token rather than silently
        rendering it as em-dash.

    No-ops when substitution_table is None / empty -- the deck
    renders the raw text as-is."""
    if not substitution_table:
        return 0
    import re as _re
    n_replacements = 0
    # Build a single combined regex of every token in the table for
    # one-pass replacement per text frame.
    tokens = [t for t in substitution_table if t.startswith("{{")
              and t.endswith("}}")]
    if not tokens:
        return 0
    token_pattern = _re.compile(
        "|".join(_re.escape(t) for t in tokens))

    def _sub(match: _re.Match) -> str:
        nonlocal n_replacements
        n_replacements += 1
        return substitution_table.get(match.group(0), match.group(0))

    for slide in prs.slides:
        # Body text frames + tables.
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for paragraph in tf.paragraphs:
                        for run in paragraph.runs:
                            if "{{" in run.text:
                                run.text = token_pattern.sub(
                                    _sub, run.text)
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    for row in tbl.rows:
                        for cell in row.cells:
                            for paragraph in (
                                    cell.text_frame.paragraphs):
                                for run in paragraph.runs:
                                    if "{{" in run.text:
                                        run.text = token_pattern.sub(
                                            _sub, run.text)
            except Exception as exc:  # noqa: BLE001
                # A single shape's substitution failure must not
                # abort the whole pass. Log + continue.
                log.warning(
                    "substitute_pptx_text_shape_failed",
                    error=str(exc))
        # Speaker notes -- same treatment.
        try:
            notes = slide.notes_slide
            if notes is not None:
                tf = notes.notes_text_frame
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        if "{{" in run.text:
                            run.text = token_pattern.sub(
                                _sub, run.text)
        except Exception as exc:  # noqa: BLE001
            log.warning(
                "substitute_pptx_text_notes_failed",
                error=str(exc))
    if n_replacements:
        log.info(
            "substitute_pptx_text_complete",
            replacements=n_replacements)
    return n_replacements


def build_presentation_deck(
    slides: Any,
    charts: dict[int, bytes | None] | None = None,
    substitution_table: dict[str, str] | None = None,
) -> bytes:
    """
    Lays out the 10-slide final presentation deck and returns the .pptx bytes.

    `slides` is the parsed AI JSON — a list of
    {slide_number, title, bullets, table_data, speaker_notes} dicts produced by
    deck_generation_prompt() through harness_narrative(). `charts` maps a slide
    number to its pre-rendered light-mode PNG (main renders them with the
    chart_render deck renderers, since that path is async — the sync builder
    only embeds the bytes).

    The builder always emits exactly ten slides with the canonical SLIDE_TITLES:
    a missing slide, a missing table, or a missing chart each degrade to a
    [DATA PENDING] note, never an exception. Every slide carries the AI DRAFT
    footer and a speaker-notes verification caveat; slide 1 additionally carries
    the review warning + submission checklist.

    June 26 2026 -- the charts dict is filtered through
    _dedupe_chart_assignments before slide rendering. Any caller-side
    bug that places the same chart on multiple slides, or places a chart
    on a non-SLIDE_CHARTS slot, is silently corrected here with a
    WARNING log naming the dropped slot.
    """
    charts = _dedupe_chart_assignments(charts)
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    total = DECK_SLIDE_COUNT

    norm = _normalize_slides(slides)
    for i, sl in enumerate(norm):
        _render_content_slide(prs, sl, charts.get(i + 1), i + 1, total)

    # Speaker notes — the AI narrative, prefixed with the verification caveat.
    # Slide 1 carries the full review warning + submission checklist (CAVEAT
    # 1 + 5); every other slide carries the per-slide verify note (CAVEAT 4).
    # The caveat is always non-empty, so every slide has speaker notes.
    for i, slide in enumerate(prs.slides):
        ai_notes = str(norm[i].get("speaker_notes") or "").strip()
        caveat = _DECK_TITLE_NOTE if i == 0 else _MOLLY_VERIFY_NOTE
        text = caveat + (f"\n\n{ai_notes}" if ai_notes else "")
        try:
            slide.notes_slide.notes_text_frame.text = text
        except Exception:  # noqa: BLE001 — a notes failure never fails the deck
            pass

    # June 27 2026 -- post-build {{...}} cleanup pass. Catches any
    # un-substituted placeholder that leaked through (e.g. when a
    # specialized renderer falls back to SLIDE_TITLES[idx-1] for a
    # missing slide title, the literal `{{REGIME_CONFIDENCE}}`
    # would otherwise render verbatim). No-op when substitution_
    # table is None (legacy callers).
    _substitute_pptx_text(prs, substitution_table)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Editor-content export ─────────────────────────────────────────────────────

# The editor canvas is a fixed 960x540 (16:9) space (migration 022). The
# exported deck uses the 10in x 5.625in 16:9 slide size, so the canvas
# maps onto the slide 1:1 — every element coordinate scales by a single
# EMU factor per axis, and the export matches the editor pixel-for-pixel.
_CANVAS_W = 960
_CANVAS_H = 540
_EDITOR_SLIDE_W = 9144000   # EMU — 10 in
_EDITOR_SLIDE_H = 5143500   # EMU — 5.625 in
# 960 canvas px span a 720 pt-wide slide — font sizes scale by 0.75.
_CANVAS_PT_FACTOR = 720.0 / _CANVAS_W


def _emu_x(x: float) -> int:
    """A canvas x-coordinate (or width), in px, to EMU."""
    return round(float(x) / _CANVAS_W * _EDITOR_SLIDE_W)


def _emu_y(y: float) -> int:
    """A canvas y-coordinate (or height), in px, to EMU."""
    return round(float(y) / _CANVAS_H * _EDITOR_SLIDE_H)


def _canvas_color(value: Any, default: RGBColor) -> RGBColor:
    """Parses a '#RRGGBB' canvas colour, falling back to `default`."""
    try:
        s = str(value or "").lstrip("#")
        if len(s) == 6:
            return RGBColor.from_string(s.upper())
    except Exception:  # noqa: BLE001 — a bad colour never fails the export
        pass
    return default


def _canvas_text(slide, el: dict[str, Any], left, top, width, height) -> None:
    """Places a canvas text element — its content, font size, weight,
    style and colour mapped onto a pptx textbox."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    size_pt = max(1, round(float(el.get("fontSize", 18)) * _CANVAS_PT_FACTOR))
    bold = str(el.get("fontWeight", "")) == "bold"
    italic = str(el.get("fontStyle", "")) == "italic"
    color = _canvas_color(el.get("color"), _INK)
    # Konva wraps a single string; explicit newlines become paragraphs.
    for i, line in enumerate(str(el.get("content") or "").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


# ── Editor-export markdown-table extraction ──────────────────────────────────
# June 26 2026 -- when a deck-draft text element's content carries a
# markdown table (3+ consecutive '| ... |' lines), the editor export
# now splits that text element into:
#   prose-before -> rendered via _canvas_text
#   table        -> rendered via _table -> real PPTX <a:tbl>
#   prose-after  -> rendered via _canvas_text
# Without this split the body text rendered as plain pipe-delimited
# strings in a textbox (zero <a:tbl> shapes in the PPTX), which
# regressed reviewability of every slide that carried a comparison
# table. The split logic looks for ONE table block per text element
# (the common case for the deck builder's _deck_slide_with_chart
# output); multiple tables in a single element would lift only the
# first and leave the rest as prose.


def _extract_markdown_table_from_text(
    text: str,
) -> tuple[str, tuple[list[str], list[list[str]]] | None, str]:
    """Returns (prose_before, (headers, rows) | None, prose_after).

    Scans the given text for the LONGEST consecutive run of lines
    matching the pipe-row regex (^\\s*\\|.*\\|\\s*$). A run of 2+
    lines is parsed via _parse_markdown_table; on success the run
    is removed from the surviving prose and returned as
    (headers, rows). The 2-line minimum is the markdown table
    bare-minimum (header + separator OR header + one data row).

    Returns (text, None, '') unchanged when no qualifying run is
    found -- the caller renders the original text as-is."""
    if not text or "|" not in text:
        return text, None, ""
    lines = text.splitlines()
    pipe_re = _PIPE_ROW_RE
    n = len(lines)

    # Find the longest consecutive run of pipe-row lines.
    best_start, best_len = -1, 0
    i = 0
    while i < n:
        if pipe_re.match(lines[i] or ""):
            j = i
            while j < n and pipe_re.match(lines[j] or ""):
                j += 1
            run_len = j - i
            if run_len > best_len:
                best_start, best_len = i, run_len
            i = j
        else:
            i += 1

    # Markdown table needs at least 2 lines (header + sep or
    # header + data); a single pipe line is not enough.
    if best_len < 2:
        return text, None, ""

    block = "\n".join(lines[best_start:best_start + best_len])
    headers, rows = _parse_markdown_table(block)
    if not headers or not rows:
        return text, None, ""

    prose_before = "\n".join(lines[:best_start]).rstrip()
    prose_after = "\n".join(
        lines[best_start + best_len:]).lstrip("\n")
    return prose_before, (headers, rows), prose_after


def _canvas_table(
    slide, headers: list[str], rows: list[list[str]],
    left, top, width, height,
) -> None:
    """Render a markdown-derived table as a native PPTX table
    inside the editor-export canvas region. Reuses the styling
    contract of _table (navy header band, alternating row fill,
    Pt 11/10 sizing) but accepts EMU-typed left/top/width to
    match the canvas coordinate flow + caps row count by the
    available height (Inches(0.32) per row, conservative)."""
    from pptx.util import Emu, Inches as _In, Pt as _Pt

    if not headers or not rows:
        return
    max_rows_for_height = max(1, int(height / _In(0.32)))
    rows = rows[: max_rows_for_height - 1] if (
        max_rows_for_height >= 2) else rows[:1]
    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, height)
    table = tbl_shape.table
    for c, label in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _NAVY
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = str(label)
        run.font.size = _Pt(11)
        run.font.bold = True
        run.font.color.rgb = _WHITE
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            if c >= n_cols:
                break
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                _WHITE if r % 2
                else RGBColor(0xF1, 0xF5, 0xF9))
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = str(value)
            run.font.size = _Pt(10)


def _canvas_text_or_table(
    slide, el: dict[str, Any], left, top, width, height,
) -> None:
    """Wraps _canvas_text: when the element's content carries a
    parseable markdown table, splits into prose-before / table /
    prose-after sub-regions and renders each in its own shape so
    the table becomes a real PPTX <a:tbl> instead of raw pipes in
    a textbox. Falls back to plain _canvas_text when no table is
    detected."""
    from pptx.util import Emu

    content = str(el.get("content") or "")
    prose_before, table, prose_after = (
        _extract_markdown_table_from_text(content))
    if table is None:
        _canvas_text(slide, el, left, top, width, height)
        return

    headers, rows = table

    # Vertical layout: split the height into (prose_before |
    # table | prose_after) proportionally to line count.
    pb_lines = max(
        prose_before.count("\n") + 1, 1) if prose_before else 0
    pa_lines = max(
        prose_after.count("\n") + 1, 1) if prose_after else 0
    table_lines = len(rows) + 1
    total_lines = pb_lines + table_lines + pa_lines
    if total_lines <= 0:
        return
    h_per_line = int(height / total_lines)
    pb_h = Emu(h_per_line * pb_lines) if pb_lines else Emu(0)
    table_h = Emu(h_per_line * table_lines)
    pa_h = Emu(h_per_line * pa_lines) if pa_lines else Emu(0)

    cursor_top = top
    if pb_lines:
        pb_el = dict(el)
        pb_el["content"] = prose_before
        _canvas_text(slide, pb_el, left, cursor_top, width, pb_h)
        cursor_top = Emu(int(cursor_top) + int(pb_h))
    _canvas_table(
        slide, headers, rows, left, cursor_top, width, table_h)
    cursor_top = Emu(int(cursor_top) + int(table_h))
    if pa_lines:
        pa_el = dict(el)
        pa_el["content"] = prose_after
        _canvas_text(slide, pa_el, left, cursor_top, width, pa_h)


def build_editor_pptx(
    draft: dict[str, Any],
    chart_pngs: dict[str, bytes] | None = None,
) -> bytes:
    """
    Renders a presentation_deck editor draft to a .pptx straight from its
    canvas content_json (migration 022). Each slide's positioned text and
    chart elements are mapped onto a 10x5.625in 16:9 slide so the export
    matches the editor canvas 1:1.

    `chart_pngs` maps a chart element id to its server-rendered PNG — the
    caller renders them, since that path is async — and a missing PNG
    degrades to a [DATA PENDING] note rather than failing the export.

    The presenter's speaker notes are carried into each slide's notes,
    prefixed with the AI-draft verification reminder (the editor export
    is a faithful WYSIWYG render, so it carries no on-slide chrome).
    """
    from pptx import Presentation
    from pptx.util import Emu

    pngs = chart_pngs or {}
    prs = Presentation()
    prs.slide_width = Emu(_EDITOR_SLIDE_W)
    prs.slide_height = Emu(_EDITOR_SLIDE_H)

    content = draft.get("content_json") or {}
    slides = content.get("slides", []) if isinstance(content, dict) else []

    for sl in slides:
        if not isinstance(sl, dict):
            continue
        s = _blank(prs)
        # Slide background — the canvas background colour, full bleed.
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                Emu(_EDITOR_SLIDE_W), Emu(_EDITOR_SLIDE_H))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _canvas_color(sl.get("background"), _WHITE)
        bg.line.fill.background()
        bg.shadow.inherit = False
        s.shapes._spTree.remove(bg._element)
        s.shapes._spTree.insert(2, bg._element)

        for el in (sl.get("elements") or []):
            if not isinstance(el, dict):
                continue
            left = Emu(_emu_x(el.get("x", 0)))
            top = Emu(_emu_y(el.get("y", 0)))
            width = Emu(_emu_x(el.get("width", 0)))
            height = Emu(_emu_y(el.get("height", 0)))
            if el.get("type") == "text":
                # June 26 2026 -- _canvas_text_or_table splits the
                # text into prose-before / markdown-table / prose-
                # after sub-regions when the content carries a
                # pipe-delimited table, so the table renders as a
                # real PPTX <a:tbl> instead of raw pipes in a
                # textbox. Falls through to plain _canvas_text for
                # text elements with no detectable table -- so
                # existing behaviour is preserved for non-table
                # bodies.
                _canvas_text_or_table(
                    s, el, left, top, width, height)
            elif el.get("type") == "chart":
                png = pngs.get(str(el.get("id")))
                if png:
                    s.shapes.add_picture(io.BytesIO(png), left, top,
                                         width=width, height=height)
                else:
                    # June 27 2026 -- spec: [DATA PENDING] must
                    # never appear in exported PPTX. Log + skip
                    # the chart element rather than emit a
                    # placeholder; the slide's other elements
                    # still render around the empty chart area.
                    log.warning(
                        "editor_pptx_chart_png_missing_skipping",
                        element_id=str(el.get("id") or ""),
                        note=("chart PNG missing; skipping element "
                              "per spec (no [DATA PENDING] "
                              "placeholder). Operator: warm "
                              "analytics caches + re-export."))
            elif el.get("type") == "table":
                # June 26 2026 -- first-class type='table' canvas
                # element. Renders as a real PPTX <a:tbl> via
                # _canvas_table, using the element's table_config
                # for headers (table_config.columns) + body rows
                # (table_config.rows). The chart_config-style
                # 'absent = fall back' contract applies: a missing
                # table_config produces a [DATA PENDING] placeholder
                # rather than raising.
                tc = el.get("table_config") or {}
                headers = list(tc.get("columns") or [])
                raw_rows = list(tc.get("rows") or [])
                # rows may be a list of lists (canonical) OR a list
                # of strategy names (when the editor's Configure
                # panel set only the strategy selection and left
                # the cell data for the renderer to look up). Only
                # the list-of-lists shape is renderable here; flat
                # strategy lists degrade to a placeholder until
                # the runtime cell-lookup layer lands.
                rows = [
                    list(r) for r in raw_rows
                    if isinstance(r, (list, tuple))
                ]
                if headers and rows:
                    _canvas_table(
                        s, headers, rows, left, top, width, height)
                else:
                    # June 27 2026 -- spec: [DATA PENDING] must
                    # never appear. Log + skip the table element.
                    log.warning(
                        "editor_pptx_table_data_missing_skipping",
                        element_id=str(el.get("id") or ""),
                        note=("table_config rows/columns missing; "
                              "skipping element per spec (no "
                              "[DATA PENDING] placeholder). "
                              "Configure rows + columns in the "
                              "slide editor."))

        # Speaker notes — the verify reminder, then the presenter's notes.
        try:
            notes = str(sl.get("speaker_notes") or "").strip()
            s.notes_slide.notes_text_frame.text = (
                _MOLLY_VERIFY_NOTE + ("\n\n" + notes if notes else ""))
        except Exception:  # noqa: BLE001 — a notes failure never fails export
            pass

    if not slides:
        s = _blank(prs)
        _textbox(s, Emu(_emu_x(120)), Emu(_emu_y(230)),
                 Emu(_emu_x(720)), Emu(_emu_y(80)),
                 "This deck draft has no slides yet.", size=18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
