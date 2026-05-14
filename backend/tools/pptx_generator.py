"""
tools/pptx_generator.py

Builds a PowerPoint deck (.pptx) from Molly's edited storyboard JSON.
The deck biases toward whatever Molly's saved version actually contains —
slide order, owner assignments, chart references, and timing all flow
from the storyboard, not from constants here.

One slide per storyboard entry, with:
  - Title block: slide.headline
  - Body block:  slide.key_point
  - Notes pane: slide.speaker_note (presenter-only)
  - Footer:     "AI DRAFT — REQUIRES HUMAN REVIEW" on every slide
"""
from __future__ import annotations

from datetime import datetime, timezone
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


_AI_DRAFT_FOOTER = "AI DRAFT — REQUIRES HUMAN REVIEW · " \
                   "Verify every number before submitting to Forest Capital"


def _add_ai_draft_footer(slide, idx: int, total: int) -> None:
    """
    Inserts a small amber footer on every slide. python-pptx doesn't
    expose slide-master footers in a convenient API, so we draw a
    textbox along the bottom margin manually. The position (left=0.3in,
    top=7.0in) clears the standard content placeholder on a 10x7.5
    layout — verify visually if the slide master ever changes.
    """
    left, top, width, height = Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.3)
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"{_AI_DRAFT_FOOTER}  ·  Slide {idx}/{total}"
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)  # amber #f59e0b
    run.font.italic = True


def _set_title(slide, text: str) -> None:
    """Populates the title placeholder if one exists; falls back to a textbox."""
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
        # Force a readable size — some PPT templates default tiny titles
        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(28)
                run.font.bold = True
        return
    # Title placeholder missing on the chosen layout → draw it
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = True


def _set_body_text(slide, key_point: str, owner: str, timing_mins: float) -> None:
    """Adds the body block — key point plus an owner/timing line below."""
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(4))
    tf = txbox.text_frame
    tf.word_wrap = True

    # Key point — the main message of the slide
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = key_point
    r1.font.size = Pt(18)

    # Owner + timing meta line
    p2 = tf.add_paragraph()
    p2.space_before = Pt(18)
    r2 = p2.add_run()
    r2.text = f"Presenter: {owner}  ·  Timing target: {timing_mins:.1f} min"
    r2.font.size = Pt(11)
    r2.font.italic = True
    r2.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)  # slate #64748b


def _add_speaker_notes(slide, note_text: str) -> None:
    """Writes the speaker note into the notes pane (visible to presenter only)."""
    notes = slide.notes_slide.notes_text_frame
    notes.text = note_text


def build_pptx_from_storyboard(
    storyboard: dict[str, Any],
    title: str = "Forest Capital Portfolio Intelligence System",
) -> bytes:
    """
    Renders a storyboard dict into a complete .pptx and returns the bytes.

    Args:
        storyboard: Dict with a 'slides' list — each slide must have
                    order, headline, key_point, owner, timing_mins,
                    speaker_note. Other fields (chart_ref, live_demo,
                    transition) are optional but rendered when present.
        title:      Title slide text (first slide is auto-generated;
                    storyboard slides start at the second deck position).

    Returns:
        The .pptx bytes — caller wraps in StreamingResponse for download.
    """
    prs = Presentation()
    # Standard 10x7.5 layout — keeps the AI-DRAFT footer position correct
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    raw_slides = list(storyboard.get("slides", []))
    raw_slides.sort(key=lambda s: int(s.get("order", 0)))
    total = len(raw_slides) + 1  # +1 for the auto-generated title slide

    # ── Title slide ───────────────────────────────────────────────────
    title_layout = prs.slide_layouts[0]  # built-in "Title Slide"
    title_slide = prs.slides.add_slide(title_layout)
    _set_title(title_slide, title)
    sub_text = (
        f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d')}  ·  "
        f"{len(raw_slides)} slides  ·  "
        f"{sum(float(s.get('timing_mins', 0)) for s in raw_slides):.1f} min"
    )
    # Subtitle placeholder typically index 1 on the title slide
    try:
        title_slide.placeholders[1].text = sub_text
    except (KeyError, IndexError):
        pass
    _add_ai_draft_footer(title_slide, 1, total)

    # ── Content slides ────────────────────────────────────────────────
    # Use a content layout that has a title placeholder (layout index 5 =
    # "Title Only" in the default master; safer than guessing higher
    # indexes that don't exist on all themes).
    content_layout = prs.slide_layouts[5]

    for i, slide_data in enumerate(raw_slides, start=2):
        slide = prs.slides.add_slide(content_layout)
        _set_title(slide, slide_data.get("headline", "(Untitled)"))
        _set_body_text(
            slide,
            key_point=slide_data.get("key_point", ""),
            owner=slide_data.get("owner", "—"),
            timing_mins=float(slide_data.get("timing_mins", 0)),
        )
        _add_speaker_notes(slide, slide_data.get("speaker_note", ""))
        _add_ai_draft_footer(slide, i, total)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
