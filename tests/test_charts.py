"""
tests/test_charts.py

Tests for the Konva canvas presentation editor's server side:
  - the chart render endpoints — GET /api/v1/charts/available and
    /api/v1/charts/render/{key} (commit 2)
  - migration 022's slide-card ↔ canvas-element conversion (commit 1)
  - the canvas-layout PPTX export — build_editor_pptx + the EMU
    coordinate mapping (commit 4)

Endpoint-contract tests run everywhere. Chart rendering degrades to a
Pillow placeholder PNG when matplotlib / the analytics caches are
unavailable, so the render tests pass without a live database.
"""
from __future__ import annotations

import asyncio
import importlib.util
import io
import os
import sys

import pytest
from fastapi.testclient import TestClient

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))
os.environ.setdefault("ENVIRONMENT", "test")
os.environ.setdefault("SECRET_KEY", "test-secret-key-at-least-32-characters-long")
os.environ.setdefault("MASTER_API_KEY", "test_master_key")
os.environ.setdefault(
    "ALLOWED_EMAILS",
    "ruurdsm@queens.edu,thaob@queens.edu,murdockm@queens.edu,panttserk@queens.edu",
)

from main import app  # noqa: E402
from auth import generate_session_token  # noqa: E402

client = TestClient(app)
TEAM = {"X-API-Key": generate_session_token("thaob@queens.edu")}
VIEWER = {"X-API-Key": generate_session_token("panttserk@queens.edu")}

_PNG_MAGIC = b"\x89PNG\r\n\x1a\n"


# ── Chart render endpoints ────────────────────────────────────────────────────

class TestChartEndpoints:
    def test_available_requires_auth(self):
        assert client.get("/api/v1/charts/available").status_code == 401

    def test_available_rejects_a_viewer(self):
        # The charts endpoints are team-gated.
        assert client.get("/api/v1/charts/available",
                           headers=VIEWER).status_code == 403

    def test_available_lists_the_renderable_charts(self):
        resp = client.get("/api/v1/charts/available", headers=TEAM)
        assert resp.status_code == 200
        charts = resp.json()
        # The library grew beyond the original five with the extended
        # renderers (regime + factors landed in Commit 2 of the chart
        # library expansion). Assert against the known-set length, which
        # tracks AVAILABLE_CHARTS exactly.
        from tools.chart_render import AVAILABLE_CHARTS
        assert isinstance(charts, list)
        assert len(charts) == len(AVAILABLE_CHARTS)
        for c in charts:
            assert {"key", "label", "description", "category"} <= set(c)

    def test_render_returns_a_png(self):
        resp = client.get("/api/v1/charts/render/rolling_correlation",
                           headers=TEAM)
        assert resp.status_code == 200
        assert resp.headers["content-type"] == "image/png"
        assert resp.content.startswith(_PNG_MAGIC)

    def test_render_unknown_key_is_404(self):
        resp = client.get("/api/v1/charts/render/not_a_real_chart",
                           headers=TEAM)
        assert resp.status_code == 404

    def test_render_requires_auth(self):
        assert client.get(
            "/api/v1/charts/render/risk_return").status_code == 401

    def test_render_populates_the_cache(self):
        from tools.chart_render import _render_cache
        client.get("/api/v1/charts/render/team_activity",
                   headers=TEAM, params={"width": 240, "height": 160})
        # The 5-minute per-(key, theme, w, h) cache holds the render.
        assert any(k[0] == "team_activity" for k in _render_cache)


# ── chart_render unit behaviour ───────────────────────────────────────────────

class TestChartRenderUnit:
    def test_is_known_chart(self):
        from tools.chart_render import is_known_chart
        assert is_known_chart("cumulative_returns")
        assert not is_known_chart("nonsense")

    def test_available_charts_keys_match_the_known_set(self):
        from tools.chart_render import AVAILABLE_CHARTS, is_known_chart
        # Every entry on AVAILABLE_CHARTS must be in the known-keys
        # frozenset _CHART_KEYS — the two are derived together and must
        # stay in sync as the library grows. The size assertion is
        # parameterised on the registry length, so a future addition to
        # AVAILABLE_CHARTS does not require touching this test.
        assert len(AVAILABLE_CHARTS) > 0
        assert all(is_known_chart(c["key"]) for c in AVAILABLE_CHARTS)
        # Categories are stable strings the chart picker groups on.
        assert all(c["category"] for c in AVAILABLE_CHARTS)

    def test_render_chart_png_returns_a_png(self):
        from tools.chart_render import render_chart_png
        png = asyncio.run(render_chart_png("risk_return", "light", 300, 200))
        assert png.startswith(_PNG_MAGIC)

    def test_render_dark_theme_falls_back_to_a_png(self):
        # The matplotlib renderers are light-only — dark still yields a PNG.
        from tools.chart_render import render_chart_png
        png = asyncio.run(render_chart_png("sensitivity", "dark", 200, 120))
        assert png.startswith(_PNG_MAGIC)

    def test_all_chart_keys_listed_on_available_endpoint(self):
        # Every key in AVAILABLE_CHARTS surfaces in the /charts/available
        # response — the API and the registry must never drift apart, or
        # the chart picker hides charts that still render server-side.
        from tools.chart_render import AVAILABLE_CHARTS
        resp = client.get("/api/v1/charts/available", headers=TEAM)
        assert resp.status_code == 200
        seen = {c["key"] for c in resp.json()}
        registered = {c["key"] for c in AVAILABLE_CHARTS}
        assert seen == registered

    def test_every_chart_renders_a_png_in_test_environment(self):
        # In the test environment the analytics caches are cold and
        # matplotlib renders against missing data — the renderers must
        # ALL degrade to the Pillow placeholder PNG rather than 500.
        # This is the chart-library fail-open contract.
        from tools.chart_render import AVAILABLE_CHARTS, render_chart_png
        for chart in AVAILABLE_CHARTS:
            png = asyncio.run(render_chart_png(
                chart["key"], "light", 240, 160))
            assert png.startswith(_PNG_MAGIC), \
                f"{chart['key']} did not return a PNG"

    def test_render_cache_hit_after_first_call_for_extended_charts(self):
        # The 5-minute per-(chart_key, theme, w, h) cache must serve the
        # extended renderers identically to the deck five. A second call
        # for the same key/theme/size is bytes-identical to the first
        # because the cache returns the same object — no re-render.
        from tools.chart_render import _render_cache, render_chart_png
        # Pick a Commit-3 key (single-strategy) and a Commit-4 key
        # (multi-strategy gates) so both compute paths are exercised.
        for key in ("drawdown_periods", "significance_journey"):
            _render_cache.clear()
            a = asyncio.run(render_chart_png(key, "light", 220, 140))
            assert any(k[0] == key for k in _render_cache)
            b = asyncio.run(render_chart_png(key, "light", 220, 140))
            # Identical bytes — second call hit the cache, not the renderer.
            assert a == b

    def test_categories_cover_every_brief_group(self):
        # The brief's chart-picker layout lists six section headers. Every
        # one must have at least one chart on the registry, or the picker
        # would render an empty section.
        from tools.chart_render import AVAILABLE_CHARTS
        cats = {c["category"] for c in AVAILABLE_CHARTS}
        # The six display groups per Commit 5/7 of the chart-library build.
        for expected in ("regime", "factors", "performance", "risk",
                         "significance", "activity"):
            assert expected in cats, f"no charts in '{expected}' category"

    def test_extended_chart_keys_registered_and_rendered(self):
        # The library expansion (commits 2-4) added canvas-only renderers
        # in three batches: regime + factors (Commit 2), performance +
        # risk (Commit 3), and significance (Commit 4). Each new key must
        # (a) be in AVAILABLE_CHARTS with the correct category,
        # (b) resolve via is_known_chart, and (c) round-trip through
        # render_chart_png to a PNG (the placeholder when source data is
        # unavailable — also a PNG).
        from tools.chart_render import (
            AVAILABLE_CHARTS, is_known_chart, render_chart_png,
        )
        registry = {c["key"]: c for c in AVAILABLE_CHARTS}
        # Categories are the chart picker's display grouping (Commit 5 of
        # the library expansion). rolling_correlation lives under
        # "performance" (the time-series block); risk_return + sensitivity
        # both live under "risk"; team_activity → "activity".
        expected = {
            # Commit 2
            "regime_signals":              "regime",
            "regime_conditional_returns":  "regime",
            "factor_loadings":             "factors",
            "factor_returns_attribution":  "factors",
            # Commit 3
            "rolling_sharpe":              "performance",
            "return_distribution":         "performance",
            "monthly_returns_heatmap":     "performance",
            "drawdown_periods":            "risk",
            # Commit 4
            "significance_journey":        "significance",
            "oos_performance":             "significance",
            "p_value_distribution":        "significance",
            # Picker regrouping (Commit 5)
            "rolling_correlation":         "performance",
            "risk_return":                 "risk",
            "sensitivity":                 "risk",
            "team_activity":               "activity",
        }
        for key, category in expected.items():
            assert is_known_chart(key), f"{key} not in _CHART_KEYS"
            assert key in registry, f"{key} not in AVAILABLE_CHARTS"
            assert registry[key]["category"] == category
            png = asyncio.run(render_chart_png(key, "light", 240, 160))
            assert png.startswith(_PNG_MAGIC), f"{key} did not render a PNG"


# ── Migration 022 — slide-card ↔ canvas conversion ────────────────────────────

def _load_migration_022():
    path = os.path.join(
        os.path.dirname(__file__), "..", "backend", "migrations",
        "versions", "022_canvas_schema_conversion.py")
    spec = importlib.util.spec_from_file_location("migration_022", path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)  # type: ignore[union-attr]
    return mod


class TestMigration022:
    mig = _load_migration_022()

    def test_slide_to_canvas_maps_card_fields_to_elements(self):
        canvas = self.mig._slide_to_canvas({
            "id": 3, "title": "Results", "content": "Body text",
            "data_points": ["Sharpe 0.63", "CAGR 7.7%"],
            "speaker_notes": "rehearsal", "verified": True,
            "notes_written": True})
        assert canvas["id"] == 3
        assert canvas["background"] == "#FFFFFF"
        assert canvas["speaker_notes"] == "rehearsal"
        # title, content and data points become three text elements.
        els = canvas["elements"]
        assert [e["id"] for e in els] == ["el_001", "el_002", "el_003"]
        assert els[0]["content"] == "Results" and els[0]["fontWeight"] == "bold"
        assert els[1]["content"] == "Body text"
        assert "Sharpe 0.63" in els[2]["content"]
        # The card-only verified / notes_written keys are dropped.
        assert "verified" not in canvas and "notes_written" not in canvas

    def test_slide_to_canvas_omits_data_point_element_when_empty(self):
        canvas = self.mig._slide_to_canvas({
            "id": 1, "title": "T", "content": "C", "data_points": [],
            "speaker_notes": ""})
        assert [e["id"] for e in canvas["elements"]] == ["el_001", "el_002"]

    def test_slide_to_canvas_is_idempotent(self):
        already = {"id": 1, "title": "T", "background": "#FFF",
                   "speaker_notes": "", "elements": [{"id": "x", "type": "text"}]}
        assert self.mig._slide_to_canvas(already) is already

    def test_slide_to_cards_downgrades_canvas_to_card_shape(self):
        card = self.mig._slide_to_cards({
            "id": 2, "title": "Findings", "background": "#FFFFFF",
            "speaker_notes": "notes",
            "elements": [
                {"id": "el_001", "type": "text", "y": 40,
                 "content": "Findings"},
                {"id": "el_002", "type": "text", "y": 140,
                 "content": "Body"},
                {"id": "el_003", "type": "text", "y": 440,
                 "content": "DP one\nDP two"}]})
        assert card["content"] == "Body"
        assert card["data_points"] == ["DP one", "DP two"]
        assert card["speaker_notes"] == "notes"
        # The downgrade restores the card-only keys.
        assert card["verified"] is False
        assert card["notes_written"] is True

    def test_slide_to_cards_is_idempotent(self):
        already = {"id": 1, "title": "T", "content": "C", "data_points": []}
        assert self.mig._slide_to_cards(already) is already

    def test_canvas_then_cards_round_trips_the_core_fields(self):
        original = {"id": 5, "title": "Round", "content": "Body",
                    "data_points": ["A", "B"], "speaker_notes": "say this"}
        back = self.mig._slide_to_cards(self.mig._slide_to_canvas(original))
        assert back["content"] == "Body"
        assert back["data_points"] == ["A", "B"]
        assert back["speaker_notes"] == "say this"


# ── Canvas PPTX export — EMU mapping + build_editor_pptx ──────────────────────

class TestCanvasPptxExport:
    def test_emu_mapping_scales_the_canvas_onto_the_slide(self):
        from tools.academic_deck import _emu_x, _emu_y
        # The 960x540 canvas maps onto a 10x5.625in (16:9) slide.
        assert _emu_x(0) == 0
        assert _emu_x(960) == 9144000
        assert _emu_x(480) == 4572000
        assert _emu_y(0) == 0
        assert _emu_y(540) == 5143500

    def test_build_editor_pptx_renders_a_canvas_draft(self):
        from pptx import Presentation
        from tools.academic_deck import build_editor_pptx

        draft = {"document_type": "presentation_deck", "content_json": {
            "slides": [{
                "id": 1, "title": "Opening", "background": "#1B2A4A",
                "speaker_notes": "UAT-NOTE rehearsal",
                "elements": [
                    {"id": "el_001", "type": "text", "x": 60, "y": 40,
                     "width": 840, "height": 80, "content": "Opening",
                     "fontSize": 36, "fontWeight": "bold",
                     "fontStyle": "normal", "color": "#FFFFFF",
                     "locked": False},
                    {"id": "c1", "type": "chart", "x": 540, "y": 160,
                     "width": 360, "height": 220, "chartKey": "risk_return",
                     "verified": True, "locked": False}]}]}}
        out = build_editor_pptx(draft, {})
        assert out.startswith(b"PK")
        prs = Presentation(io.BytesIO(out))
        # The editor slide is the 10x5.625in 16:9 size.
        assert prs.slide_width == 9144000
        assert prs.slide_height == 5143500
        # The chart had no PNG — it degrades to a [DATA PENDING] note,
        # never a failed export. The speaker notes carry through.
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "UAT-NOTE rehearsal" in notes

    def test_build_editor_pptx_embeds_a_rendered_chart_png(self):
        from pptx import Presentation
        from tools.academic_deck import build_editor_pptx
        from tools.chart_render import render_chart_png

        png = asyncio.run(render_chart_png("risk_return", "light", 400, 240))
        draft = {"document_type": "presentation_deck", "content_json": {
            "slides": [{"id": 1, "title": "T", "background": "#FFFFFF",
                        "speaker_notes": "",
                        "elements": [{"id": "c1", "type": "chart",
                                      "x": 100, "y": 100, "width": 400,
                                      "height": 240, "chartKey": "risk_return",
                                      "verified": False, "locked": False}]}]}}
        out = build_editor_pptx(draft, {"c1": png})
        prs = Presentation(io.BytesIO(out))
        # background rect + the chart picture.
        assert len(prs.slides[0].shapes) >= 2

    def test_build_editor_pptx_handles_an_empty_deck(self):
        from pptx import Presentation
        from tools.academic_deck import build_editor_pptx
        out = build_editor_pptx(
            {"document_type": "presentation_deck", "content_json": {"slides": []}})
        prs = Presentation(io.BytesIO(out))
        assert len(prs.slides) == 1
