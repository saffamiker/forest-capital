"""tests/test_deck_specialized_renderers.py -- PR B (June 27 2026).

Pins the per-slide specialized renderers added in PR B + the
cell_style_hook on _table + the canvas-mirror dispatch in
editor_content. The renderers themselves don't have rich return
values; these tests walk the pptx Presentation and assert shape
counts, text content, fills, and table cell styling.

Test groups:

  TestDispatch
    _render_content_slide must route slides 1, 3, 4/6/7/9/11, 8, 10
    to their specialized renderers and let slides 2 / 5 fall through
    to the generic layout.

  TestTitleSlide
    Slide 1 carries the title chrome -- navy header bar, teal accent
    rule, subtitle, presenter line. No body bullets.

  TestSplitPanelSlide
    Slide 3 carries a left stat-card column + a right OOS table +
    the "+85% vs benchmark" callout. The Regime-Conditional row in
    the OOS table is highlighted teal.

  TestCardGridSlide
    Card-grid slides without a chart spread the cards horizontally;
    slide 4 (chart in SLIDE_CHARTS) stacks them vertically on the
    left so the chart fits on the right.

  TestScorecardSlide
    Slide 8 applies BEAR/TRANSITION/BULL signal-color cell fills
    and shades 'Added value' rows in light teal.

  TestFeatureRowsSlide
    Slide 10 draws three full-width rows with numeric badge icons
    1 / 2 / 3 + per-row alternating color backgrounds.

  TestTableCellStyleHook
    _table cell_style_hook applies per-cell fills + per-row fills +
    bold + text colors as the hook returns them; missing hook is
    backward-compatible.

  TestCanvasMirror
    deck_slides_to_editor produces title-chrome elements for slide
    1 and a left-body / right-table split for slide 3.
"""
from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.util import Inches


def _build_deck(overrides: dict[int, dict] | None = None):
    """Build a full 11-slide pptx with sane defaults; lets each test
    override one or two slides while keeping the deck buildable.
    Returns the parsed Presentation."""
    from tools.academic_deck import build_presentation_deck

    base = [
        {"slide_number": 1, "title": (
            "Yes -- Regime-Conditional Beats 100% Equity OOS"),
         "bullets": [], "table_data": None, "speaker_notes": ""},
        {"slide_number": 2, "title": "Agenda",
         "bullets": ["A", "B", "C"],
         "table_data": None, "speaker_notes": ""},
        {"slide_number": 3,
         "title": "The Investment Case: Setup and OOS Verdict",
         "bullets": [
             "Static (60/40) | IS Sharpe 0.71",
             "Dynamic (Regime-Conditional) | IS Sharpe 0.93",
             "Benchmark (100% S&P 500) | IS Sharpe 0.58",
         ],
         "table_data": {
             "headers": [
                 "Strategy", "OOS Sharpe", "Max DD", "OOS CAGR"],
             "rows": [
                 ["Dynamic Blend (Regime-Conditional)",
                  "0.86", "-15%", "11.4%"],
                 ["Classic 60/40", "0.62", "-22%", "7.8%"],
                 ["100% Equity Benchmark",
                  "0.43", "-34%", "6.9%"],
             ]},
         "speaker_notes": ""},
        {"slide_number": 4, "title": "Why Static Failed in 2022",
         "bullets": [
             "Correlation Break | sign flip 2022",
             "Static Misallocation | 60/40 lost both",
             "Regime Signal | HMM caught 3mo early"],
         "table_data": None, "speaker_notes": ""},
        {"slide_number": 5, "title": "Capital Preservation",
         "bullets": ["Half the drawdown"],
         "table_data": None, "speaker_notes": ""},
        {"slide_number": 6, "title": "Does It Hold Up OOS? Yes.",
         "bullets": [
             "53-month window | post-2022",
             "Sharpe held | within 8% of IS",
             "Walk-forward | re-fit yearly"],
         "table_data": None, "speaker_notes": ""},
        {"slide_number": 7, "title": "Live Regime Signal",
         "bullets": [
             "VIX | 13.2 low",
             "Yield Curve | +25 bps",
             "Credit Spread | 95 bps"],
         "table_data": None, "speaker_notes": ""},
        {"slide_number": 8, "title": "What the Model Gets Wrong",
         "bullets": ["9 stress scenarios", "2 of 9 missed"],
         "table_data": {
             "headers": ["Scenario", "Council Signal", "Outcome"],
             "rows": [
                 ["GFC 2008", "BEAR", "Added value"],
                 ["Euro 2011", "TRANSITION", "Added value"],
                 ["Taper 2013", "BULL", "Did not add value"],
                 ["China 2015", "BEAR", "Did not add value"],
                 ["Brexit 2016", "TRANSITION", "Added value"],
                 ["Vol 2018", "BEAR", "Added value"],
                 ["COVID 2020", "BEAR", "Added value"],
                 ["Rate 2022", "TRANSITION", "Added value"],
                 ["Bank 2023", "BULL", "Added value"],
             ]},
         "speaker_notes": ""},
        {"slide_number": 9, "title": "How We Used AI",
         "bullets": [
             "Generator-Evaluator | Sonnet + Opus",
             "Multi-model panel | Gemini + Grok",
             "Deterministic recompute | Python verified"],
         "table_data": None, "speaker_notes": ""},
        {"slide_number": 10, "title": "Live Demo",
         "bullets": [
             "Regime Classifier | Real-time HMM",
             "Strategy Comparison | Backtest subsets",
             "CIO Recommendation | Plain-English"],
         "table_data": None, "speaker_notes": ""},
        {"slide_number": 11, "title": "The Answer: Yes, With Conditions",
         "bullets": [
             "Recommendation | adopt blend",
             "Conditions | monthly monitor + 2pp gate",
             "Risk Watch | HMM degrades sideways"],
         "table_data": None, "speaker_notes": ""},
    ]
    if overrides:
        for n, ov in overrides.items():
            base[n - 1].update(ov)
    pptx = build_presentation_deck(base, charts={})
    return Presentation(io.BytesIO(pptx))


def _all_text(shape) -> str:
    """Concatenate every run's text from a shape (or every cell's
    text from a table shape). Returns '' for non-text shapes (no
    text frame and no table)."""
    if getattr(shape, "has_table", False):
        parts: list[str] = []
        tbl = shape.table
        for r in range(len(tbl.rows)):
            for c in range(len(tbl.columns)):
                parts.append(tbl.cell(r, c).text_frame.text)
        return " ".join(parts).strip()
    try:
        tf = shape.text_frame
    except (AttributeError, ValueError):
        return ""
    parts2: list[str] = []
    for p in tf.paragraphs:
        for r in p.runs:
            parts2.append(r.text)
    return " ".join(parts2).strip()


def _slide_text(slide) -> str:
    """Concatenate every shape's text on a slide."""
    return " ".join(_all_text(s) for s in slide.shapes).strip()


# ── Dispatch ─────────────────────────────────────────────────────────


class TestDispatch:
    """_render_content_slide must dispatch by idx to the specialized
    renderers added in PR B. We verify by reading the slide-text on
    each slide: the title chrome / split panel / card grid / scorecard
    / feature rows each produce content the generic renderer would
    not."""

    def test_slide_1_uses_title_chrome(self):
        prs = _build_deck()
        text = _slide_text(prs.slides[0])
        assert "McColl School of Business" in text
        assert "Bob Thao" in text

    def test_slide_3_uses_split_panel(self):
        prs = _build_deck()
        text = _slide_text(prs.slides[2])
        assert "+85% Sharpe improvement vs benchmark" in text

    def test_slide_4_uses_card_grid(self):
        prs = _build_deck()
        # Card-grid headers from the '|' split land as bold text in
        # rounded rect shapes.
        text = _slide_text(prs.slides[3])
        assert "Correlation Break" in text
        assert "Regime Signal" in text

    def test_slide_8_uses_scorecard(self):
        prs = _build_deck()
        text = _slide_text(prs.slides[7])
        # The scorecard renderer produces the same table content,
        # but with colored cell fills (verified separately).
        assert "GFC 2008" in text and "Bank 2023" in text

    def test_slide_10_uses_feature_rows(self):
        prs = _build_deck()
        # Feature rows carry numeric badges 1 / 2 / 3.
        text = _slide_text(prs.slides[9])
        assert "1" in text and "2" in text and "3" in text
        assert "Regime Classifier" in text

    def test_slide_2_falls_through_to_generic(self):
        prs = _build_deck()
        # Generic agenda renderer produces a title-bar + a bullets
        # textbox with the items. No card-grid header chrome.
        text = _slide_text(prs.slides[1])
        assert "Agenda" in text


# ── Title slide ──────────────────────────────────────────────────────


class TestTitleSlide:

    def test_no_bullets_body(self):
        """The title slide does NOT render the body bullets even
        when the AI emitted any."""
        prs = _build_deck(overrides={
            1: {"bullets": ["A bullet that should NOT render"]},
        })
        text = _slide_text(prs.slides[0])
        assert "A bullet that should NOT render" not in text

    def test_carries_subtitle(self):
        prs = _build_deck()
        assert "McColl School of Business" in _slide_text(prs.slides[0])

    def test_carries_presenter_line(self):
        prs = _build_deck()
        assert "Molly Murdock" in _slide_text(prs.slides[0])

    def test_navy_header_bar_full_width(self):
        from tools.academic_deck import _MOLLY_NAVY, _SLIDE_W
        prs = _build_deck()

        def is_navy(s):
            try:
                return s.fill.fore_color.rgb == _MOLLY_NAVY
            except (AttributeError, TypeError):
                return False

        bars = [
            s for s in prs.slides[0].shapes
            if is_navy(s)
            and (s.width or 0) >= int(_SLIDE_W) - 1000
        ]
        assert bars, "title slide missing navy header bar"


# ── Split-panel slide ────────────────────────────────────────────────


class TestSplitPanelSlide:

    def test_has_oos_table(self):
        prs = _build_deck()
        tables = [
            s for s in prs.slides[2].shapes if s.has_table]
        assert tables, "slide 3 missing OOS results table"
        cols = [
            tables[0].table.cell(0, c).text_frame.text
            for c in range(len(tables[0].table.columns))]
        assert any("Sharpe" in c for c in cols)

    def test_regime_row_highlighted_teal(self):
        from tools.academic_deck import _MOLLY_TEAL
        prs = _build_deck()
        tables = [
            s for s in prs.slides[2].shapes if s.has_table]
        assert tables
        tbl = tables[0].table
        regime_row = None
        for r in range(1, len(tbl.rows)):
            if "regime" in tbl.cell(r, 0).text_frame.text.lower():
                regime_row = r
                break
        assert regime_row is not None
        # Every cell on the regime row must have the teal fill.
        for c in range(len(tbl.columns)):
            assert tbl.cell(regime_row, c).fill.fore_color.rgb == \
                _MOLLY_TEAL

    def test_callout_present(self):
        prs = _build_deck()
        text = _slide_text(prs.slides[2])
        assert "+85% Sharpe improvement vs benchmark" in text

    def test_left_cards_carry_strategy_names(self):
        prs = _build_deck()
        text = _slide_text(prs.slides[2])
        # Card headers come from the '|' split of each bullet.
        assert "Static (60/40)" in text
        assert "Dynamic (Regime-Conditional)" in text
        assert "Benchmark (100% S&P 500)" in text


# ── Card-grid slides ─────────────────────────────────────────────────


class TestCardGridSlide:
    """Slides 4 / 6 / 7 / 9 / 11. When SLIDE_CHARTS pins a chart for
    the slide (slide 4 today), cards stack vertically in the left
    50%; otherwise they spread horizontally."""

    def test_slide_4_chart_slot_empty_when_png_missing(self):
        # June 27 2026 -- spec: NO placeholder text in exported PPTX.
        # SLIDE_CHARTS[4] = rolling_correlation; with no chart png
        # supplied the renderer now logs + skips the chart slot
        # rather than emitting a [Chart unavailable] placeholder
        # textbox. Pin the absence of the placeholder text.
        prs = _build_deck()
        text = _slide_text(prs.slides[3])
        assert "Chart unavailable" not in text
        assert "[DATA PENDING]" not in text
        # The card content still renders even when chart is missing.
        assert "Correlation Break" in text

    def test_slide_6_three_horizontal_cards_no_chart(self):
        """No chart for slide 6 -> three side-by-side cards. Each
        card draws as a rounded-rectangle shape with header text."""
        prs = _build_deck()
        text = _slide_text(prs.slides[5])
        for hdr in ("53-month window", "Sharpe held",
                    "Walk-forward"):
            assert hdr in text, f"slide 6 missing card header: {hdr}"

    def test_handles_uneven_bullet_count(self):
        """Two bullets -> two cards. One bullet -> one card. No
        crashes on undershoot."""
        prs = _build_deck(overrides={
            7: {"bullets": [
                "VIX | low", "Yield Curve | flat"]}})
        text = _slide_text(prs.slides[6])
        assert "VIX" in text and "Yield Curve" in text


# ── Scorecard slide ──────────────────────────────────────────────────


class TestScorecardSlide:

    def test_bear_cells_colored_red(self):
        from tools.academic_deck import _SIGNAL_RED
        prs = _build_deck()
        tables = [s for s in prs.slides[7].shapes if s.has_table]
        assert tables
        tbl = tables[0].table
        # Find the council signal column (index 1 per the test
        # fixture: Scenario, Council Signal, Outcome).
        signal_col = 1
        bear_cells = [
            r for r in range(1, len(tbl.rows))
            if "BEAR" in tbl.cell(r, signal_col).text_frame.text
        ]
        assert bear_cells
        for r in bear_cells:
            assert tbl.cell(r, signal_col).fill.fore_color.rgb == \
                _SIGNAL_RED, f"row {r}: BEAR not red"

    def test_bull_cells_colored_green(self):
        from tools.academic_deck import _SIGNAL_GREEN
        prs = _build_deck()
        tbl = next(
            s for s in prs.slides[7].shapes
            if s.has_table).table
        bull_cells = [
            r for r in range(1, len(tbl.rows))
            if "BULL" in tbl.cell(r, 1).text_frame.text
        ]
        assert bull_cells
        for r in bull_cells:
            assert tbl.cell(r, 1).fill.fore_color.rgb == \
                _SIGNAL_GREEN

    def test_transition_cells_colored_amber(self):
        from tools.academic_deck import _SIGNAL_AMBER
        prs = _build_deck()
        tbl = next(
            s for s in prs.slides[7].shapes
            if s.has_table).table
        amber_cells = [
            r for r in range(1, len(tbl.rows))
            if "TRANSITION" in tbl.cell(r, 1).text_frame.text
        ]
        assert amber_cells
        for r in amber_cells:
            assert tbl.cell(r, 1).fill.fore_color.rgb == \
                _SIGNAL_AMBER

    def test_added_value_rows_shaded_teal(self):
        from tools.academic_deck import _TEAL_LIGHT
        prs = _build_deck()
        tbl = next(
            s for s in prs.slides[7].shapes
            if s.has_table).table
        # 'Added value' rows -- every cell EXCEPT the signal column
        # (which gets its own fill) must be _TEAL_LIGHT.
        for r in range(1, len(tbl.rows)):
            outcome = tbl.cell(r, 2).text_frame.text
            if "Added value" not in outcome:
                continue
            # Scenario column (col 0) reflects the row fill.
            assert tbl.cell(r, 0).fill.fore_color.rgb == \
                _TEAL_LIGHT, (
                    f"row {r} 'Added value' scenario cell not "
                    "teal-shaded")


# ── Feature-rows slide ───────────────────────────────────────────────


class TestFeatureRowsSlide:

    def test_three_oval_icons(self):
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = _build_deck()
        ovals = [
            s for s in prs.slides[9].shapes
            if (s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and _all_text(s) in ("1", "2", "3"))
        ]
        assert len(ovals) == 3, (
            f"expected 3 numeric icon badges, got {len(ovals)}")

    def test_feature_titles_present(self):
        prs = _build_deck()
        text = _slide_text(prs.slides[9])
        for title in (
                "Regime Classifier",
                "Strategy Comparison",
                "CIO Recommendation"):
            assert title in text, f"slide 10 missing feature: {title}"


# ── _table cell_style_hook ───────────────────────────────────────────


class TestTableCellStyleHook:
    """The hook signature is hardened so a missing / None / failing
    hook degrades gracefully to the default alternating-row fill.
    Tests probe the contract by exercising _table directly against
    a single-slide pptx."""

    @pytest.fixture
    def one_slide(self):
        from pptx import Presentation as _P
        from tools.academic_deck import _blank
        prs = _P()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        return prs, _blank(prs)

    def test_default_alternating_fill_unchanged(self, one_slide):
        from tools.academic_deck import _table
        prs, s = one_slide
        _table(s, ["A", "B"],
               [["1", "x"], ["2", "y"], ["3", "z"]],
               left=Inches(0.5), top=Inches(0.5),
               width=Inches(6))
        tbl = next(
            sh for sh in s.shapes if sh.has_table).table
        # Default fill: r % 2 -> _WHITE else slate (#F1F5F9).
        # Row 1 (r=1, r%2 truthy) -> _WHITE.
        from tools.academic_deck import _WHITE
        assert tbl.cell(1, 0).fill.fore_color.rgb == _WHITE

    def test_hook_per_cell_fill(self, one_slide):
        from pptx.dml.color import RGBColor
        from tools.academic_deck import _table
        prs, s = one_slide
        red = RGBColor(0xFF, 0x00, 0x00)

        def hook(*, row_idx, col_idx, header, value,
                 row_values, **_):
            if value == "TARGET":
                return {"fill": red, "bold": True}
            return None
        _table(s, ["A", "B"],
               [["1", "TARGET"], ["2", "y"]],
               left=Inches(0.5), top=Inches(0.5),
               width=Inches(6), cell_style_hook=hook)
        tbl = next(
            sh for sh in s.shapes if sh.has_table).table
        assert tbl.cell(1, 1).fill.fore_color.rgb == red
        assert tbl.cell(1, 0).fill.fore_color.rgb != red

    def test_hook_row_fill_applies_to_entire_row(self, one_slide):
        from pptx.dml.color import RGBColor
        from tools.academic_deck import _table
        prs, s = one_slide
        teal = RGBColor(0x02, 0x80, 0x90)

        def hook(*, row_idx, col_idx, header, value,
                 row_values, **_):
            if "ROW2" in row_values:
                return {"row_fill": teal}
            return None
        _table(s, ["A", "B", "C"],
               [["a", "b", "c"], ["ROW2", "y", "z"]],
               left=Inches(0.5), top=Inches(0.5),
               width=Inches(6), cell_style_hook=hook)
        tbl = next(
            sh for sh in s.shapes if sh.has_table).table
        # Every cell on the ROW2 row gets teal; other rows do not.
        for c in range(3):
            assert tbl.cell(2, c).fill.fore_color.rgb == teal
            assert tbl.cell(1, c).fill.fore_color.rgb != teal

    def test_hook_exception_degrades_to_default(self, one_slide):
        """A broken hook must not crash the deck -- the cell falls
        back to the default alternating fill."""
        from tools.academic_deck import _table, _WHITE
        prs, s = one_slide

        def hook(**_):
            raise ValueError("boom")

        _table(s, ["A", "B"],
               [["1", "x"]],
               left=Inches(0.5), top=Inches(0.5),
               width=Inches(6), cell_style_hook=hook)
        tbl = next(
            sh for sh in s.shapes if sh.has_table).table
        # Row 1 (r=1, odd) -> _WHITE default.
        assert tbl.cell(1, 0).fill.fore_color.rgb == _WHITE


# ── Canvas mirror ────────────────────────────────────────────────────


class TestCanvasMirror:
    """deck_slides_to_editor mirrors the PPTX dispatch for slides 1
    and 3 so the in-browser preview matches the export. Slides 4 /
    6 / 7 / 8 / 9 / 10 / 11 continue through the generic canvas
    builder because the existing text + chart + table element shape
    is already structurally faithful."""

    def test_slide_1_emits_title_chrome_no_body(self):
        from tools.editor_content import deck_slides_to_editor
        content_json, _ = deck_slides_to_editor([])
        s1 = content_json["slides"][0]
        ids = [e["id"] for e in s1["elements"]]
        assert "s1_title" in ids
        assert "s1_subtitle" in ids
        assert "s1_presenters" in ids
        assert "s1_body" not in ids
        assert all(e["type"] == "text" for e in s1["elements"])

    def test_slide_3_emits_split_panel(self):
        from tools.editor_content import deck_slides_to_editor
        content_json, _ = deck_slides_to_editor([{
            "slide_number": 3,
            "title": "The Investment Case",
            "bullets": ["A | B"],
            "table_data": {
                "headers": ["Strategy", "Sharpe"],
                "rows": [["Dynamic", "0.86"]],
            },
        }])
        s3 = content_json["slides"][2]
        body = next(
            e for e in s3["elements"] if e["id"] == "s3_body")
        table = next(
            e for e in s3["elements"] if e["id"] == "s3_table")
        # Body on the left half, table on the right half.
        assert body["x"] < table["x"]
        assert body["x"] + body["width"] <= table["x"]
