"""
tests/test_document_generation.py

Tests for the academic document-generation endpoints — the three graded
deliverables assembled from real platform data, light-mode charts and
AI-generated narrative:

  POST /api/v1/export/midpoint-paper      → 3-page midpoint paper (.docx)
  POST /api/v1/export/executive-brief     → 5-page executive brief (.docx)
  POST /api/v1/export/presentation-deck   → 6-slide final deck (.pptx)

Two tiers, the same pattern as test_export_package.py:
  - Endpoint-contract tests run everywhere including CI. In the test
    environment the analytics caches are cold and no academic documents
    are stored, so these tests double as the graceful-degradation tests:
    every section falls back to a [DATA PENDING] marker and the document
    still assembles into a valid, parseable file.
  - One DB round-trip confirms a document-generation run logs to
    agent_interactions for a team email and is gated out for a non-team
    email; it skips cleanly when no live PostgreSQL is reachable.
"""
from __future__ import annotations

import asyncio
import io
import os
import sys
import uuid

import pytest
from docx import Document
from fastapi.testclient import TestClient
from pptx import Presentation

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))
os.environ.setdefault("ENVIRONMENT", "test")
os.environ.setdefault("SECRET_KEY", "test-secret-key-at-least-32-characters-long")
os.environ.setdefault("MASTER_API_KEY", "test_master_key")
os.environ.setdefault(
    "ALLOWED_EMAILS",
    "ruurdsm@queens.edu,thaob@queens.edu,murdockm@queens.edu,panttserk@queens.edu",
)

from main import app  # noqa: E402
from auth import generate_session_token  # noqa: E402

client = TestClient(app)
TEAM_EMAIL = "ruurdsm@queens.edu"
NON_TEAM_EMAIL = "panttserk@queens.edu"
SESSION_HEADERS = {"X-API-Key": generate_session_token(TEAM_EMAIL)}

MIDPOINT = "/api/v1/export/midpoint-paper"
BRIEF = "/api/v1/export/executive-brief"
DECK = "/api/v1/export/presentation-deck"
APPENDIX = "/api/v1/export/analytical-appendix"

_DOCX_CT = "wordprocessingml"
_PPTX_CT = "presentationml"


def _run(coro):
    return asyncio.run(coro)


# ── Brief-grounding stub fixture (June 21 2026, PR #364) ──────────────────
#
# PR #364 added 409 gates on _generate_deck_document and
# _generate_appendix_document requiring a brief draft (and for
# the deck, also an appendix draft) to exist on the user's
# account. CI runs against a fresh DB with no drafts seeded, so
# the gates fire and the legitimate deck / appendix contract
# tests fail.
#
# Fix: an autouse fixture that monkeypatches the two grounding
# helpers to return stub payloads. The deck / appendix tests
# don't actually need real drafts -- they exercise the
# generator's downstream behaviour (PPTX assembly, eight-section
# appendix shape, etc.). The stub payloads carry valid content
# strings so the upstream Pass-1 Opus call composes a non-empty
# system prompt; the deck Pass-1 call still hits the test-env
# Anthropic short-circuit (which the existing tests already
# tolerated via [DATA PENDING] fallbacks).
#
# The fixture is intentionally NOT applied at module scope --
# pytest's monkeypatch fixture is function-scoped. Tests that
# want to exercise the GATE (i.e. the brief-grounding tests in
# test_brief_grounding.py) can stub the grounding helpers
# directly without this fixture interfering.


@pytest.fixture(autouse=True)
def _stub_brief_appendix_grounding(monkeypatch):
    """Stub brief + appendix grounding helpers so the 409 gates
    in _generate_deck_document / _generate_appendix_document do
    NOT fire during the document-generation contract tests. CI
    runs against a fresh DB with no drafts seeded; the stub
    bypasses the gate without requiring DB fixtures."""
    async def _fake_brief():
        return {
            "content_text": (
                "## Executive Summary\n\nThe blend outperforms "
                "benchmark on OOS Sharpe.\n\n"
                "## Methodology Overview\n\nWe use HMM regime "
                "detection.\n\n"
                "## Key Findings and Insights\n\nDrawdown "
                "reduction of 50% versus benchmark.\n\n"
                "## Limitations and Risks\n\nSample size 40 "
                "months.\n\n"
                "## Final Recommendations\n\nWe recommend the "
                "regime-conditional blend.\n\n"
                "## Visuals\n\nFour charts demonstrate the "
                "findings.\n"),
            "content_hash": "test_brief_hash",
            "draft_id": 1,
        }

    async def _fake_appendix():
        return {
            "content_text": (
                "## Data Sources and Methodology\n\nS&P 500 + "
                "AGG + HYG.\n\n"
                "## Portfolio Construction\n\nFull 10-strategy "
                "rules.\n\n"
                "## Calculations and Models\n\nFDR + Carhart + "
                "bootstrap.\n\n"
                "## Performance Metrics\n\n10-strategy table.\n\n"
                "## Sensitivity and Robustness\n\n10/15/20bp "
                "cost sensitivity.\n"),
            "content_hash": "test_appendix_hash",
            "draft_id": 2,
        }

    monkeypatch.setattr(
        "tools.brief_grounding.get_brief_for_grounding",
        _fake_brief)
    monkeypatch.setattr(
        "tools.brief_grounding.get_appendix_for_grounding",
        _fake_appendix)

    # June 21 2026 (PR #365) -- pre-flight cache gate on appendix
    # generation 409s when strategy_results / bootstrap_ci_sharpe /
    # factor_loadings / cost_sensitivity are empty. CI has cold
    # caches so the gate fires and the legitimate appendix
    # content tests fail with 409 instead of producing the
    # artifact. Stub the data-gather fn so the gate sees populated
    # caches; downstream behaviour (8-section assembly,
    # reproducibility line, editor-draft creation) still
    # exercises against the [DATA PENDING] fallbacks where the
    # generators haven't been stubbed.
    async def _fake_appendix_data(data_hash=None):
        return {
            "available": True,
            "study_period": {"start": "2002-07-31",
                             "end": "2026-05-31",
                             "n_months": 287},
            "summary_statistics": [],
            "regime_conditional": [],
            "drawdown_comparison": [],
            "factor_loadings": [{"strategy": "STUB"}],
            "cumulative_returns": {},
            "rolling_correlation": {},
            "strategy_results": {"STUB": {"sharpe_ratio": 0.5}},
            "strategy_metadata": {},
            "risk_free_rate": None,
            "team_summary": {},
            "last_review_text": None,
            "academic_docs": [],
            "audit_disclosures": None,
            "bootstrap_ci_sharpe": [{"strategy": "STUB"}],
            "crisis_performance": None,
            "cost_sensitivity": {"net_sharpe_15bp": 0.5},
            "invariant_summary": None,
            "data_hash": "test_hash_for_appendix_fixture",
        }

    monkeypatch.setattr(
        "tools.academic_export.gather_analytical_appendix_data",
        _fake_appendix_data)
    # The appendix call site re-imports from the same path inside
    # the function, so the monkeypatch at the import name is
    # sufficient.

    # June 22 2026 (hash-match strictness PR) -- the appendix
    # pre-flight gate now ALSO computes the canonical strategy
    # hash via _compute_data_hash(n_rows, last_date, 10) and
    # verifies strategy_results_cache + each analytics metric
    # carries a row AT THAT HASH. Stub the inputs so the gate
    # passes for these document-generation contract tests
    # (which exercise downstream behaviour, not the gate).
    import pandas as pd

    async def _fake_history_async():
        dates = pd.date_range(
            end="2026-05-31", periods=287, freq="ME")
        equity = pd.Series([1.0] * 287, index=dates)
        return {"equity_monthly": equity}

    monkeypatch.setattr(
        "tools.data_fetcher.get_full_history_async",
        _fake_history_async)

    async def _fake_strategy_cache(_h):
        return {"STUB": {"sharpe_ratio": 0.5}}

    async def _fake_metric_by_hash(kind, _h):
        if kind == "academic_analytics":
            return {"bootstrap_ci_sharpe": [{"strategy": "STUB"}],
                    "factor_loadings": [{"strategy": "STUB"}]}
        if kind == "oos_cost_sensitivity":
            return {"net_sharpe_15bp": 0.5}
        return None

    monkeypatch.setattr(
        "tools.cache.get_strategy_cache", _fake_strategy_cache)
    monkeypatch.setattr(
        "tools.precomputed_analytics.get_metric_by_hash",
        _fake_metric_by_hash)


def _docx_text(content: bytes) -> str:
    """All header, paragraph and table text from a .docx, for content checks."""
    doc = Document(io.BytesIO(content))
    parts: list[str] = []
    for section in doc.sections:
        parts.extend(p.text for p in section.header.paragraphs)
    parts.extend(p.text for p in doc.paragraphs)
    for table in doc.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def _pptx_text(content: bytes) -> str:
    """All shape text from a .pptx, for [DATA PENDING] checks."""
    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.extend(c.text for c in row.cells)
    return "\n".join(parts)


# ── Endpoint contract — runs in CI ────────────────────────────────────────────

class TestDocumentGenerationContract:
    # The endpoints are async: a POST creates a job and returns 202 with
    # a job_id; the file is produced by a background task. The document
    # CONTENT checks call the generation helpers directly — the
    # background task does not complete under Starlette's TestClient.

    def test_midpoint_paper_returns_410_gone_after_retirement(self):
        # PR-B (June 2026) retired the midpoint pipeline. The endpoint
        # is preserved as a 410 stub so existing clients receive a
        # clear "this existed and is now gone" signal rather than a
        # 404 connection error.
        resp = client.post(MIDPOINT, headers=SESSION_HEADERS)
        assert resp.status_code == 410
        body = resp.json()
        assert body["error"] == "gone"
        assert "Midpoint paper generation has been retired" \
            in body["message"]
        assert body["canonical_path"] == (
            "/api/v1/export/executive-brief")

    def test_executive_brief_returns_202_job(self):
        resp = client.post(BRIEF, headers=SESSION_HEADERS)
        assert resp.status_code == 202
        assert resp.json()["status"] == "pending"
        assert resp.json()["job_id"]

    def test_presentation_deck_returns_202_job(self):
        resp = client.post(DECK, headers=SESSION_HEADERS)
        assert resp.status_code == 202
        assert resp.json()["status"] == "pending"
        assert resp.json()["job_id"]

    # PR-B (June 2026) deleted _generate_midpoint_document alongside
    # the midpoint endpoint retirement; the corresponding docx-render
    # smoke test is gone with it.

    def test_executive_brief_document_is_a_valid_docx_with_headings(self):
        """June 18 2026 -- the brief was rewritten to the FNA 670
        rubric's six required sections in rubric order. The earlier
        (June 6) structure carried non-rubric "Five Human Decisions"
        and "Part II preview" sections; rubric review against the
        FNA 670 spec retired both. The six section headings below
        match the rubric exactly."""
        import main
        docx_bytes, filename, media, _draft = _run(
            main._generate_brief_document(TEAM_EMAIL))
        assert _DOCX_CT in media
        assert filename.endswith(".docx")
        text = _docx_text(docx_bytes)
        # All six new section headings present in order.
        assert "1. Executive Summary" in text
        assert "2. Methodology Overview" in text
        assert "3. Key Findings and Insights" in text
        assert "4. Limitations and Risks" in text
        assert "5. Final Recommendations" in text
        assert "6. Visuals to Demonstrate the Insights" in text
        # And the section order is preserved.
        idx = [text.index(h) for h in (
            "1. Executive Summary",
            "2. Methodology Overview",
            "3. Key Findings and Insights",
            "4. Limitations and Risks",
            "5. Final Recommendations",
            "6. Visuals to Demonstrate the Insights",
        )]
        assert idx == sorted(idx)
        # The retired pre-rebuild [[BOB]] callouts must not creep back.
        for banned in (
            "BOB — YOUR FRAMING",
            "BOB — YOUR JUDGEMENT",
            "BOB — YOUR RECOMMENDATION",
        ):
            assert banned not in text, f"unexpected placeholder: {banned}"
        assert "[[BOB]]" not in text
        assert "[DATA PENDING]" in text

    def test_presentation_deck_document_is_a_valid_11_slide_pptx(self):
        # June 27 2026 -- back to 11 slides after the Investment
        # Case merge collapsed the old slide 3+4 setup + verdict
        # into one split-panel slide to match Molly's reference
        # deck. (Was 12 between 2026-06-22 agenda insert and
        # the 2026-06-27 collapse.)
        import main
        from tools.academic_deck import DECK_SLIDE_COUNT
        pptx_bytes, filename, media, _draft = _run(
            main._generate_deck_document(TEAM_EMAIL))
        assert _PPTX_CT in media
        assert filename.endswith(".pptx")
        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) == DECK_SLIDE_COUNT
        assert DECK_SLIDE_COUNT == 11  # June 27 2026 Molly-aligned collapse
        # Cold caches / no matplotlib in the test env must not fail the deck.
        assert "[DATA PENDING]" in _pptx_text(pptx_bytes)

    def test_presentation_deck_has_canonical_titles_and_notes(self):
        import main
        from tools.academic_deck import SLIDE_TITLES
        pptx_bytes, *_ = _run(main._generate_deck_document(TEAM_EMAIL))
        prs = Presentation(io.BytesIO(pptx_bytes))
        text = _pptx_text(pptx_bytes)
        # The six canonical slide titles are always present (the builder
        # falls back to them when the AI JSON is absent — the test env case).
        for title in SLIDE_TITLES:
            assert title in text, f"missing slide title: {title}"
        # Every slide carries non-empty speaker notes (the verify caveat at
        # minimum); slide 1 additionally carries the submission checklist.
        for s in prs.slides:
            assert s.has_notes_slide
            assert s.notes_slide.notes_text_frame.text.strip()
        assert "SUBMISSION CHECKLIST" in \
            prs.slides[0].notes_slide.notes_text_frame.text

    def test_build_presentation_deck_embeds_charts_on_4_5_11(self):
        # June 7 2026 (bridges #98 / #100) -- chart-bearing slides
        # are 4, 5, 11 per the eleven-slide rebuild (rolling
        # correlation on slide 4, OOS Sharpe / cumulative comparison
        # proxy on slide 5, efficient frontier with the live blend
        # marker on slide 11). The other eight slides carry stat
        # cards, comparison tables, or pure prose -- no matplotlib.
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from tools.academic_deck import (
            DECK_SLIDE_COUNT, SLIDE_CHARTS, build_presentation_deck)
        from tools.chart_render import _placeholder
        slides = [{"slide_number": n, "title": f"T{n}", "bullets": ["b1", "b2"],
                   "table_data": None, "speaker_notes": f"notes {n}"}
                  for n in range(1, DECK_SLIDE_COUNT + 1)]
        # June 26 2026 -- build a DISTINCT PNG per slide. The
        # production renderer (_render_deck_slide_charts) produces
        # three different matplotlib outputs (rolling_correlation,
        # strategy_comparison, efficient_frontier), one per
        # SLIDE_CHARTS role. The bytes-identity dedup added by
        # the same PR that this test exercises would otherwise
        # collapse all three placeholder PNGs to a single slide
        # under the 'same bytes -> one slot' rule. The dedup is
        # working as intended; the test setup is updated to match
        # the production shape.
        charts = {
            n: _placeholder(240 + i, 150)
            for i, n in enumerate(sorted(SLIDE_CHARTS))
        }
        out = build_presentation_deck(slides, charts)
        prs = Presentation(io.BytesIO(out))
        assert len(prs.slides) == DECK_SLIDE_COUNT

        def _has_pic(s):
            return any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                       for sh in s.shapes)
        # Chart slides per SLIDE_CHARTS keys.
        for n in SLIDE_CHARTS:
            assert _has_pic(prs.slides[n - 1]), f"slide {n} missing chart"
        # No-chart slides — all slides 1..N not in SLIDE_CHARTS.
        chart_keys = set(SLIDE_CHARTS.keys())
        for n in range(1, DECK_SLIDE_COUNT + 1):
            if n in chart_keys:
                continue
            assert not _has_pic(prs.slides[n - 1]), (
                f"slide {n} unexpected chart")
        for s in prs.slides:
            assert s.notes_slide.notes_text_frame.text.strip()

    def test_analytical_appendix_returns_202_job(self):
        resp = client.post(APPENDIX, headers=SESSION_HEADERS)
        assert resp.status_code == 202
        assert resp.json()["status"] == "pending"
        assert resp.json()["job_id"]

    def test_analytical_appendix_document_has_eight_sections(self):
        """The appendix must carry every one of the eight sections
        regardless of cache state — a cold deploy renders the same
        document shell with [DATA PENDING] markers in place of live
        figures, so every section heading appears verbatim."""
        import main
        docx_bytes, filename, media, _draft = _run(
            main._generate_appendix_document(TEAM_EMAIL))
        assert _DOCX_CT in media
        assert filename.endswith(".docx")
        assert "analytical-appendix" in filename
        text = _docx_text(docx_bytes)
        section_headings = [
            "A. Data and Methodology",
            "B. Full Strategy Performance",
            "C. Statistical Tests",
            "D. Bootstrap Confidence Intervals",
            "E. Factor Loadings",
            "F. Crisis Window Performance",
            "G. Transaction Cost Sensitivity",
            "H. Validation Audit Summary",
        ]
        for h in section_headings:
            assert h in text, f"missing section heading: {h}"
        # And the section order is preserved A → H.
        idx = [text.index(h) for h in section_headings]
        assert idx == sorted(idx), (
            f"section headings out of order: {idx}")
        assert "AI DRAFT" in text                # mandatory banner
        assert "Submission Checklist" in text    # CAVEAT 5

    def test_analytical_appendix_carries_reproducibility_line(self):
        """The Reproducibility line must surface the data hash + the
        generation timestamp. In the test env the cache is cold so the
        hash renders as [DATA PENDING]; the heading and the
        explanatory copy are still present."""
        import main
        docx_bytes, *_ = _run(
            main._generate_appendix_document(TEAM_EMAIL))
        text = _docx_text(docx_bytes)
        assert "Reproducibility" in text
        assert "Data hash:" in text
        assert "Generated at:" in text
        # The cache traceability sentence — the part of the line that
        # actually explains what the hash anchors.
        assert "strategy_results_cache" in text

    def test_analytical_appendix_renders_data_hash_when_supplied(self):
        """Build the document directly to confirm a supplied data hash
        appears verbatim. The endpoint reads the hash from the cache;
        the builder renders whatever it is given."""
        from tools.academic_docx import build_analytical_appendix
        narratives = {
            f"appendix_{ltr}": f"Intro for section {ltr.upper()}."
            for ltr in "abcdefgh"
        }
        data = {
            "summary_statistics": [],
            "strategy_results": {},
            "bootstrap_ci_sharpe": [],
            "factor_loadings": [],
            "crisis_performance": None,
            "cost_sensitivity": None,
            "invariant_summary": None,
            "audit_disclosures": None,
            "data_hash": "deadbeef12345678",
        }
        docx_bytes = build_analytical_appendix(data, narratives)
        text = _docx_text(docx_bytes)
        assert "deadbeef12345678" in text

    def test_analytical_appendix_creates_editor_draft_or_falls_open(self):
        """The endpoint creates an editor_drafts row alongside the
        .docx so the document opens in the in-platform editor. In the
        test env with no live DB the create_draft call fails open and
        the helper returns draft_id=None — the file is still produced.
        Either outcome (an integer id or None) is acceptable; what is
        NOT acceptable is the helper raising."""
        import main
        _bytes, _name, _media, draft_id = _run(
            main._generate_appendix_document(TEAM_EMAIL))
        assert draft_id is None or isinstance(draft_id, int)

    def test_all_four_require_authentication(self):
        for endpoint in (MIDPOINT, BRIEF, DECK, APPENDIX):
            assert client.post(endpoint).status_code == 401


# ── DB round-trip — skips without a live database ─────────────────────────────

_db_ready_cache: bool | None = None


async def _fresh_session():
    """Disposes the pooled engine and returns a session on the current loop."""
    from database import engine, AsyncSessionLocal
    if engine is not None:
        await engine.dispose()
    return AsyncSessionLocal()  # type: ignore[union-attr]


def _db_ready() -> bool:
    """True when a live PostgreSQL with the activity tables is reachable."""
    global _db_ready_cache
    if _db_ready_cache is not None:
        return _db_ready_cache
    try:
        from tools.cache import _DB_AVAILABLE
        if not _DB_AVAILABLE:
            _db_ready_cache = False
            return False
        from sqlalchemy import text

        async def _probe() -> bool:
            async with await _fresh_session() as s:
                await s.execute(text("SELECT 1 FROM agent_interactions LIMIT 1"))
            return True

        _db_ready_cache = _run(_probe())
    except Exception:
        _db_ready_cache = False
    return _db_ready_cache


class TestPMAudienceEvaluator:
    """The PM audience evaluator fires as a SECOND evaluator pass on
    every document-section narrative — midpoint paper, executive brief,
    deck narrative blocks all flow through harness_narrative(). The
    presentation script writer, council specialists, and triage agent
    do NOT use the PM evaluator (different audiences). This class pins
    the prompt content and the wiring contract."""

    def test_pm_evaluator_prompt_contains_all_five_criteria(self):
        # The five PM criteria must all be named in the prompt — a
        # missing criterion would let the evaluator silently drop a
        # whole dimension of the rubric.
        from agents.evaluator_prompts import academic_export_evaluator_pm_prompt
        prompt = academic_export_evaluator_pm_prompt()
        for marker in (
            "PM_CRITERION_1", "PM_CRITERION_2", "PM_CRITERION_3",
            "PM_CRITERION_4", "PM_CRITERION_5",
            "INSIGHT BEYOND THE OBVIOUS",
            "MECHANISM NOT JUST OBSERVATION",
            "ACTIONABLE SIGNAL IDENTIFICATION",
            "CONTRADICTIONS ACKNOWLEDGED AND PRESSED",
            "SO WHAT / EXPLICIT IMPLICATION",
        ):
            assert marker in prompt, f"missing PM criterion marker: {marker}"

    def test_pm_evaluator_prompt_emits_verdict_to_overall_mapping(self):
        # The harness reads `overall` as a number; the PM verdict must
        # map STRONG → 9.0, DEVELOPING → 7.5, NEEDS WORK → 3.0 so the
        # 7.0 threshold retries only on NEEDS WORK.
        from agents.evaluator_prompts import academic_export_evaluator_pm_prompt
        prompt = academic_export_evaluator_pm_prompt()
        assert "9.0 (STRONG)" in prompt
        assert "7.5 (DEVELOPING)" in prompt
        assert "3.0 (NEEDS WORK)" in prompt

    def test_harness_narrative_passes_pm_evaluator_as_secondary(self):
        # Wiring contract — harness_narrative supplies the PM evaluator
        # as the secondary so every document section is dual-evaluated.
        # We patch the harness class to capture how it was called.
        from unittest.mock import MagicMock, patch
        from tools import academic_export
        from agents.harness import HarnessResult

        captured: dict = {}

        class _FakeHarness:
            def __init__(self):
                pass

            def run(self, *args, **kwargs):
                captured["kwargs"] = kwargs
                return HarnessResult(
                    response="ok", final_score=9.0, attempts=1,
                    improved=False, feedback_applied="", initial_score=9.0,
                )

        # The function early-exits in the test environment; bypass that
        # guard so the harness path actually runs.
        with patch.object(academic_export, "ENVIRONMENT", "production"), \
             patch("agents.harness.GeneratorEvaluatorHarness", _FakeHarness):
            academic_export.harness_narrative(
                "midpoint_paper_intro", "draft this section", {"x": 1})

        # The secondary evaluator must be the PM prompt — same string
        # the prompt builder returns. Comparing equality rather than
        # 'contains' so a regression that swaps the prompt is caught.
        from agents.evaluator_prompts import academic_export_evaluator_pm_prompt
        assert captured["kwargs"].get("secondary_evaluator_prompt") \
            == academic_export_evaluator_pm_prompt()

    def test_presentation_script_does_NOT_pass_pm_evaluator(self):
        # The script generator passes only the presentation_script
        # evaluator — no secondary. Spoken delivery is a different
        # audience; the PM rubric doesn't apply to a 16-slide oral
        # script the way it applies to a written analytical document.
        from unittest.mock import patch
        from tools import script_generation
        from agents.harness import HarnessResult

        captured: dict = {}

        class _FakeHarness:
            def __init__(self):
                pass

            def run(self, *args, **kwargs):
                captured["kwargs"] = kwargs
                return HarnessResult(
                    response="script body", final_score=9.0, attempts=1,
                    improved=False, feedback_applied="", initial_score=9.0,
                )

        with patch("agents.harness.GeneratorEvaluatorHarness", _FakeHarness):
            script_generation._run_harness(
                slides=[{"slide_number": 1, "title": "T", "speaker": "Molly",
                         "content_text": "x"}],
                exec_brief_text=None, midpoint_text=None)

        # Either the kwarg is absent OR it is explicitly None.
        secondary = captured["kwargs"].get("secondary_evaluator_prompt")
        assert secondary is None


class TestAcademicWriterAudiencePrompt:
    """The Academic Writer system prompt sets the audience expectation —
    primary reader is a portfolio manager who wants insight beyond
    standard metrics. Pins the audience guidance + the so-what
    instruction so a future prompt edit doesn't silently drop them."""

    def test_writer_prompt_names_portfolio_manager_audience(self):
        from agents.academic_writer import _SYSTEM_PROMPT
        assert "PORTFOLIO MANAGER" in _SYSTEM_PROMPT
        assert "primary reader" in _SYSTEM_PROMPT.lower()

    def test_writer_prompt_requires_explicit_so_what(self):
        from agents.academic_writer import _SYSTEM_PROMPT
        # "so what?" + "implication" are the load-bearing phrases.
        assert "so what" in _SYSTEM_PROMPT.lower()
        assert "implication" in _SYSTEM_PROMPT.lower()

    def test_writer_prompt_names_the_2022_break_and_contradictions(self):
        from agents.academic_writer import _SYSTEM_PROMPT
        assert "2022" in _SYSTEM_PROMPT
        assert "contradict" in _SYSTEM_PROMPT.lower()


class TestArbiterDualVerdict:
    """The Academic Review arbiter verdict opens with TWO top-level
    summary lines (Academic rigour + Portfolio Manager insight) so the
    user sees both lenses at a glance. The five rubric sections still
    follow."""

    def test_arbiter_prompt_requires_two_top_level_verdict_lines(self):
        from agents.academic_review import _ARBITER_INSTRUCTIONS
        # Both lenses are named in the instructions block.
        assert "**Academic rigour:**" in _ARBITER_INSTRUCTIONS
        assert "**Portfolio Manager insight:**" in _ARBITER_INSTRUCTIONS

    def test_arbiter_prompt_lists_pm_criteria(self):
        # The arbiter must apply the same five PM criteria the harness
        # uses on document sections, so its PM verdict is consistent
        # with the writer's dual-evaluator feedback.
        from agents.academic_review import _ARBITER_INSTRUCTIONS
        for marker in (
            "Insight beyond the obvious",
            "The 2022 break",
            "Actionable signal identification",
            "Contradictions acknowledged and pressed",
            "So what / explicit implication",
        ):
            assert marker in _ARBITER_INSTRUCTIONS, (
                f"arbiter prompt missing PM criterion: {marker}")

    def test_arbiter_prompt_still_lists_five_rubric_sections(self):
        # PR #194 renamed sections 1-4 to match the FNA 670 midpoint
        # rubric (Data and Methodology / Preliminary Results and
        # Diagnostics / Roles and Division of Labor / Next Steps and
        # Open Questions). Section 5 keeps the literal title "Overall
        # Academic Readiness" so the existing truncation detector and
        # fallback (UAT #53/#59/#125/#128) keep working.
        from agents.academic_review import _ARBITER_INSTRUCTIONS
        for marker in (
            "### 1. Data and Methodology (1p, 33%)",
            "### 2. Preliminary Results and Diagnostics (1p, 33%)",
            "### 3. Roles and Division of Labor (0.5p, 17%)",
            "### 4. Next Steps and Open Questions (0.5p, 17%)",
            "### 5. Overall Academic Readiness",
        ):
            assert marker in _ARBITER_INSTRUCTIONS


class TestDocumentGenerationLogging:
    def test_document_generation_logged_and_team_gated(self):
        """Document generation records a run via log_agent_interaction with
        interaction_type='export' and a deliverable in metadata. Exercises
        that data layer directly — the single-loop pattern test_activity.py
        uses — so it verifies a team user's generation logs a row and a
        non-team user is gated out."""
        if not _db_ready():
            pytest.skip("no live database")
        from tools.activity_log import log_agent_interaction
        from sqlalchemy import text

        async def scenario():
            from database import engine, AsyncSessionLocal
            await engine.dispose()
            sid = str(uuid.uuid4())
            try:
                ok = await log_agent_interaction(
                    user_email=TEAM_EMAIL, session_id=sid,
                    session_type="analytical", interaction_type="export",
                    agents_involved=["academic_writer"],
                    response_summary="Midpoint paper generated",
                    metadata={"deliverable": "midpoint_paper"})
                assert ok is True
                ok2 = await log_agent_interaction(
                    user_email=NON_TEAM_EMAIL, session_id=sid,
                    session_type="analytical", interaction_type="export",
                    metadata={"deliverable": "midpoint_paper"})
                assert ok2 is False
                async with AsyncSessionLocal() as s:
                    row = await s.execute(
                        text("SELECT COUNT(*) FROM agent_interactions "
                             "WHERE session_id = :sid AND "
                             "interaction_type = 'export'"),
                        {"sid": sid})
                    assert row.scalar() == 1
            finally:
                async with AsyncSessionLocal() as s:
                    await s.execute(
                        text("DELETE FROM agent_interactions WHERE session_id = :sid"),
                        {"sid": sid})
                    await s.commit()

        _run(scenario())
