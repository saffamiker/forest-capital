"""
tests/test_editor_drafts.py

Tests for the in-platform document editor — the editor_drafts /
editor_draft_versions data layer (migration 021), the
/api/v1/documents/drafts endpoints, the draft created on document
generation, and Academic Review reading the editor draft.

Endpoint-contract tests (auth gating) run everywhere. The CRUD
round-trips need a live database and skip cleanly without one — the
same pattern as the rest of the suite. Every DB test uses the
clean_editor_drafts fixture so it leaves no rows behind.
"""
from __future__ import annotations

import asyncio
import os
import sys

import pytest
from fastapi.testclient import TestClient

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))
os.environ.setdefault("ENVIRONMENT", "test")
os.environ.setdefault("SECRET_KEY", "test-secret-key-at-least-32-characters-long")
os.environ.setdefault("MASTER_API_KEY", "test_master_key")
os.environ.setdefault(
    "ALLOWED_EMAILS",
    "ruurdsm@queens.edu,thaob@queens.edu,murdockm@queens.edu,panttserk@queens.edu",
)

from main import app  # noqa: E402
from auth import generate_session_token  # noqa: E402

client = TestClient(app)
TEAM = {"X-API-Key": generate_session_token("thaob@queens.edu")}     # team member
VIEWER = {"X-API-Key": generate_session_token("panttserk@queens.edu")}  # viewer
DRAFTS = "/api/v1/documents/drafts"


def _run(coro):
    return asyncio.run(coro)


_db_ready_cache: bool | None = None


def _db_ready() -> bool:
    global _db_ready_cache
    if _db_ready_cache is not None:
        return _db_ready_cache
    try:
        from tools.cache import _DB_AVAILABLE
        if not _DB_AVAILABLE:
            _db_ready_cache = False
            return False
        from sqlalchemy import text
        from database import engine, AsyncSessionLocal

        async def _probe() -> bool:
            if engine is not None:
                await engine.dispose()
            async with AsyncSessionLocal() as s:  # type: ignore[union-attr]
                await s.execute(text("SELECT 1 FROM editor_drafts LIMIT 1"))
            return True

        _db_ready_cache = _run(_probe())
    except Exception:
        _db_ready_cache = False
    return _db_ready_cache


# ── Endpoint gating ───────────────────────────────────────────────────────────

class TestEditorEndpointGating:
    def test_list_drafts_requires_auth(self):
        assert client.get(DRAFTS).status_code == 401

    def test_list_drafts_rejects_a_viewer(self):
        # A viewer has no team_member permission.
        assert client.get(DRAFTS, headers=VIEWER).status_code == 403

    def test_list_drafts_admits_a_team_member(self):
        resp = client.get(DRAFTS, headers=TEAM)
        assert resp.status_code == 200
        assert "drafts" in resp.json()

    def test_create_rejects_a_bad_document_type(self):
        resp = client.post(DRAFTS, headers=TEAM,
                            json={"document_type": "novel", "title": "x"})
        assert resp.status_code == 422


# ── Marker helpers — pure logic ───────────────────────────────────────────────

class TestEditorContentBuilders:
    def test_midpoint_to_editor_builds_tiptap_and_text(self):
        from tools.editor_content import midpoint_to_editor
        cj, ct = midpoint_to_editor({
            "methodology": "Method para.", "results": "Results para.",
            "roles": "Roles para.", "next_steps": "Next steps para."})
        assert cj["type"] == "doc"
        # Four H1 section headings.
        headings = [n for n in cj["content"] if n.get("type") == "heading"]
        assert len(headings) == 4
        # The Roles and Next Steps [[BOB]] callouts are embedded.
        assert ct.count("[[BOB:") == 2

    def test_deck_to_editor_builds_sixteen_canvas_slides(self):
        from tools.editor_content import deck_to_editor
        cj, ct = deck_to_editor({"conclusions": "- one", "thesis": "A thesis."})
        assert len(cj["slides"]) == 16
        # Canvas schema (migration 022) — every slide carries an elements
        # array, a background, and empty speaker notes.
        for s in cj["slides"]:
            assert isinstance(s["elements"], list)
            assert s["background"] == "#FFFFFF"
            assert s["speaker_notes"] == ""
        assert "Slide 1:" in ct

    def test_deck_to_editor_populates_every_slide_body(self):
        # Every one of the 16 slides carries a non-empty body text element
        # — the five narrative-keyed slides from the generated narratives,
        # the rest from the static seed. No slide opens blank.
        from tools.editor_content import deck_to_editor
        cj, _ = deck_to_editor({
            "thesis": "A thesis.", "conclusions": "- one",
            "recommendations": "- rec", "ai_leverage": "AI narrative."})
        assert len(cj["slides"]) == 16
        for s in cj["slides"]:
            body = next(e for e in s["elements"] if e["id"] == "el_002")
            assert body["type"] == "text"
            assert body["content"].strip()

    def test_deck_to_editor_seeds_keyless_slides_when_narratives_empty(self):
        # Even with no narratives at all, every slide still has a body
        # element seeded from its static description.
        from tools.editor_content import deck_to_editor
        cj, _ = deck_to_editor({})
        for s in cj["slides"]:
            body = next(e for e in s["elements"] if e["id"] == "el_002")
            assert body["content"].strip()

    def test_executive_brief_to_editor_builds_tiptap_and_text(self):
        from tools.editor_content import executive_brief_to_editor
        cj, ct = executive_brief_to_editor({
            "exec_summary": "Summary para.", "methodology": "Method para.",
            "finding_1": "F1.", "finding_2": "F2.", "finding_3": "F3.",
            "finding_4": "F4.", "limitations": "Limits.",
            "recommendations": "Recs."})
        assert cj["type"] == "doc"
        # Eight H1 section headings.
        headings = [n for n in cj["content"] if n.get("type") == "heading"]
        assert len(headings) == 8
        # The brief carries no [[BOB]] callouts.
        assert "[[BOB:" not in ct
        assert "Executive Summary" in ct


# ── CRUD round-trips — skip without a live database ───────────────────────────

class TestEditorCRUD:
    def test_create_get_patch_version_restore_delete(self, clean_editor_drafts):
        if not _db_ready():
            pytest.skip("no live database with the migration-021 tables")

        # Create — word count is computed from content_text.
        created = client.post(DRAFTS, headers=TEAM, json={
            "document_type": "midpoint_paper", "title": "Round-trip",
            "content_json": {"type": "doc", "content": []},
            "content_text": "one two three"})
        assert created.status_code == 201
        draft = created.json()
        did = draft["id"]
        assert draft["word_count"] == 3
        assert draft["version"] == 1

        # Get.
        got = client.get(f"{DRAFTS}/{did}", headers=TEAM)
        assert got.status_code == 200 and got.json()["id"] == did

        # Auto-save (PATCH) — updates content, word count updates on edit,
        # and does NOT create a version.
        patched = client.patch(f"{DRAFTS}/{did}", headers=TEAM, json={
            "content_text": "one two three four five"})
        assert patched.status_code == 200
        after = client.get(f"{DRAFTS}/{did}", headers=TEAM).json()
        assert after["word_count"] == 5
        assert after["version"] == 1   # auto-save creates no version
        assert client.get(f"{DRAFTS}/{did}/versions",
                           headers=TEAM).json()["versions"] == []

        # Save a named version checkpoint.
        ver = client.post(f"{DRAFTS}/{did}/versions", headers=TEAM,
                          json={"version_label": "checkpoint one"})
        assert ver.status_code == 201
        version_id = ver.json()["id"]
        assert ver.json()["version_label"] == "checkpoint one"
        versions = client.get(f"{DRAFTS}/{did}/versions",
                              headers=TEAM).json()["versions"]
        assert len(versions) == 1

        # Edit again, then restore the saved version.
        client.patch(f"{DRAFTS}/{did}", headers=TEAM,
                     json={"content_text": "completely different text now"})
        restored = client.post(f"{DRAFTS}/{did}/restore/{version_id}",
                                headers=TEAM)
        assert restored.status_code == 200
        assert restored.json()["content_text"] == "one two three four five"

        # Soft delete — the draft then 404s.
        assert client.delete(f"{DRAFTS}/{did}",
                             headers=TEAM).status_code == 200
        assert client.get(f"{DRAFTS}/{did}", headers=TEAM).status_code == 404

    def test_create_sets_is_current_and_unsets_siblings(self, clean_editor_drafts):
        if not _db_ready():
            pytest.skip("no live database")
        first = client.post(DRAFTS, headers=TEAM, json={
            "document_type": "executive_brief", "title": "First",
            "content_text": "first"}).json()
        second = client.post(DRAFTS, headers=TEAM, json={
            "document_type": "executive_brief", "title": "Second",
            "content_text": "second"}).json()
        # The newest draft of a type is current; the previous one is not.
        drafts = {d["id"]: d for d in
                  client.get(DRAFTS, headers=TEAM).json()["drafts"]}
        assert drafts[second["id"]]["is_current"] is True
        assert drafts[first["id"]]["is_current"] is False


class TestDraftOnGeneration:
    # Generation is async — the endpoint returns 202 and a background
    # task creates the draft. The generation helper is exercised directly
    # (it creates the editor draft and returns its id).
    def test_midpoint_generation_creates_a_draft(self, clean_editor_drafts):
        if not _db_ready():
            pytest.skip("no live database")
        import main
        _bytes, _fn, _media, draft_id = _run(
            main._generate_midpoint_document("thaob@queens.edu"))
        assert draft_id is not None
        got = client.get(f"{DRAFTS}/{draft_id}", headers=TEAM)
        assert got.status_code == 200
        assert got.json()["document_type"] == "midpoint_paper"
        assert got.json()["created_from"] == "generated"

    def test_executive_brief_generation_creates_a_draft(
        self, clean_editor_drafts,
    ):
        if not _db_ready():
            pytest.skip("no live database")
        import main
        _bytes, _fn, _media, draft_id = _run(
            main._generate_brief_document("thaob@queens.edu"))
        assert draft_id is not None
        got = client.get(f"{DRAFTS}/{draft_id}", headers=TEAM)
        assert got.status_code == 200
        assert got.json()["document_type"] == "executive_brief"
        assert got.json()["created_from"] == "generated"


class TestInEditorExport:
    def test_export_docx_from_a_paper_draft(self, clean_editor_drafts):
        if not _db_ready():
            pytest.skip("no live database")
        created = client.post(DRAFTS, headers=TEAM, json={
            "document_type": "midpoint_paper", "title": "Export draft",
            "content_json": {"type": "doc", "content": [
                {"type": "heading", "attrs": {"level": 1},
                 "content": [{"type": "text", "text": "1. Methodology"}]},
                {"type": "paragraph",
                 "content": [{"type": "text", "text": "Body paragraph."}]},
            ]},
            "content_text": "1. Methodology Body paragraph."})
        did = created.json()["id"]
        resp = client.post("/api/v1/export/midpoint-paper", headers=TEAM,
                            json={"editor_draft_id": did})
        assert resp.status_code == 200
        # A valid .docx is a ZIP — the PK magic bytes.
        assert resp.content[:2] == b"PK"
        assert "wordprocessingml" in resp.headers.get("content-type", "")

    def test_export_pptx_carries_editor_speaker_notes(self, clean_editor_drafts):
        if not _db_ready():
            pytest.skip("no live database")
        import io
        from pptx import Presentation

        created = client.post(DRAFTS, headers=TEAM, json={
            "document_type": "presentation_deck", "title": "Deck export",
            "content_json": {"slides": [
                {"id": 1, "title": "Opening", "content": "Body",
                 "data_points": [], "speaker_notes": "UAT-NOTE rehearsal line",
                 "verified": False, "notes_written": True}]},
            "content_text": "Opening"})
        did = created.json()["id"]
        resp = client.post("/api/v1/export/presentation-deck", headers=TEAM,
                            json={"editor_draft_id": did})
        assert resp.status_code == 200
        assert resp.content[:2] == b"PK"
        prs = Presentation(io.BytesIO(resp.content))
        # The presenter's speaker notes survive into the exported file.
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "UAT-NOTE rehearsal line" in notes

    def test_export_unknown_draft_is_404(self):
        resp = client.post("/api/v1/export/midpoint-paper", headers=TEAM,
                            json={"editor_draft_id": 999999999})
        assert resp.status_code == 404


class TestEditorChat:
    CHAT = "/api/v1/documents/drafts/{}/chat"

    def test_chat_requires_auth(self):
        assert client.post(self.CHAT.format(1),
                           json={"message": "hi"}).status_code == 401

    def test_writing_assistant_is_a_logged_interaction_type(self):
        from tools.activity_log import _INTERACTION_TYPES
        assert "writing_assistant" in _INTERACTION_TYPES

    def test_chat_unknown_draft_is_404(self):
        resp = client.post(self.CHAT.format(999999999), headers=TEAM,
                           json={"message": "hi"})
        assert resp.status_code == 404

    def test_chat_returns_a_response(self, clean_editor_drafts):
        if not _db_ready():
            pytest.skip("no live database")
        did = client.post(DRAFTS, headers=TEAM, json={
            "document_type": "midpoint_paper", "title": "Chat draft",
            "content_text": "the document body"}).json()["id"]
        resp = client.post(self.CHAT.format(did), headers=TEAM,
                           json={"message": "Is my conclusion clear?",
                                 "history": [], "selection": None})
        assert resp.status_code == 200
        assert isinstance(resp.json()["response"], str)

    def test_chat_empty_message_is_422(self, clean_editor_drafts):
        if not _db_ready():
            pytest.skip("no live database")
        did = client.post(DRAFTS, headers=TEAM, json={
            "document_type": "midpoint_paper", "title": "Chat draft",
            "content_text": "body"}).json()["id"]
        resp = client.post(self.CHAT.format(did), headers=TEAM,
                           json={"message": "   "})
        assert resp.status_code == 422

    def test_chat_on_a_draft_not_owned_is_404(self, clean_editor_drafts):
        if not _db_ready():
            pytest.skip("no live database")
        # Created by Bob (TEAM) ...
        did = client.post(DRAFTS, headers=TEAM, json={
            "document_type": "midpoint_paper", "title": "Bob's draft",
            "content_text": "body"}).json()["id"]
        # ... and another team member cannot chat against it.
        molly = {"X-API-Key": generate_session_token("murdockm@queens.edu")}
        resp = client.post(self.CHAT.format(did), headers=molly,
                           json={"message": "hi"})
        assert resp.status_code == 404

    def test_chat_prompt_carries_content_selection_and_history(
            self, clean_editor_drafts, monkeypatch):
        if not _db_ready():
            pytest.skip("no live database")
        did = client.post(DRAFTS, headers=TEAM, json={
            "document_type": "midpoint_paper", "title": "Prompt draft",
            "content_text": "DOC-BODY-MARKER unique document text"}).json()["id"]

        captured: dict[str, str] = {}

        def _fake_call_claude(model, system_prompt, user_message,
                              max_tokens=600, tools=None):
            captured["system"] = system_prompt
            captured["user"] = user_message
            return "fake assistant reply"

        # Leave the test-env short-circuit so the real prompt is built.
        monkeypatch.setattr("main.ENVIRONMENT", "production")
        monkeypatch.setattr("agents.base.call_claude", _fake_call_claude)

        resp = client.post(self.CHAT.format(did), headers=TEAM, json={
            "message": "Tighten this",
            "history": [{"role": "user", "content": "earlier question"},
                        {"role": "assistant", "content": "earlier answer"}],
            "selection": "the passage I selected"})
        assert resp.status_code == 200
        assert resp.json()["response"] == "fake assistant reply"
        # The draft's content_text is in the system prompt.
        assert "DOC-BODY-MARKER" in captured["system"]
        # The selection is quoted into the user message.
        assert "> the passage I selected" in captured["user"]
        # The history is flattened in.
        assert "earlier question" in captured["user"]
        assert "Tighten this" in captured["user"]


class TestAcademicReviewReadsDraft:
    def test_review_context_overlays_the_editor_draft(self, clean_editor_drafts):
        if not _db_ready():
            pytest.skip("no live database")
        from agents.academic_review import gather_review_context
        from tools.editor_drafts import create_draft
        from database import engine

        async def _scenario():
            if engine is not None:
                await engine.dispose()
            await create_draft(
                "midpoint_paper", "thaob@queens.edu", "Review draft",
                {"type": "doc", "content": []},
                "DRAFT-MARKER unique editor content for review",
                created_from="generated")
            return await gather_review_context(
                reviewer_email="thaob@queens.edu")

        ctx = _run(_scenario())
        midpoint = ctx["documents_by_type"].get("midpoint_draft", [])
        # The editor draft stands in for the uploaded midpoint document.
        assert any("DRAFT-MARKER" in (d.get("content_text") or "")
                   for d in midpoint)
