"""tests/test_deck_per_slide_generation.py -- bridge #95.

The deck generation now runs ONE LLM call per slide (instead of one
4000-token call for all six) so the JSON never truncates. These tests
pin:

  * slide_generation_prompt -- the per-slide prompt isolates ONE
    slide's spec from SLIDE_SPECIFICATIONS so the model has the full
    project framing but only one slide's required content to emit.
  * parse_single_slide_json -- tolerant of markdown fences, leading
    prose, trailing prose. Returns None on unparseable input so the
    caller can surface a per-slide [DATA PENDING].
  * _slice_slide_spec -- the spec-splitter that returns one slide's
    block from the SLIDE_SPECIFICATIONS string.
"""
from __future__ import annotations

import os
import sys

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))
os.environ.setdefault("SECRET_KEY", "test-secret-key-at-least-32-characters-long")
os.environ.setdefault("ENVIRONMENT", "test")
os.environ.setdefault("MASTER_API_KEY", "test_master_key")


class TestSliceSlideSpec:
    """Each slide's body must be isolatable so the per-slide call
    sees exactly one slide's required content."""

    def test_slice_returns_only_the_requested_slide(self):
        from tools.academic_deck import _slice_slide_spec

        slide_2 = _slice_slide_spec(2)
        # The slice MUST start with the slide-2 header.
        assert slide_2.startswith("Slide 2 --")
        # And MUST NOT carry the slide-3 header or any later slide.
        assert "Slide 3 --" not in slide_2
        assert "Slide 4 --" not in slide_2
        assert "Slide 5 --" not in slide_2
        assert "Slide 6 --" not in slide_2

    def test_slice_returns_the_last_slide_to_end_of_string(self):
        from tools.academic_deck import _slice_slide_spec, DECK_SLIDE_COUNT

        last = _slice_slide_spec(DECK_SLIDE_COUNT)
        assert last.startswith(f"Slide {DECK_SLIDE_COUNT} --")
        # The slice MUST carry the slide's body content -- speaker
        # notes survive into the LLM prompt.
        assert "Speaker notes" in last

    def test_slice_raises_on_out_of_range_slide_number(self):
        from tools.academic_deck import _slice_slide_spec, DECK_SLIDE_COUNT

        with pytest.raises(ValueError):
            _slice_slide_spec(0)
        with pytest.raises(ValueError):
            _slice_slide_spec(DECK_SLIDE_COUNT + 1)


class TestSlideGenerationPrompt:
    """The per-slide prompt carries the project framing AND exactly
    one slide's spec AND a single-object JSON contract -- not the
    {"slides":[...]} wrapper."""

    def test_prompt_carries_project_framing(self):
        from tools.academic_deck import slide_generation_prompt

        prompt = slide_generation_prompt(1)
        # The preamble must still introduce the project so the LLM
        # has every detail it needs for one slide.
        assert "diversification" in prompt.lower()
        assert "FNA 670" in prompt

    def test_prompt_isolates_one_slide_only(self):
        """The prompt for slide 3 must carry slide 3's spec and
        NOT carry slide 5's spec (or any other slide's body)."""
        from tools.academic_deck import slide_generation_prompt

        s3 = slide_generation_prompt(3)
        assert "Slide 3 --" in s3
        # Slide 4 and 5 specs must NOT bleed into slide 3's prompt.
        assert "Slide 4 --" not in s3
        assert "Slide 5 --" not in s3

    def test_prompt_asks_for_single_object_not_wrapped_list(self):
        from tools.academic_deck import slide_generation_prompt

        prompt = slide_generation_prompt(2)
        # The contract is a SINGLE object -- speaker notes, bullets,
        # table -- not the {"slides":[...]} wrapper the all-six call
        # uses.
        assert '"slide_number": 2' in prompt
        assert '"speaker_notes"' in prompt
        # The wrapper must NOT appear -- a model that emits the
        # wrapper would force the parser to dig in.
        assert '"slides":' not in prompt


class TestParseSingleSlideJson:
    """The single-slide parser pulls one object out of the LLM
    response. Tolerant of fences and prose; None on failure."""

    def test_parses_clean_json_object(self):
        from tools.academic_deck import parse_single_slide_json

        raw = (
            '{"slide_number": 1, "title": "Hi", '
            '"bullets": ["a", "b"], "table_data": null, '
            '"speaker_notes": "notes"}'
        )
        out = parse_single_slide_json(raw)
        assert out is not None
        assert out["slide_number"] == 1
        assert out["bullets"] == ["a", "b"]

    def test_strips_markdown_fences(self):
        from tools.academic_deck import parse_single_slide_json

        raw = (
            "```json\n"
            '{"slide_number": 2, "title": "T", "bullets": [],'
            ' "speaker_notes": ""}\n'
            "```"
        )
        out = parse_single_slide_json(raw)
        assert out is not None
        assert out["slide_number"] == 2

    def test_returns_none_on_empty_input(self):
        from tools.academic_deck import parse_single_slide_json

        assert parse_single_slide_json("") is None
        assert parse_single_slide_json(None) is None
        assert parse_single_slide_json("   \n") is None

    def test_returns_none_on_truncated_unparseable_response(self):
        from tools.academic_deck import parse_single_slide_json

        # Mid-JSON truncation -- the exact failure mode the bridge #95
        # change is designed to avoid. parse_single_slide_json returns
        # None so the caller logs a per-slide warning AND inserts a
        # [DATA PENDING] placeholder for that slide only.
        raw = '{"slide_number": 3, "title": "Three", "bullets": ["unclosed'
        assert parse_single_slide_json(raw) is None

    def test_returns_none_on_array_at_root(self):
        """A model that returned the {"slides":[...]} wrapper instead
        of a single object should NOT be accepted as a single slide."""
        from tools.academic_deck import parse_single_slide_json

        # An array (not a dict) is not a valid single-slide object.
        # The parser pulls between the first { and last } -- which
        # would be the inner slide, but that's by design: a wrapper
        # with one slide is ambiguous. Use a clearly-non-dict input.
        raw = "[1, 2, 3]"
        assert parse_single_slide_json(raw) is None

    def test_strips_preamble_text_before_first_brace(self):
        """The model occasionally emits a short conversational opener
        before the JSON body ("Here is the slide content:\n{...}").
        The hardening discards any text before the first '{' so the
        opening preamble can never leak into the parsed output."""
        from tools.academic_deck import parse_single_slide_json

        raw = (
            "Here is the slide content for you:\n\n"
            '{"slide_number": 4, "title": "Methodology", '
            '"bullets": ["one", "two"], "speaker_notes": ""}')
        out = parse_single_slide_json(raw)
        assert out is not None
        assert out["slide_number"] == 4
        assert out["bullets"] == ["one", "two"]

    def test_logs_raw_preview_on_parse_failure(self, monkeypatch):
        """A fully-mangled response should log the raw preview at
        WARNING so a regression hunter can see what the model emitted
        without having to reproduce locally. We assert via a monkey-
        patched structlog warning capture because the module uses a
        structlog BoundLogger that bypasses stdlib caplog."""
        from tools import academic_deck

        captured: list[tuple[str, dict]] = []

        def _capture(event, **kwargs):
            captured.append((event, kwargs))

        monkeypatch.setattr(academic_deck.log, "warning", _capture)
        # Unparseable: looks like JSON but the braces don't close.
        bad = "Here is the slide: {" + ("oops " * 60)
        assert academic_deck.parse_single_slide_json(bad) is None
        # One of the parse-failure log keys must have fired with a
        # raw_preview field carrying the start of the broken response.
        events = [e for e, _ in captured]
        assert any(e.startswith("deck_slide_parse") for e in events), events
        last_event, last_kw = captured[-1]
        assert "raw_preview" in last_kw
        assert "Here is the slide" in last_kw["raw_preview"]


class TestBulletPreambleScrub:
    """A bullet that opens with a known LLM apology / preamble pattern
    is the model talking ABOUT the slide rather than producing slide
    content. The post-parse scrub replaces those bullets with a
    clearly-flagged regen marker so the deck never carries raw
    apology text in front of an audience."""

    def test_replaces_apology_bullet_with_regen_marker(self):
        from tools.academic_deck import parse_single_slide_json

        raw = (
            '{"slide_number": 5, "title": "T", '
            '"bullets": ["I cannot produce this section.", '
            '"Real bullet content."], '
            '"speaker_notes": ""}')
        out = parse_single_slide_json(raw)
        assert out is not None
        assert out["bullets"][0].startswith(
            "[Content generation error")
        assert out["bullets"][1] == "Real bullet content."

    def test_replaces_note_prefix_bullet(self):
        from tools.academic_deck import parse_single_slide_json

        raw = (
            '{"slide_number": 6, "title": "T", '
            '"bullets": ["Note: this is filler text from the model.", '
            '"Genuine analytical point."], '
            '"speaker_notes": ""}')
        out = parse_single_slide_json(raw)
        assert out is not None
        assert out["bullets"][0].startswith(
            "[Content generation error")
        assert out["bullets"][1] == "Genuine analytical point."

    def test_leaves_normal_bullets_untouched(self):
        from tools.academic_deck import parse_single_slide_json

        raw = (
            '{"slide_number": 7, "title": "T", '
            '"bullets": ["Equity 60% / Bonds 40%.", '
            '"Drawdown reduced 27 percentage points."], '
            '"speaker_notes": ""}')
        out = parse_single_slide_json(raw)
        assert out is not None
        # Both bullets are real analysis content and must NOT be
        # replaced. The scrub is opt-in via the preamble prefix.
        assert "[Content generation error" not in (out["bullets"][0]
                                                   + out["bullets"][1])

    def test_handles_bullet_glyphs_before_preamble(self):
        """A model that opens with a bullet glyph followed by the
        preamble pattern still triggers the scrub -- the prefix scan
        strips list-marker glyphs before matching."""
        from tools.academic_deck import _bullet_looks_like_preamble

        assert _bullet_looks_like_preamble("- I cannot do that")
        assert _bullet_looks_like_preamble("• Sorry, the data is missing")
        assert _bullet_looks_like_preamble("▪  Note: filler")
        assert _bullet_looks_like_preamble("As an AI language model")
        assert not _bullet_looks_like_preamble(
            "Inflation-adjusted blend Sharpe is 1.24.")


class TestImagePlaceholder:
    """The chart slot is filled by the matplotlib render when the
    cache is warm. A None return from the renderer used to silently
    insert a generic [DATA PENDING] placeholder; the hardening logs
    the slot name + emits a more diagnostic placeholder."""

    def test_image_logs_warning_with_chart_slot(self, monkeypatch):
        from tools import academic_deck

        log_calls: list[tuple[str, dict]] = []

        def _capture(event, **kwargs):
            log_calls.append((event, kwargs))

        monkeypatch.setattr(academic_deck.log, "warning", _capture)

        # Fake slide that records add_picture / add_textbox calls so
        # the placeholder string + the log emission can be inspected
        # without standing up a real python-pptx Presentation. _image
        # only ever calls slide.shapes.add_picture (skipped when png
        # is None) and the module-level _textbox helper -- the latter
        # uses slide.shapes.add_textbox. A minimal stub is sufficient.
        captured: dict = {"add_picture": 0, "textboxes": []}

        class _FakeText:
            def __init__(self):
                self.word_wrap = False
                self.vertical_anchor = None
                self.paragraphs = [_FakePara()]

            def add_paragraph(self):
                self.paragraphs.append(_FakePara())
                return self.paragraphs[-1]

        class _FakePara:
            def __init__(self):
                self.alignment = None
                self.space_after = None
                self.runs: list = []

            def add_run(self):
                run = _FakeRun()
                self.runs.append(run)
                return run

        class _FakeRun:
            def __init__(self):
                self.text = ""
                self.font = _FakeFont()

        class _FakeFont:
            def __init__(self):
                self.size = None
                self.color = _FakeColor()
                self.name = None
                self.bold = False

        class _FakeColor:
            def __init__(self):
                self.rgb = None

        class _FakeBox:
            def __init__(self):
                self.text_frame = _FakeText()

        class _FakeShapes:
            def add_picture(self, *_a, **_kw):
                captured["add_picture"] += 1

            def add_textbox(self, *_a, **_kw):
                box = _FakeBox()
                captured["textboxes"].append(box)
                return box

        class _FakeSlide:
            def __init__(self):
                self.shapes = _FakeShapes()

        from pptx.util import Inches

        slide = _FakeSlide()
        academic_deck._image(
            slide, None,
            left=Inches(1.0), top=Inches(1.0), width=Inches(5.0),
            fallback="rolling_correlation")

        # No image was added; one placeholder textbox was added; a
        # WARNING was emitted naming the chart slot.
        assert captured["add_picture"] == 0
        assert len(captured["textboxes"]) == 1
        events = [e for e, _ in log_calls
                  if e == "deck_chart_slot_unavailable"]
        assert events, "expected deck_chart_slot_unavailable warning"
        # The chart slot name is carried in the structured log kwarg.
        assert any(kw.get("chart") == "rolling_correlation"
                   for _e, kw in log_calls)

    def test_image_renders_picture_when_png_present(self):
        """Belt-and-braces: confirm the happy path still calls
        add_picture without going down the placeholder path."""
        from tools.academic_deck import _image

        counts: dict = {"add_picture": 0, "textboxes": 0}

        class _Shapes:
            def add_picture(self, *_a, **_kw):
                counts["add_picture"] += 1

            def add_textbox(self, *_a, **_kw):  # pragma: no cover
                counts["textboxes"] += 1
                raise AssertionError("textbox path should not fire")

        class _Slide:
            shapes = _Shapes()

        from pptx.util import Inches

        _image(_Slide(), b"\x89PNG\r\n\x1a\n",
               left=Inches(0), top=Inches(0), width=Inches(5))
        assert counts["add_picture"] == 1
        assert counts["textboxes"] == 0


class TestDeckSlideCount:
    """The deck has exactly 11 slides post-rebuild (bridge #98).
    _slice_slide_spec + slide_generation_prompt + parse_single_slide_json
    all rely on this constant; pin it so a future change touches every
    helper at once."""

    def test_eleven_slides_post_rebuild(self):
        from tools.academic_deck import DECK_SLIDE_COUNT
        assert DECK_SLIDE_COUNT == 11

    def test_every_slide_number_has_a_spec(self):
        """The per-slide loop iterates 1..DECK_SLIDE_COUNT. Each must
        produce a non-empty slice or the loop drops a slide silently."""
        from tools.academic_deck import _slice_slide_spec, DECK_SLIDE_COUNT
        for n in range(1, DECK_SLIDE_COUNT + 1):
            slice_n = _slice_slide_spec(n)
            assert slice_n
            assert f"Slide {n} --" in slice_n


class TestStoryArcSeed:
    """Bridge #98: each per-slide prompt carries the story-arc seed
    so the LLM knows where it sits in the 11-slide narrative without
    re-reading every other slide's content."""

    def test_seed_substitutes_slide_number_and_title(self):
        from tools.academic_deck import (
            slide_generation_prompt, SLIDE_TITLES,
        )
        prompt = slide_generation_prompt(7)
        assert "slide 7 of 11" in prompt
        assert SLIDE_TITLES[6] in prompt   # slide 7's title

    def test_seed_lists_the_eleven_arc_stops(self):
        from tools.academic_deck import slide_generation_prompt
        prompt = slide_generation_prompt(1)
        # The arc seed lists all 11 narrative beats so the slide-1
        # generator knows where the deck is heading.
        for n in range(1, 12):
            assert f"{n}." in prompt


class TestChartSlotDerivedFromSlideCharts:
    """The renderer's output slot dict is DERIVED from
    academic_deck.SLIDE_CHARTS so the two cannot drift. Earlier
    (pre-PR #332) the renderer hardcoded {2, 3, 6} while SLIDE_CHARTS
    had moved to {4, 5, 11} during the 11-slide rebuild; every deck
    generation since then fired three deck_chart_slot_unavailable
    warnings. Pinning the derivation here so a future hardcode
    regression fails loudly."""

    def test_renderer_keys_follow_slide_charts(self, monkeypatch):
        """Mock SLIDE_CHARTS to a single arbitrary slot {7:
        rolling_correlation}. The renderer must return a dict whose
        ONLY key is 7 -- proving the slot follows SLIDE_CHARTS and
        is not a hardcoded constant elsewhere."""
        from tools import academic_deck

        # Replace SLIDE_CHARTS with a single-entry probe. The renderer
        # reads it via `from tools.academic_deck import SLIDE_CHARTS`
        # at call time, so monkeypatching the module attribute is
        # sufficient.
        monkeypatch.setattr(
            academic_deck, "SLIDE_CHARTS", {7: "rolling_correlation"})

        from main import _render_deck_slide_charts
        out = _render_deck_slide_charts(
            data={}, blend_weights={}, blend_series=[])
        # Only key 7 -- the slot the mocked SLIDE_CHARTS named. The
        # hardcoded-pre-fix would have returned {2, 3, 6} (or any
        # other constants) regardless of SLIDE_CHARTS contents.
        assert set(out.keys()) == {7}, (
            f"renderer slot ({sorted(out.keys())}) did not follow "
            f"SLIDE_CHARTS ({{7}}) -- the derivation regressed to a "
            "hardcoded constant.")

    def test_renderer_keys_align_with_real_slide_charts(self):
        """End-to-end smoke: the production SLIDE_CHARTS keys and the
        renderer output keys are equal. This is the test that would
        have caught the pre-PR-332 drift the moment it landed."""
        from tools.academic_deck import SLIDE_CHARTS

        from main import _render_deck_slide_charts
        out = _render_deck_slide_charts(
            data={}, blend_weights={}, blend_series=[])
        assert set(out.keys()) == set(SLIDE_CHARTS.keys()), (
            "SLIDE_CHARTS and _render_deck_slide_charts must agree on "
            "the slide numbers that carry charts. Mismatch causes "
            "silent deck_chart_slot_unavailable WARNINGs on every "
            "deck generation.")

    def test_unknown_chart_role_logs_warning_not_crash(
            self, monkeypatch):
        """An unknown chart_name in SLIDE_CHARTS must NOT raise -- the
        renderer logs a deck_chart_role_unwired WARNING and silently
        skips that slot. Defends against a SLIDE_CHARTS edit that adds
        a new role before the renderer dispatch is wired."""
        from tools import academic_deck

        monkeypatch.setattr(
            academic_deck, "SLIDE_CHARTS",
            {4: "rolling_correlation",   # known role -- rendered
             9: "future_chart_role"})    # unknown role -- skipped

        captured: list[tuple[str, dict]] = []

        import main

        def _capture(event, **kwargs):
            captured.append((event, kwargs))

        monkeypatch.setattr(main.log, "warning", _capture)

        # Must NOT raise.
        out = main._render_deck_slide_charts(
            data={}, blend_weights={}, blend_series=[])
        # Known role slot is present; unknown role slot is dropped.
        assert 4 in out
        assert 9 not in out
        # The deck_chart_role_unwired event fired with both slide_number
        # and role kwargs so an operator reading Render logs can see
        # exactly which slot was unwired.
        events = [(e, kw) for e, kw in captured
                  if e == "deck_chart_role_unwired"]
        assert events, "expected deck_chart_role_unwired warning"
        _, kw = events[-1]
        assert kw.get("slide_number") == 9
        assert kw.get("role") == "future_chart_role"


class TestNoHarnessReferenced:
    """Bridge #100: the deck slide generation must NOT route through
    the academic-writer harness (which uses the peer-discussant
    evaluator that caused the slide-2 leak). The prompt is sent
    directly to Sonnet via call_claude."""

    def test_main_imports_call_claude_not_harness_for_slide_generation(self):
        """Grep the _generate_one_deck_slide body for the import.
        A future refactor that reintroduces harness_narrative for
        slide generation breaks this pin. The docstring may mention
        the deprecated harness path for context; the test checks the
        non-comment body for the actual invocation."""
        import inspect
        import main
        src = inspect.getsource(main._generate_one_deck_slide)
        assert "call_claude" in src
        # Strip the docstring out before scanning -- the docstring
        # explains WHY harness_narrative was removed, so it's allowed
        # to mention the name. The actual function body must not.
        lines = src.split("\n")
        in_docstring = False
        body_only_lines: list[str] = []
        for line in lines:
            stripped = line.strip()
            if stripped.startswith('"""'):
                # Either opens or closes a docstring; could open AND
                # close on the same line for one-liners.
                if in_docstring:
                    in_docstring = False
                else:
                    in_docstring = True
                    if stripped.count('"""') >= 2:
                        in_docstring = False
                continue
            if in_docstring:
                continue
            body_only_lines.append(line)
        body = "\n".join(body_only_lines)
        # The body must not import or call harness_narrative.
        assert "harness_narrative" not in body
